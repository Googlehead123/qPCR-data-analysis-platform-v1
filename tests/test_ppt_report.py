"""Tests for the PPT generation path the app actually uses.

Covers PPTGenerator: the Cosmax-branded, template-driven writer -- display
names, template detail fields, Korean font runs, and the direction-checked
efficacy verdict.

The old from-scratch ReportGenerator slide builders were removed along with
their tests: they had no app call site and never received the correctness
fixes this path did, so they would have emitted a verdict that ignores
expected_direction if ever re-enabled.
"""

import io
from unittest.mock import patch, MagicMock
import pandas as pd
import pytest


# Create a minimal PNG bytes for mocking to_image
MOCK_PNG_BYTES = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00"
    b"\x00\x01\x01\x00\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
)


class TestPPTGeneratorDisplayName:
    """A rename of a gene in the Graphs tab must reach the PPT slide title."""

    def test_create_gene_slide_uses_display_name_in_title(self, mock_streamlit, processed_gene_data):
        from importlib import import_module
        from pptx import Presentation
        from pptx.util import Emu
        import plotly.graph_objects as go

        spec = import_module("streamlit qpcr analysis v1")
        PPTGenerator = spec.PPTGenerator

        fig = go.Figure(data=[go.Bar(x=["A", "B"], y=[1, 2])])
        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)

        with patch.object(go.Figure, "to_image", return_value=MOCK_PNG_BYTES):
            PPTGenerator.create_gene_slide(
                prs, "COL1A1", fig, processed_gene_data,
                {"Housekeeping_Gene": "GAPDH"},
                display_name="Collagen I",
            )

        slide = prs.slides[0]
        title_texts = [
            shape.text_frame.text for shape in slide.shapes
            if shape.has_text_frame
        ]
        assert any("Collagen I Expression" in t for t in title_texts), (
            f"Expected 'Collagen I Expression' in slide titles, got: {title_texts}"
        )
        # The raw gene name must NOT appear in the title — that's the bug we fix.
        assert not any("COL1A1 Expression" in t for t in title_texts), (
            f"Raw gene name leaked into slide title: {title_texts}"
        )

    def test_create_gene_slide_falls_back_to_raw_name_without_display(self, mock_streamlit, processed_gene_data):
        """When no display_name is passed, behavior is unchanged from before."""
        from importlib import import_module
        from pptx import Presentation
        from pptx.util import Emu
        import plotly.graph_objects as go

        spec = import_module("streamlit qpcr analysis v1")
        PPTGenerator = spec.PPTGenerator

        fig = go.Figure(data=[go.Bar(x=["A", "B"], y=[1, 2])])
        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)

        with patch.object(go.Figure, "to_image", return_value=MOCK_PNG_BYTES):
            PPTGenerator.create_gene_slide(
                prs, "COL1A1", fig, processed_gene_data,
                {"Housekeeping_Gene": "GAPDH"},
            )

        slide = prs.slides[0]
        title_texts = [
            shape.text_frame.text for shape in slide.shapes
            if shape.has_text_frame
        ]
        assert any("COL1A1 Expression" in t for t in title_texts)


class TestTemplateDetailFields:
    """Pins the template detail box on the RIGHT source for every field.

    `_populate_gene_slide` finds this box by geometry (top-right of the slide),
    then rewrites each paragraph whose first run starts with a known label.
    Two of those labels come from operator input rather than from the data, and
    "Sample concentration" was silently wrong for months: the widget feeding
    `analysis_params["concentration"]` was removed while the writer kept a
    `"1 ppm"` fallback, so every slide printed a plausible, unverifiable dose.
    These tests fail if either half of that regression comes back.
    """

    DETAIL_LABELS = [
        "Date: ",
        "Cell line: ",
        "Sample concentration: ",
        "Positive control: ",
        "Inducer: ",
        "Treatment time: ",
        "Test method:",
    ]

    def _slide_with_detail_box(self):
        """A slide carrying a stand-in for the template's top-right detail box.

        Geometry must satisfy the writer's `left > 3500000 and top < 500000`
        test, and the text must avoid the earlier `효능 평가` / `Results`
        branches, or the box is claimed by one of those instead.
        """
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        box = slide.shapes.add_textbox(Emu(7000000), Emu(200000), Emu(4000000), Emu(2000000))
        tf = box.text_frame
        for i, label in enumerate(self.DETAIL_LABELS):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.add_run().text = label
        return prs, slide

    def _detail_text(self, slide):
        """The detail box's text, as {label: value} for the labels we wrote."""
        from pptx.util import Emu

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.left > 3500000 and shape.top < 500000:
                lines = [p.runs[0].text for p in shape.text_frame.paragraphs if p.runs]
                out = {}
                for line in lines:
                    for label in self.DETAIL_LABELS:
                        if line.startswith(label):
                            out[label.strip()] = line[len(label):]
                            break
                return out
        raise AssertionError("detail box not found on slide")

    def _populate(self, spec, prs, slide, analysis_params):
        import plotly.graph_objects as go

        fig = go.Figure(data=[go.Bar(x=["A"], y=[1])])
        with patch.object(go.Figure, "to_image", return_value=MOCK_PNG_BYTES):
            spec.PPTGenerator._populate_gene_slide(
                prs, slide, "MMP1", fig, pd.DataFrame(), analysis_params, {},
            )

    def test_operator_supplied_concentration_reaches_the_slide(self, mock_streamlit):
        """Whatever the operator typed must be what the slide says."""
        from importlib import import_module

        spec = import_module("streamlit qpcr analysis v1")
        prs, slide = self._slide_with_detail_box()

        self._populate(spec, prs, slide, {
            "Efficacy_Type": "광노화",
            "Date": "2026-07-31 12:00:00",
            "concentration": "0.05%",
        })

        fields = self._detail_text(slide)
        assert fields["Sample concentration:"] == "0.05%"
        assert fields["Date:"] == "2026-07-31"

    def test_missing_concentration_is_blank_not_guessed(self, mock_streamlit):
        """No input → an empty field, never a default dose.

        A blank is visible to whoever reviews the deck; "1 ppm" is not. This is
        the assertion that keeps the old fallback from being reintroduced as a
        convenience.
        """
        from importlib import import_module

        spec = import_module("streamlit qpcr analysis v1")

        for params in (
            {"Efficacy_Type": "광노화", "Date": "2026-07-31"},          # key absent
            {"Efficacy_Type": "광노화", "Date": "2026-07-31", "concentration": ""},
        ):
            prs, slide = self._slide_with_detail_box()
            self._populate(spec, prs, slide, params)
            fields = self._detail_text(slide)
            assert fields["Sample concentration:"] == "", (
                f"expected a blank concentration, got {fields['Sample concentration:']!r}"
            )

    def test_inducer_prints_the_display_value_when_the_catalog_has_one(self, mock_streamlit):
        """Where the ledger wording can't be the matching key, the slide still
        shows the dose: 광노화 matches on "UVB only" but must print
        "UVB 40 mj/cm2"."""
        from importlib import import_module
        from qpcr.constants import EFFICACY_CONFIG

        spec = import_module("streamlit qpcr analysis v1")
        controls = EFFICACY_CONFIG["광노화"]["controls"]

        prs, slide = self._slide_with_detail_box()
        self._populate(spec, prs, slide, {
            "Efficacy_Type": "광노화", "Date": "2026-07-31", "concentration": "1 ppm",
        })

        fields = self._detail_text(slide)
        assert fields["Inducer:"] == controls["negative_display"] == "UVB 40 mj/cm2"
        # And the matching key itself is NOT what got printed.
        assert fields["Inducer:"] != controls["negative"]

    def test_inducer_falls_back_to_the_matching_key_without_a_display_value(
        self, mock_streamlit
    ):
        """Categories with no display override behave exactly as before.

        장벽 has no `negative_display`, so the slide shows "Non-treated" — the
        backward-compatible path, and the reason adding this field to 7 of 21
        categories did not need the other 14 touched.
        """
        from importlib import import_module
        from qpcr.constants import EFFICACY_CONFIG

        spec = import_module("streamlit qpcr analysis v1")
        controls = EFFICACY_CONFIG["장벽"]["controls"]
        assert "negative_display" not in controls, "fixture assumes no override here"

        prs, slide = self._slide_with_detail_box()
        self._populate(spec, prs, slide, {
            "Efficacy_Type": "장벽", "Date": "2026-07-31", "concentration": "1 ppm",
        })

        fields = self._detail_text(slide)
        assert fields["Inducer:"] == controls["negative"] == "Non-treated"
        assert fields["Positive control:"] == "Retinoic acid"

    def test_controls_come_from_the_efficacy_catalog(self, mock_streamlit):
        """Positive control and inducer are catalog-owned, not operator-typed.

        Unlike concentration these ARE knowable from the chosen efficacy type,
        so they must be filled from `EFFICACY_CONFIG` — including any dose
        recorded there.
        """
        from importlib import import_module
        from qpcr.constants import EFFICACY_CONFIG

        spec = import_module("streamlit qpcr analysis v1")
        controls = EFFICACY_CONFIG["광노화"]["controls"]

        prs, slide = self._slide_with_detail_box()
        self._populate(spec, prs, slide, {
            "Efficacy_Type": "광노화", "Date": "2026-07-31", "concentration": "1 ppm",
        })

        fields = self._detail_text(slide)
        assert fields["Positive control:"] == controls["positive"]
        # Inducer resolves through negative_display when set (see the two tests
        # above for both branches); 광노화 has one.
        assert fields["Inducer:"] == (
            controls.get("negative_display") or controls["negative"]
        )
        assert fields["Cell line:"] == EFFICACY_CONFIG["광노화"]["cell"]


class TestKoreanRunsNameAKoreanFace:
    """Every Hangul-bearing PPT run must declare an East-Asian face.

    All template runs were Arial, which has zero Hangul coverage, so 진정 효능
    평가 and Results: 효능 有 rendered by host substitution — visibly a different
    typeface from the chart image's Korean inside the same slide. python-pptx's
    font.name writes a:latin only; a:ea is what East-Asian runs read.
    """

    NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    def test_set_korean_face_writes_both_latin_and_ea(self, mock_streamlit):
        from importlib import import_module
        from pptx import Presentation
        from pptx.util import Emu
        from qpcr.constants import KOREAN_FONT_NAME

        spec = import_module("streamlit qpcr analysis v1")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000000), Emu(500000))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "진정 효능 평가"

        spec._set_korean_face(run)

        assert run.font.name == KOREAN_FONT_NAME
        rPr = run._r.find(f"{self.NS}rPr")
        ea = rPr.find(f"{self.NS}ea")
        assert ea is not None, "a:ea missing — East-Asian runs read that one"
        assert ea.get("typeface") == KOREAN_FONT_NAME

    def test_deck_sweep_covers_untouched_template_korean(self, mock_streamlit):
        """The template's own Korean is not written by any writer."""
        from importlib import import_module
        from pptx import Presentation
        from pptx.util import Emu
        from qpcr.constants import KOREAN_FONT_NAME

        spec = import_module("streamlit qpcr analysis v1")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(1000000), Emu(500000))
        p = box.text_frame.paragraphs[0]
        kr = p.add_run()
        kr.text = "효능 평가"
        en = p.add_run()
        en.text = "Results only"

        spec._apply_korean_face_to_deck(prs)

        rPr = kr._r.find(f"{self.NS}rPr")
        ea = rPr.find(f"{self.NS}ea")
        assert ea is not None and ea.get("typeface") == KOREAN_FONT_NAME

        # A Latin-only run is left alone — no need to name a Korean face on it.
        en_rPr = en._r.find(f"{self.NS}rPr")
        assert en_rPr is None or en_rPr.find(f"{self.NS}ea") is None


class TestEfficacyVerdictRespectsDirection:
    """The client slide's 효능 有/無 must agree with expected_direction.

    The verdict tested only that a significance marker existed on a test-article
    row, so a significant move AGAINST the configured mechanism was stamped
    효능 有: MMP1 is configured "down" for 광노화, and a 4x RISE read as a
    success. Approved 2026-08-24: require significance AND the right direction,
    with a third verdict when a significant marker went the wrong way.
    """

    def _slide_with_results_box(self):
        from pptx import Presentation
        from pptx.util import Emu

        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Left/top keep it out of the top-right detail-box branch.
        box = slide.shapes.add_textbox(Emu(500000), Emu(3000000), Emu(3000000), Emu(500000))
        box.text_frame.paragraphs[0].add_run().text = "Results: 효능 有"
        return prs, slide

    def _verdict(self, gene, rows):
        from importlib import import_module
        import plotly.graph_objects as go

        spec = import_module("streamlit qpcr analysis v1")
        prs, slide = self._slide_with_results_box()
        fig = go.Figure(data=[go.Bar(x=["A"], y=[1])])
        params = {
            "Efficacy_Type": "광노화",
            "Reference_Sample": "Non-treated",
            "Compare_To": "Non-treated",
        }
        with patch.object(go.Figure, "to_image", return_value=MOCK_PNG_BYTES):
            spec.PPTGenerator._populate_gene_slide(
                prs, slide, gene, fig, pd.DataFrame(rows), params, {},
            )
        for shape in slide.shapes:
            if shape.has_text_frame and "Results" in shape.text_frame.text:
                return shape.text_frame.text
        raise AssertionError("results box not found")

    def test_significant_but_wrong_direction_is_not_efficacy(self):
        """MMP1 is configured 'down'; a significant 4x RISE is not 효능 有."""
        text = self._verdict("MMP1", [
            {"Condition": "Test article", "Fold_Change": 4.0,
             "p_value": 0.001, "significance": "**"},
        ])
        assert "有" not in text, (
            f"a significant move against expected_direction must not read as "
            f"efficacy; got {text!r}"
        )
        assert "재검토" in text

    def test_significant_in_the_expected_direction_is_efficacy(self):
        text = self._verdict("MMP1", [
            {"Condition": "Test article", "Fold_Change": 0.4,
             "p_value": 0.001, "significance": "**"},
        ])
        assert "有" in text

    def test_no_significant_marker_is_no_efficacy(self):
        text = self._verdict("MMP1", [
            {"Condition": "Test article", "Fold_Change": 0.4,
             "p_value": 0.4, "significance": ""},
        ])
        assert "無" in text

    def test_unconfigured_marker_keeps_the_historical_verdict(self):
        """A gene with no expected_direction entry has nothing to check.

        Decision (Min, 2026-08-24): fall back to significance-only rather than
        3rd-verdict it. EFFICACY_CONFIG covers 73/83 markers and the gaps
        cluster — 립 색상 has none of its six — so 3rd-verdicting the
        unconfigured case would have flipped whole categories of report.
        """
        text = self._verdict("NOTAGENE", [
            {"Condition": "Test article", "Fold_Change": 4.0,
             "p_value": 0.001, "significance": "**"},
        ])
        assert "有" in text
        assert "재검토" not in text

    def test_a_whole_category_without_directions_is_unaffected(self):
        """립 색상 configures no direction for any of its six markers."""
        from qpcr.constants import EFFICACY_CONFIG
        from qpcr.auto import expected_direction_for

        cfg = EFFICACY_CONFIG["립 색상"]
        assert not any(
            expected_direction_for(g, cfg.get("expected_direction", {}))
            for g in cfg["genes"]
        ), "fixture assumption changed — 립 색상 now has directions configured"
