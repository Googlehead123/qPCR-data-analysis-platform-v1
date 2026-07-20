"""Report generation — ReportGenerator (chart images) + PPTGenerator (decks).

Single source of truth (extracted from the monolith). PPTGenerator keeps a small
Streamlit coupling (session-state reads for lazy figure render + a user-facing
error) which is acceptable for this app.
"""
import io
import os
import logging
from datetime import datetime
from typing import Dict

import numpy as np
import pandas as pd
import plotly.graph_objects as go
import streamlit as st

from qpcr.constants import PLOTLY_FONT_FAMILY, CM_TO_PX, CM_TO_EMU, EFFICACY_CONFIG
from qpcr.export_utils import export_figure_to_bytes
from qpcr.graph import GraphGenerator


class ReportGenerator:
    SLIDE_WIDTH_INCHES = 13.333
    SLIDE_HEIGHT_INCHES = 7.5

    @staticmethod
    def _fig_to_image(fig: go.Figure, format: str = "png", scale: int = 2,
                      width: int = None, height: int = None) -> bytes:
        """Convert a Plotly figure to image bytes.

        Delegates to the shared, environment-hardened renderer in
        ``qpcr.export_utils`` (browser auto-detection + Chrome-download fallback).
        """
        return export_figure_to_bytes(fig, fmt=format, scale=scale, width=width, height=height)

    @staticmethod
    def create_presentation(
        graphs: Dict[str, go.Figure],
        processed_data: Dict[str, pd.DataFrame],
        analysis_params: dict,
        layout: str = "one_per_slide",
        include_title_slide: bool = True,
        include_summary: bool = True,
        method_text: str = "",
    ) -> bytes:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Emu
            from pptx.dml.color import RGBColor
        except ImportError:
            raise ImportError(
                "python-pptx is required. Install with: pip install python-pptx"
            )

        prs = Presentation()
        prs.slide_width = Inches(ReportGenerator.SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(ReportGenerator.SLIDE_HEIGHT_INCHES)

        blank_layout = prs.slide_layouts[6]

        if include_title_slide:
            ReportGenerator._add_title_slide(prs, analysis_params)

        for gene, fig in graphs.items():
            gene_data = processed_data.get(gene)
            ReportGenerator._add_gene_slide(
                prs, blank_layout, gene, fig, gene_data, method_text
            )

        if include_summary:
            ReportGenerator._add_summary_slide(
                prs, blank_layout, processed_data, analysis_params
            )

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()

    @staticmethod
    def _add_title_slide(prs, analysis_params: dict):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Set slide background to white
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White

        # Add logo placeholder in top-left corner
        logo_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(2), Inches(0.8)
        )
        logo_frame = logo_box.text_frame
        logo_frame.word_wrap = True
        logo_para = logo_frame.paragraphs[0]
        logo_para.text = "[Add Logo Here]"
        logo_para.font.size = Pt(12)
        logo_para.font.color.rgb = RGBColor(0xC1, 0xC6, 0xC7)  # Frost Grey

        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.33), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "qPCR Gene Expression Analysis"
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle with efficacy type
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.8), Inches(12.33), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        efficacy = analysis_params.get("Efficacy_Type", "Analysis")
        subtitle_para.text = f"{efficacy} Efficacy Study"
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Info box with date, HK gene, and reference sample
        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5), Inches(12.33), Inches(1)
        )
        info_frame = info_box.text_frame
        info_para = info_frame.paragraphs[0]
        date = analysis_params.get("Date", datetime.now().strftime("%Y-%m-%d"))
        hk = analysis_params.get("Housekeeping_Gene", "N/A")
        ref = analysis_params.get("Reference_Sample", "N/A")
        info_para.text = f"Date: {date}  |  HK Gene: {hk}  |  Reference: {ref}"
        info_para.font.size = Pt(14)
        info_para.font.color.rgb = RGBColor(0xC1, 0xC6, 0xC7)  # Frost Grey
        info_para.alignment = PP_ALIGN.CENTER

        # Add Cosmax Red accent bar at bottom (full width, 0.5" height)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(0xEA, 0x1D, 0x22)  # Cosmax Red
        accent_bar.line.fill.background()  # No border

    @staticmethod
    def _add_gene_slide(
        prs,
        layout,
        gene: str,
        fig: go.Figure,
        gene_data: pd.DataFrame = None,
        method_text: str = "",
    ):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(layout)

        # Set slide background to white
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White

        # Gene title (Black, bold)
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.2), Inches(12.73), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = f"{gene} Expression"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black
        title_para.alignment = PP_ALIGN.LEFT

        # Method text (if provided) - Frost Grey, small font
        graph_top = Inches(0.9)
        if method_text:
            method_box = slide.shapes.add_textbox(
                Inches(0.3), Inches(0.7), Inches(12.73), Inches(0.4)
            )
            method_frame = method_box.text_frame
            method_frame.word_wrap = True
            method_para = method_frame.paragraphs[0]
            method_para.text = method_text
            method_para.font.size = Pt(11)
            method_para.font.color.rgb = RGBColor(0xC1, 0xC6, 0xC7)  # Frost Grey
            method_para.alignment = PP_ALIGN.LEFT
            # Adjust graph position if method text is present
            graph_top = Inches(1.2)

        fig_copy = go.Figure(fig)
        # Preserve the figure's bottom margin (may be large for angled labels)
        orig_margin = fig.layout.margin
        orig_b = orig_margin.b if orig_margin and orig_margin.b else 140
        fig_copy.update_layout(
            width=1000,
            height=550 + max(0, orig_b - 80),
            margin=dict(l=60, r=60, t=60, b=max(80, orig_b)),
            font=dict(size=14, family=PLOTLY_FONT_FAMILY, color="black"),
        )

        img_bytes = ReportGenerator._fig_to_image(fig_copy, format="png", scale=2)
        img_stream = io.BytesIO(img_bytes)

        left = Inches(0.5)
        top = graph_top
        width = Inches(9.5)
        slide.shapes.add_picture(img_stream, left, top, width=width)

        if gene_data is not None and not gene_data.empty:
            stats_box = slide.shapes.add_textbox(
                Inches(10.2), Inches(1.0), Inches(2.8), Inches(5.5)
            )
            stats_frame = stats_box.text_frame
            stats_frame.word_wrap = True

            header_para = stats_frame.paragraphs[0]
            header_para.text = "Statistics"
            header_para.font.size = Pt(14)
            header_para.font.bold = True

            for _, row in gene_data.iterrows():
                cond = row.get("Condition", "N/A")
                fc = row.get("Fold_Change", row.get("Relative_Expression", 0))
                pval = row.get("p_value", float("nan"))
                sig = row.get("significance", "")

                para = stats_frame.add_paragraph()
                para.text = f"{cond[:12]}"
                para.font.size = Pt(10)
                para.font.bold = True

                para2 = stats_frame.add_paragraph()
                fc_str = f"FC: {fc:.2f}" if pd.notna(fc) else "FC: N/A"
                p_str = f"p={pval:.3f}" if pd.notna(pval) else ""
                para2.text = f"  {fc_str} {sig}"
                para2.font.size = Pt(9)

        # Add Cosmax Red accent bar at bottom (full width, 0.5" height)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(0xEA, 0x1D, 0x22)  # Cosmax Red
        accent_bar.line.fill.background()  # No border

    @staticmethod
    def _add_dual_gene_slide(prs, layout, gene_pairs: list, processed_data: dict):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(layout)

        positions = [
            (Inches(0.3), Inches(0.8), Inches(6.2)),
            (Inches(6.8), Inches(0.8), Inches(6.2)),
        ]

        for idx, (gene, fig) in enumerate(gene_pairs):
            if idx >= 2:
                break

            left, top, width = positions[idx]

            title_box = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.5))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = f"{gene}"
            title_para.font.size = Pt(20)
            title_para.font.bold = True

            fig_copy = go.Figure(fig)
            fig_copy.update_layout(
                width=600,
                height=450,
                margin=dict(l=50, r=30, t=40, b=60),
                font=dict(size=11, family=PLOTLY_FONT_FAMILY, color="black"),
            )

            img_bytes = ReportGenerator._fig_to_image(fig_copy, format="png", scale=2)
            img_stream = io.BytesIO(img_bytes)

            slide.shapes.add_picture(img_stream, left, top, width=width)

    @staticmethod
    def _add_grid_slide(prs, layout, graphs: dict):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.1), Inches(12.73), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Gene Expression Overview"
        title_para.font.size = Pt(24)
        title_para.font.bold = True

        gene_list = list(graphs.items())
        n_genes = len(gene_list)

        if n_genes <= 2:
            cols, rows = 2, 1
        elif n_genes <= 4:
            cols, rows = 2, 2
        elif n_genes <= 6:
            cols, rows = 3, 2
        else:
            cols, rows = 4, 2

        cell_width = 12.5 / cols
        cell_height = 6.5 / rows

        for idx, (gene, fig) in enumerate(gene_list[: cols * rows]):
            col_idx = idx % cols
            row_idx = idx // cols

            left = Inches(0.4 + col_idx * cell_width)
            top = Inches(0.7 + row_idx * cell_height)

            fig_copy = go.Figure(fig)
            fig_copy.update_layout(
                width=350,
                height=280,
                margin=dict(l=40, r=20, t=35, b=40),
                font=dict(size=9, family=PLOTLY_FONT_FAMILY, color="black"),
                title=dict(text=gene, font=dict(size=12, family=PLOTLY_FONT_FAMILY, color="black")),
            )

            img_bytes = ReportGenerator._fig_to_image(fig_copy, format="png", scale=2)
            img_stream = io.BytesIO(img_bytes)

            slide.shapes.add_picture(
                img_stream, left, top, width=Inches(cell_width - 0.2)
            )

    @staticmethod
    def _add_summary_slide(prs, layout, processed_data: dict, analysis_params: dict):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.2), Inches(12.73), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Analysis Summary"
        title_para.font.size = Pt(28)
        title_para.font.bold = True

        all_results = (
            pd.concat(processed_data.values(), ignore_index=True)
            if processed_data
            else pd.DataFrame()
        )

        summary_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.0), Inches(6), Inches(5.5)
        )
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True

        header = summary_frame.paragraphs[0]
        header.text = "Experimental Parameters"
        header.font.size = Pt(16)
        header.font.bold = True

        params_text = [
            f"Efficacy Type: {analysis_params.get('Efficacy_Type', 'N/A')}",
            f"Housekeeping Gene: {analysis_params.get('Housekeeping_Gene', 'N/A')}",
            f"Reference Condition: {analysis_params.get('Reference_Sample', 'N/A')}",
            f"Comparison (*): {analysis_params.get('Compare_To', 'N/A')}",
        ]
        cmp2 = analysis_params.get("Compare_To_2")
        if cmp2:
            params_text.append(f"Comparison (#): {cmp2}")
        cmp3 = analysis_params.get("Compare_To_3")
        if cmp3:
            params_text.append(f"Comparison (\u2020): {cmp3}")
        params_text.extend([
            f"Genes Analyzed: {len(processed_data)}",
            f"Analysis Date: {analysis_params.get('Date', 'N/A')}",
        ])

        for text in params_text:
            para = summary_frame.add_paragraph()
            para.text = text
            para.font.size = Pt(12)

        if not all_results.empty:
            results_box = slide.shapes.add_textbox(
                Inches(7), Inches(1.0), Inches(5.8), Inches(5.5)
            )
            results_frame = results_box.text_frame
            results_frame.word_wrap = True

            header2 = results_frame.paragraphs[0]
            header2.text = "Key Findings"
            header2.font.size = Pt(16)
            header2.font.bold = True

            for gene, gene_df in processed_data.items():
                if gene_df.empty:
                    continue

                para = results_frame.add_paragraph()
                para.text = f"\n{gene}:"
                para.font.size = Pt(12)
                para.font.bold = True

                sig_results = gene_df[
                    gene_df.get("significance", pd.Series([""])).str.len() > 0
                ]
                if not sig_results.empty:
                    for _, row in sig_results.head(3).iterrows():
                        cond = row.get("Condition", "N/A")
                        fc = row.get("Fold_Change", row.get("Relative_Expression", 0))
                        sig = row.get("significance", "")
                        para2 = results_frame.add_paragraph()
                        para2.text = f"  {cond}: {fc:.2f}x {sig}"
                        para2.font.size = Pt(10)
                else:
                    para2 = results_frame.add_paragraph()
                    para2.text = "  No significant changes"
                    para2.font.size = Pt(10)


# ==================== PPT GENERATOR (NEW) ====================
class PPTGenerator:
    NAVY_BLUE = (0, 0, 0)  # Black (was navy #1B365D)
    WHITE = (255, 255, 255)

    @staticmethod
    def _get_color_rgb(rgb_tuple):
        from pptx.dml.color import RGBColor

        return RGBColor(*rgb_tuple)

    @staticmethod
    def create_title_slide(prs, analysis_params):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PPTGenerator._get_color_rgb(PPTGenerator.WHITE)

        # Title "유전자 발현 분석" (Gene Expression Analysis)
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.33), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "유전자 발현 분석 (Gene Expression Analysis)"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.name = "Malgun Gothic"
        p.alignment = PP_ALIGN.CENTER

        # Subtitle (Efficacy Type)
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.8), Inches(11.33), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        efficacy = analysis_params.get("Efficacy_Type", "Analysis")
        p.text = f"{efficacy} Efficacy Test"
        p.font.size = Pt(28)
        p.font.name = "Malgun Gothic"
        p.alignment = PP_ALIGN.CENTER

        # Date and Info
        info_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.0), Inches(11.33), Inches(1)
        )
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        date = analysis_params.get("Date", "")
        p.text = f"Date: {date}"
        p.font.size = Pt(18)
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

        # Logo Placeholder (Bottom Right)
        # X: 11.5, Y: 6.5 approx
        logo_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(6.5), Inches(1.5), Inches(0.5)
        )
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[Logo]"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.RIGHT

    @staticmethod
    def create_gene_slide(prs, gene, fig, gene_data, analysis_params, graph_settings=None, display_name=None):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        display_name = display_name or gene

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PPTGenerator._get_color_rgb(PPTGenerator.WHITE)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{display_name} Expression"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.name = "Malgun Gothic"

        # Navy Blue Line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = PPTGenerator._get_color_rgb(PPTGenerator.NAVY_BLUE)
        line.line.fill.background()

        # Graph — use per-gene dimensions from graph_settings if available
        gs = graph_settings or {}
        fb_w = max(int(gs.get(f"{gene}_figure_width", gs.get("figure_width", 28)) * CM_TO_PX), 800)
        fb_h = max(int(gs.get(f"{gene}_figure_height", gs.get("figure_height", 16)) * CM_TO_PX), 500)
        try:
            orig_m = fig.layout.margin
            extra_b = max(0, (orig_m.b if orig_m and orig_m.b else 0) - 120)
            img_bytes = ReportGenerator._fig_to_image(fig, format="png", scale=2, width=fb_w, height=fb_h + extra_b)
            image_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(
                image_stream,
                Inches(0.5),
                Inches(1.5),
                width=Inches(6.0),
                height=Inches(4.0),
            )
        except Exception as e:
            err_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(6.0), Inches(4.0)
            )
            err_box.text = f"Graph Error: {str(e)}"

        # Data Table (limit to 5 rows to prevent slide overflow — slide height is 7.5")
        if gene_data is not None and not gene_data.empty:
            max_table_rows = 5
            display_data = gene_data.head(max_table_rows) if len(gene_data) > max_table_rows else gene_data
            overflow_note = f" (+{len(gene_data) - max_table_rows} more in Excel)" if len(gene_data) > max_table_rows else ""
            rows = len(display_data) + 1
            cols = 4  # Sample, Condition, Fold Change, P-value
            table_height = Inches(0.3 * rows)
            table_shape = slide.shapes.add_table(
                rows, cols, Inches(0.5), Inches(5.8), Inches(6.0), table_height
            )
            table = table_shape.table

            # Headers
            headers = ["시료명 (Sample)", "Condition", "Fold Change", "P-value"]
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = h
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.bold = True

            # Data
            for i, row in enumerate(display_data.itertuples()):
                r = i + 1
                sample_val = getattr(row, "Original_Sample", getattr(row, "Sample", ""))
                cond_val = getattr(row, "Condition", "")
                fc_val = getattr(
                    row, "Fold_Change", getattr(row, "Relative_Expression", 0)
                )
                pval = getattr(row, "p_value", None)
                sig = getattr(row, "significance", "")

                table.cell(r, 0).text = str(sample_val)[:30]
                table.cell(r, 1).text = str(cond_val)[:30]
                table.cell(r, 2).text = f"{fc_val:.2f}" if pd.notna(fc_val) else "-"
                table.cell(r, 3).text = f"{pval:.4f} {sig}" if pd.notna(pval) else "-"

                for j in range(4):
                    table.cell(r, j).text_frame.paragraphs[0].font.size = Pt(9)

            if overflow_note:
                note_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(5.8 + 0.3 * rows), Inches(6.0), Inches(0.3)
                )
                note_box.text_frame.paragraphs[0].text = f"Showing {max_table_rows}/{len(gene_data)} conditions{overflow_note}"
                note_box.text_frame.paragraphs[0].font.size = Pt(8)
                note_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        # Metadata Box
        meta_box = slide.shapes.add_textbox(
            Inches(7.0), Inches(1.5), Inches(2.5), Inches(5.0)
        )
        tf = meta_box.text_frame
        tf.word_wrap = True

        def add_meta_line(text, bold=False):
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(11)
            p.font.bold = bold

        add_meta_line("Analysis Parameters", bold=True)
        add_meta_line(f"HK Gene: {analysis_params.get('Housekeeping_Gene', '-')}")
        add_meta_line(f"Ref Sample: {analysis_params.get('Reference_Sample', '-')}")
        add_meta_line(f"Compare (*): {analysis_params.get('Compare_To', '-')}")
        compare_2 = analysis_params.get("Compare_To_2")
        if compare_2:
            add_meta_line(f"Compare (#): {compare_2}")
        compare_3 = analysis_params.get("Compare_To_3")
        if compare_3:
            add_meta_line(f"Compare (\u2020): {compare_3}")
        add_meta_line("")
        add_meta_line("통계적 유의성 (Significance)", bold=True)
        add_meta_line("* p < 0.05   ** p < 0.01   *** p < 0.001")
        if compare_2:
            add_meta_line("# p < 0.05   ## p < 0.01   ### p < 0.001")
        if compare_3:
            add_meta_line("\u2020 p < 0.05   \u2020\u2020 p < 0.01   \u2020\u2020\u2020 p < 0.001")

        # Logo Placeholder
        logo_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(6.5), Inches(1.5), Inches(0.5)
        )
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[Logo]"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.RIGHT

    @staticmethod
    def _delete_slide(prs, index):
        slides_list = prs.slides._sldIdLst
        rId = slides_list[index].attrib[
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        ]
        prs.part.drop_rel(rId)
        slides_list.remove(slides_list[index])

    @staticmethod
    def _safe_copy_slide(prs, source_index):
        """Copy a slide by duplicating shapes via the spTree proxy.

        Uses python-pptx's spTree property to ensure shapes are properly
        recognized after copy. Copies relationships for images/media,
        then deep-copies each shape element from the source spTree.
        """
        import copy

        source = prs.slides[source_index]
        new_slide = prs.slides.add_slide(source.slide_layout)

        # Map source relationship IDs to new ones
        rid_map = {}
        for rid, rel in source.part.rels.items():
            rel_type = str(rel.reltype)
            if 'slideLayout' in rel_type:
                for new_rid, new_rel in new_slide.part.rels.items():
                    if 'slideLayout' in str(new_rel.reltype):
                        rid_map[rid] = new_rid
                        break
                continue
            if 'notesSlide' in rel_type:
                continue
            try:
                if rel.is_external:
                    new_rid = new_slide.part.rels.get_or_add_ext_rel(
                        rel.reltype, rel.target_ref
                    )
                else:
                    new_rid = new_slide.part.relate_to(
                        rel.target_part, rel.reltype
                    )
                rid_map[rid] = new_rid
            except Exception as e:
                import logging
                logging.warning(f"PPT slide copy: skipped relationship {rid}: {e}")
                # Surface to the user so a broken slide isn't silently shipped.
                try:
                    st.warning(
                        f"⚠️ PPT slide copy skipped a relationship "
                        f"({rel.reltype}). The exported slide may be missing an "
                        f"image or link. Details: {e}"
                    )
                except Exception:
                    pass

        source_spTree = source._element.spTree
        target_spTree = new_slide._element.spTree

        # Remove layout placeholder shapes from target
        shape_tags = ['}sp', '}pic', '}cxnSp', '}grpSp']
        for child in list(target_spTree):
            if any(child.tag.endswith(t) for t in shape_tags):
                target_spTree.remove(child)

        # Deep-copy shapes from source, remapping relationship IDs
        for child in source_spTree:
            if any(child.tag.endswith(t) for t in shape_tags):
                new_child = copy.deepcopy(child)
                for elem in new_child.iter():
                    for attr_key in list(elem.attrib.keys()):
                        val = elem.attrib[attr_key]
                        if val in rid_map:
                            elem.attrib[attr_key] = rid_map[val]
                target_spTree.append(new_child)

        return new_slide

    @staticmethod
    def _move_slide_to_end(prs, slide_index):
        slides_list = prs.slides._sldIdLst
        el = slides_list[slide_index]
        slides_list.remove(el)
        slides_list.append(el)

    @staticmethod
    def _populate_gene_slide(prs, slide, gene, fig, gene_data, analysis_params, graph_settings, display_name=None):
        """Populate a template-duplicated gene slide with gene name and graph image.

        `display_name` overrides `gene` for ALL user-visible text on the slide
        (title, "효능 평가" prefix). When the user renamed the gene in the Graphs
        tab, the rename is honored here so the slide matches the on-screen graph.
        """
        from pptx.util import Inches, Pt, Emu

        display_name = display_name or gene

        # Update text boxes on the template slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text

            # TextBox 17: "효능 평가" → "{display_name} 효능 평가"
            if "효능 평가" in full_text and "Results" not in full_text and "有" not in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "효능" in run.text:
                            run.text = run.text.replace("효능 평가", f"{display_name} 효능 평가")
                            break
                    break

            # TextBox 14: "Results: 효능 有/無" → update with gene result
            elif "Results" in full_text or "효능 有" in full_text:
                has_efficacy = False
                if gene_data is not None and not gene_data.empty:
                    sig_vals = gene_data.get("significance", pd.Series())
                    has_efficacy = sig_vals.str.len().gt(0).any() if not sig_vals.empty else False
                result_text = f"Results: 효능 {'有' if has_efficacy else '無'}"
                for para in shape.text_frame.paragraphs:
                    runs = list(para.runs)
                    for i, run in enumerate(runs):
                        if i == 0:
                            run.text = result_text
                        else:
                            run.text = ""
                    break

            # TextBox 2: experiment details (top-right area) — fill template fields
            elif (shape.left is not None and shape.top is not None
                  and shape.left > 3500000 and shape.top < 500000):
                efficacy_t = analysis_params.get("Efficacy_Type", "")
                eff_cfg = EFFICACY_CONFIG.get(efficacy_t, {})
                field_values = {
                    "Date: ": analysis_params.get("Date", "")[:10],
                    "Cell line: ": eff_cfg.get("cell", ""),
                    "Sample concentration: ": analysis_params.get("concentration", "1 ppm"),
                    "Positive control: ": eff_cfg.get("controls", {}).get("positive", ""),
                    "Inducer: ": eff_cfg.get("controls", {}).get("negative", ""),
                    "Treatment time: ": analysis_params.get("treatment_time", "24 h"),
                    "Test method:": " qPCR (DDCt)",
                }
                for para in shape.text_frame.paragraphs:
                    if not para.runs:
                        continue
                    label = para.runs[0].text
                    for field_label, value in field_values.items():
                        if label.strip().startswith(field_label.strip()):
                            para.runs[0].text = f"{field_label}{value}"
                            for extra_run in para.runs[1:]:
                                extra_run.text = ""
                            break

        # Add graph image
        gs = graph_settings or {}
        w_cm = gs.get(f"{gene}_figure_width", gs.get("figure_width", 28))
        h_cm = gs.get(f"{gene}_figure_height", gs.get("figure_height", 16))

        fig_copy = go.Figure(fig)
        fig_copy.update_layout(
            width=max(int(w_cm * CM_TO_PX), 800),
            height=max(int(h_cm * CM_TO_PX), 500),
        )

        try:
            img_bytes = ReportGenerator._fig_to_image(fig_copy, format="png", scale=2)
            img_stream = io.BytesIO(img_bytes)

            w_emu = int(w_cm * CM_TO_EMU)
            h_emu = int(h_cm * CM_TO_EMU)
            # Available area: ~1.0" to ~6.6" vertically, ~12" wide
            max_w_emu = int(11.5 * 914400)
            max_h_emu = int(5.2 * 914400)

            if w_emu > max_w_emu:
                scale_f = max_w_emu / w_emu
                w_emu = max_w_emu
                h_emu = int(h_emu * scale_f)
            if h_emu > max_h_emu:
                scale_f = max_h_emu / h_emu
                h_emu = max_h_emu
                w_emu = int(w_emu * scale_f)

            slide_w = int(prs.slide_width)
            left = max(0, int((slide_w - w_emu) / 2))
            top = int(1.0 * 914400)  # 1.0" below header

            slide.shapes.add_picture(
                img_stream, left, top, width=Emu(w_emu), height=Emu(h_emu)
            )
        except Exception as e:
            err_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(11.0), Inches(4.0)
            )
            err_box.text_frame.paragraphs[0].text = f"Graph Error: {str(e)}"
            err_box.text_frame.paragraphs[0].font.size = Pt(14)

    @staticmethod
    def generate_presentation(graphs, processed_data, analysis_params, graph_settings=None, gene_display_names=None):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Emu
        except ImportError:
            st.error("python-pptx not installed")
            return None

        import os

        script_dir = os.path.dirname(os.path.abspath(__file__))
        template_path = os.path.join(
            script_dir, "251215 효능평가 결과 TEMPLATE.pptx"
        )
        use_template = os.path.isfile(template_path)

        gs = graph_settings or {}
        gene_display_names = gene_display_names or {}
        graphs = graphs or {}

        # Source of truth is processed_data (the analysis output). Stale entries
        # in `graphs` from earlier runs are ignored; missing figures are
        # rendered lazily below so PPT works even if the user never opened
        # the Graphs tab.
        gene_list = [
            g for g in (processed_data.keys() if processed_data else [])
            if processed_data.get(g) is not None and not processed_data[g].empty
        ]
        n_genes = len(gene_list)

        def _get_fig(gene):
            """Return a Plotly figure for `gene`, rendering on demand if missing."""
            fig = graphs.get(gene)
            if fig is not None:
                return fig
            try:
                gene_data_local = processed_data.get(gene)
                if gene_data_local is None or gene_data_local.empty:
                    return None
                disp = gene_display_names.get(gene, gene)
                return GraphGenerator.create_gene_graph(
                    gene_data_local,
                    gene,
                    gs or {},
                    EFFICACY_CONFIG.get(analysis_params.get("Efficacy_Type", ""), {}),
                    sample_order=st.session_state.get("sample_order"),
                    display_gene_name=disp,
                    ref_condition=st.session_state.get("analysis_ref_condition"),
                )
            except Exception:
                return None

        if use_template and n_genes > 0:
            prs = Presentation(template_path)

            # Template slides: [0]=Title, [1]=Description, [2]=Gene template, [3]=Closing
            # Step 1: Update title slide — replace "소재" with efficacy type
            efficacy = analysis_params.get("Efficacy_Type", "소재")
            title_slide = prs.slides[0]
            for shape in title_slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "소재" in run.text:
                                run.text = run.text.replace("소재", efficacy)

            # Step 2: Copy gene template (index 2) BEFORE deleting to avoid
            # ZIP part-name conflicts (delete frees a name that add_slide reuses)
            for _ in range(n_genes - 1):
                PPTGenerator._safe_copy_slide(prs, 2)
            # Now: [0]=Title, [1]=Desc, [2]=Gene, [3]=Closing, [4..N+2]=GeneCopies

            # Step 3: Delete slide 1 (description overview)
            PPTGenerator._delete_slide(prs, 1)
            # Now: [0]=Title, [1]=Gene, [2]=Closing, [3..N+1]=GeneCopies

            # Step 4: Move closing slide (index 2) to end
            PPTGenerator._move_slide_to_end(prs, 2)
            # Now: [0]=Title, [1]=Gene, [2..N]=GeneCopies, [N+1]=Closing

            # Step 5: Populate each gene slide with gene name and graph
            for i, gene in enumerate(gene_list):
                slide = prs.slides[1 + i]
                fig = _get_fig(gene)
                gene_data = processed_data.get(gene)
                PPTGenerator._populate_gene_slide(
                    prs, slide, gene, fig, gene_data, analysis_params, gs,
                    display_name=gene_display_names.get(gene, gene),
                )
        else:
            # Fallback: no template available
            prs = Presentation()
            prs.slide_width = Emu(12192000)
            prs.slide_height = Emu(6858000)
            PPTGenerator.create_title_slide(prs, analysis_params)

            for gene in gene_list:
                fig = _get_fig(gene)
                gene_data = processed_data.get(gene)
                PPTGenerator.create_gene_slide(
                    prs, gene, fig, gene_data, analysis_params, graph_settings=gs,
                    display_name=gene_display_names.get(gene, gene),
                )

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()
