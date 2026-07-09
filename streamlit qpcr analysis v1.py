import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
from scipy import stats
import gc
import hashlib
import io
import warnings
import json
import re
import zipfile
from datetime import datetime
from typing import Dict, Tuple

from qpcr.constants import GRAPH_PRESETS, FIGURE_SIZE_PRESETS
from qpcr.export_utils import export_figure_to_bytes, build_zip
from qpcr.parser import QPCRParser
from qpcr.quality_control import QualityControl
from qpcr.graph import GraphGenerator
from qpcr.analysis import AnalysisEngine as _CoreAnalysisEngine

try:
    from streamlit_sortables import sort_items
except ImportError:
    sort_items = None


# ==================== UTILITY FUNCTIONS ====================
def natural_sort_key(sample_name):
    """Extract numbers from sample name for natural sorting (e.g., Sample2 < Sample10)"""
    parts = re.split(r"(\d+)", str(sample_name))
    return [int(part) if part.isdigit() else part.lower() for part in parts]


WORKFLOW_STEPS = [
    ("Upload", "upload"),
    ("QC", "qc"),
    ("Mapping", "mapping"),
    ("Analysis", "analysis"),
    ("Graphs", "graphs"),
    ("Export", "export"),
]


def render_step_indicator(current_step: str):
    """Render a compact horizontal step progress indicator."""
    has_data = st.session_state.get("data") is not None
    has_analysis = bool(st.session_state.get("processed_data"))
    has_graphs = bool(st.session_state.get("graphs"))
    mapping_done = st.session_state.get("mapping_finalized", False)

    qc_done = has_data and (
        sum(len(v) for v in st.session_state.get("excluded_wells", {}).values()) > 0
        or st.session_state.get("qc_reviewed", False)
    )

    step_status = {
        "upload": "done" if has_data else "pending",
        "qc": "done" if qc_done else "pending",
        "mapping": "done" if mapping_done else "pending",
        "analysis": "done" if has_analysis else "pending",
        "graphs": "done" if has_graphs else "pending",
        "export": "done" if has_graphs else "pending",
    }

    parts = []
    for i, (label, key) in enumerate(WORKFLOW_STEPS):
        if key == current_step:
            cls = "active"
        elif step_status[key] == "done":
            cls = "done"
        else:
            cls = ""
        parts.append(f'<span class="step-item {cls}"><span class="step-dot"></span>{label}</span>')
        if i < len(WORKFLOW_STEPS) - 1:
            parts.append('<span class="step-conn"></span>')

    st.markdown(f'<div class="step-progress">{"".join(parts)}</div>', unsafe_allow_html=True)


def render_sidebar_rail():
    """Persistent Cosmax-branded workflow rail + run context, shown on every tab.

    Progress is derived from real session state (not merely that a tab rendered),
    so the checkmarks are honest. The active tab is already highlighted by
    Streamlit's tab bar, so this rail shows completion + context rather than a
    redundant 'active' cue.
    """
    data = st.session_state.get("data")
    has_data = data is not None
    status = {
        "upload": has_data,
        "qc": st.session_state.get("qc_reviewed", False) or (
            has_data and sum(len(v) for v in st.session_state.get("excluded_wells", {}).values()) > 0
        ),
        "mapping": st.session_state.get("mapping_finalized", False),
        "analysis": bool(st.session_state.get("processed_data")),
        "graphs": bool(st.session_state.get("graphs")),
        "export": bool(st.session_state.get("graphs")),
    }
    with st.sidebar:
        st.markdown(
            '<div class="cosmax-brand">qPCR <span class="accent">Analysis</span> Suite</div>'
            '<div class="cosmax-sub">Cosmax · Efficacy Analytics</div>',
            unsafe_allow_html=True,
        )
        rail = "".join(
            f'<div class="rail-step {"done" if status.get(key) else ""}">'
            f'<span class="rdot"></span>{label}</div>'
            for label, key in WORKFLOW_STEPS
        )
        st.markdown(rail, unsafe_allow_html=True)
        if has_data:
            try:
                n_genes = int(data["Target"].nunique())
                n_samples = int(data["Sample"].nunique())
            except Exception:
                n_genes = n_samples = 0
            n_excl = sum(len(v) for v in st.session_state.get("excluded_wells", {}).values())
            eff = st.session_state.get("selected_efficacy") or "—"
            hk = st.session_state.get("hk_gene") or "—"
            st.markdown("<div style='height:14px'></div>", unsafe_allow_html=True)
            st.markdown(
                '<div class="rail-ctx">'
                f'<div class="row"><span class="k">Efficacy</span><span class="v">{eff}</span></div>'
                f'<div class="row"><span class="k">Ref gene</span><span class="v">{hk}</span></div>'
                f'<div class="row"><span class="k">Genes</span><span class="v">{n_genes}</span></div>'
                f'<div class="row"><span class="k">Samples</span><span class="v">{n_samples}</span></div>'
                f'<div class="row"><span class="k">Excluded wells</span><span class="v">{n_excl}</span></div>'
                '</div>',
                unsafe_allow_html=True,
            )


def build_provenance(*, efficacy, hk_gene, ref_condition, cmp_conditions, ttest_type,
                     excluded_wells, excluded_samples, n_genes, n_samples, timestamp,
                     app_version="qPCR Analysis Suite v3.1"):
    """Build a reproducibility/provenance record for an analysis run.

    Pure function (no Streamlit / global state) so it is unit-testable. Captures
    the parameters a reviewer needs to reproduce the run (MIQE-style): method,
    reference gene/condition, comparisons, stats test, and every excluded well
    with its gene/sample.
    """
    excl_list = []
    if isinstance(excluded_wells, dict):
        for key, wells in excluded_wells.items():
            if isinstance(key, (tuple, list)) and len(key) == 2:
                gene, sample = key
            else:
                gene, sample = "?", str(key)
            for w in sorted(wells):
                excl_list.append({"gene": str(gene), "sample": str(sample), "well": str(w)})
    excl_list.sort(key=lambda d: (d["gene"], d["sample"], d["well"]))
    return {
        "generated": timestamp,
        "software": app_version,
        "method": "Livak 2^-ddCt, single reference gene",
        "fdr_correction": "Benjamini-Hochberg",
        "efficacy_type": efficacy or None,
        "reference_gene": hk_gene or None,
        "reference_condition": ref_condition or None,
        "comparison_conditions": [c for c in (cmp_conditions or []) if c],
        "statistical_test": "Welch t-test" if ttest_type == "welch" else "Student t-test",
        "n_genes": int(n_genes),
        "n_samples": int(n_samples),
        "excluded_samples": sorted(excluded_samples) if excluded_samples else [],
        "excluded_wells_count": len(excl_list),
        "excluded_wells": excl_list,
    }


def format_provenance_text(prov: dict) -> str:
    """Human-readable multi-line summary of a provenance record (pure)."""
    cmp = ", ".join(prov.get("comparison_conditions") or []) or "—"
    excl_s = ", ".join(prov.get("excluded_samples") or []) or "none"
    lines = [
        "qPCR ANALYSIS PROVENANCE",
        f"Generated:            {prov.get('generated')}",
        f"Software:             {prov.get('software')}",
        f"Method:               {prov.get('method')}",
        f"FDR correction:       {prov.get('fdr_correction')}",
        f"Efficacy type:        {prov.get('efficacy_type')}",
        f"Reference gene:       {prov.get('reference_gene')}",
        f"Reference condition:  {prov.get('reference_condition')}",
        f"Comparison(s):        {cmp}",
        f"Statistical test:     {prov.get('statistical_test')}",
        f"Genes analysed:       {prov.get('n_genes')}",
        f"Samples:              {prov.get('n_samples')}",
        f"Excluded samples:     {excl_s}",
        f"Excluded wells:       {prov.get('excluded_wells_count')}",
    ]
    for e in prov.get("excluded_wells") or []:
        lines.append(f"    - {e['gene']} / {e['sample']} / well {e['well']}")
    return "\n".join(lines)


def _current_provenance():
    """Assemble a provenance record from current session state (UI helper)."""
    data = st.session_state.get("data")
    try:
        n_genes = int(data["Target"].nunique()) if data is not None else 0
        n_samples = int(data["Sample"].nunique()) if data is not None else 0
    except Exception:
        n_genes = n_samples = 0
    return build_provenance(
        efficacy=st.session_state.get("selected_efficacy"),
        hk_gene=st.session_state.get("hk_gene"),
        ref_condition=st.session_state.get("analysis_ref_condition"),
        cmp_conditions=[
            st.session_state.get("analysis_cmp_condition"),
            st.session_state.get("analysis_cmp_condition_2"),
            st.session_state.get("analysis_cmp_condition_3"),
        ],
        ttest_type=st.session_state.get("ttest_type", "welch"),
        excluded_wells=st.session_state.get("excluded_wells", {}),
        excluded_samples=st.session_state.get("excluded_samples", set()),
        n_genes=n_genes,
        n_samples=n_samples,
        timestamp=datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    )


# ==================== PAGE CONFIG ====================
st.set_page_config(
    page_title="qPCR Analysis Suite",
    page_icon="🧬",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ==================== GLOBAL THEME (Cosmax-branded) ====================
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

    :root {
        --cosmax-red: #EA1D22;
        --cosmax-red-dark: #C41419;
        --ink: #1d1d1f;
        --ink-soft: #6e6e73;
        --ink-faint: #86868b;
        --line: #e5e5e7;
        --line-strong: #d2d2d7;
        --surface: #ffffff;
        --surface-alt: #f5f5f7;
        --lab-white: #f3f0ed;
        --frost: #c1c6c7;
        --radius: 12px;
        --radius-sm: 8px;
    }

    /* Base typography — Inter for latin, Korean fallbacks for efficacy labels */
    html, body, [class*="css"] {
        font-family: 'Inter', 'Pretendard', 'Noto Sans KR', -apple-system,
                     BlinkMacSystemFont, 'SF Pro Display', 'Segoe UI', sans-serif;
    }

    /* Clean headers with a subtle red tick under h1 */
    h1, h2, h3 { font-weight: 600; letter-spacing: -0.02em; color: var(--ink); }
    h1 { font-size: 1.9rem; }
    h2 { font-size: 1.4rem; }
    h3 { font-size: 1.12rem; }

    /* App title accent bar */
    .block-container h1:first-of-type {
        border-bottom: 3px solid var(--cosmax-red);
        padding-bottom: 8px;
        display: inline-block;
    }

    /* Metric cards — lab-white surface */
    [data-testid="stMetric"] {
        background: var(--lab-white);
        border: 1px solid var(--line);
        border-radius: var(--radius);
        padding: 16px 18px;
    }
    [data-testid="stMetricValue"] { font-weight: 600; color: var(--ink); }
    [data-testid="stMetricLabel"] { color: var(--ink-faint); font-weight: 500; font-size: 0.85rem; }

    /* Tabs — active tab gets the Cosmax-red underline */
    .stTabs [data-baseweb="tab-list"] {
        gap: 0;
        border-bottom: 1px solid var(--line);
        background: transparent;
    }
    .stTabs [data-baseweb="tab"] {
        padding: 12px 22px;
        font-weight: 500;
        font-size: 0.9rem;
        color: var(--ink-faint);
        border-bottom: 2px solid transparent;
        background: transparent;
    }
    .stTabs [aria-selected="true"] {
        color: var(--ink);
        border-bottom: 2px solid var(--cosmax-red);
        background: transparent;
    }

    /* Buttons — secondary (outline) default, primary is Cosmax red */
    .stButton > button {
        border-radius: var(--radius-sm);
        font-weight: 500;
        font-size: 0.85rem;
        border: 1px solid var(--line-strong);
        background: var(--surface);
        color: var(--ink);
        transition: background .15s ease, border-color .15s ease;
        padding: 8px 16px;
    }
    .stButton > button:hover {
        background: var(--surface-alt);
        border-color: var(--frost);
    }
    .stButton > button[kind="primary"] {
        background: var(--cosmax-red);
        color: #ffffff;
        border-color: var(--cosmax-red);
    }
    .stButton > button[kind="primary"]:hover {
        background: var(--cosmax-red-dark);
        border-color: var(--cosmax-red-dark);
    }
    /* Download buttons read as primary-adjacent actions */
    .stDownloadButton > button {
        border-radius: var(--radius-sm);
        font-weight: 500;
        border: 1px solid var(--cosmax-red);
        color: var(--cosmax-red);
        background: var(--surface);
        transition: background .15s ease;
    }
    .stDownloadButton > button:hover { background: #fdecec; }

    /* Expanders */
    .streamlit-expanderHeader {
        font-weight: 500; font-size: 0.95rem; color: var(--ink);
        background: transparent; border: none;
    }
    .streamlit-expanderContent {
        border: 1px solid var(--line);
        border-radius: 0 0 var(--radius-sm) var(--radius-sm);
        padding: 16px;
    }

    /* Dataframes */
    [data-testid="stDataFrame"] { border: 1px solid var(--line); border-radius: var(--radius-sm); }

    /* Input labels */
    [data-testid="stSelectbox"] label,
    [data-testid="stNumberInput"] label,
    [data-testid="stTextInput"] label { color: var(--ink-faint); font-weight: 500; font-size: 0.85rem; }

    /* Plotly chart container */
    [data-testid="stPlotlyChart"] { border-radius: var(--radius); overflow: hidden; }

    /* Form submit — Cosmax red */
    .stForm [data-testid="stFormSubmitButton"] > button {
        background: var(--cosmax-red); color: #fff; border: none;
        border-radius: var(--radius-sm); font-weight: 600;
    }
    .stForm [data-testid="stFormSubmitButton"] > button:hover { background: var(--cosmax-red-dark); }

    /* Alerts, dividers, captions */
    .stAlert { border-radius: var(--radius-sm); border: none; }
    hr { border: none; border-top: 1px solid var(--line); margin: 22px 0; }
    .stCaption { color: var(--ink-faint); }

    /* ---- Sidebar workflow rail ---- */
    [data-testid="stSidebar"] {
        background: var(--lab-white);
        border-right: 1px solid var(--line);
    }
    [data-testid="stSidebar"] .block-container { padding-top: 1.2rem; }
    .cosmax-brand {
        font-weight: 700; font-size: 1.05rem; color: var(--ink);
        letter-spacing: -0.01em; margin-bottom: 2px;
    }
    .cosmax-brand .accent { color: var(--cosmax-red); }
    .cosmax-sub { font-size: 0.72rem; color: var(--ink-faint); margin-bottom: 14px;
                  text-transform: uppercase; letter-spacing: 0.06em; }
    .rail-step {
        display: flex; align-items: center; gap: 10px;
        padding: 7px 10px; border-radius: var(--radius-sm);
        font-size: 0.85rem; color: var(--ink-faint); font-weight: 500;
    }
    .rail-step .rdot {
        width: 9px; height: 9px; border-radius: 50%;
        background: var(--line-strong); flex-shrink: 0;
    }
    .rail-step.done { color: var(--ink-soft); }
    .rail-step.done .rdot { background: var(--cosmax-red); }
    .rail-ctx {
        background: var(--surface); border: 1px solid var(--line);
        border-radius: var(--radius-sm); padding: 10px 12px; font-size: 0.8rem;
    }
    .rail-ctx .row { display: flex; justify-content: space-between; padding: 3px 0; }
    .rail-ctx .k { color: var(--ink-faint); }
    .rail-ctx .v { color: var(--ink); font-weight: 600; }

    /* Step progress indicator (top-of-tab breadcrumb) */
    .step-progress { display: flex; align-items: center; gap: 0; padding: 6px 0 12px 0; margin-bottom: 4px; }
    .step-progress .step-item { display: flex; align-items: center; gap: 6px; font-size: 0.75rem;
                                color: var(--line-strong); font-weight: 500; white-space: nowrap; }
    .step-progress .step-item.active { color: var(--cosmax-red); font-weight: 600; }
    .step-progress .step-item.done { color: var(--ink-faint); }
    .step-progress .step-dot { width: 7px; height: 7px; border-radius: 50%; background: var(--line); flex-shrink: 0; }
    .step-progress .step-item.active .step-dot { background: var(--cosmax-red);
                                                 box-shadow: 0 0 0 3px rgba(234,29,34,0.15); }
    .step-progress .step-item.done .step-dot { background: var(--ink-faint); }
    .step-progress .step-conn { width: 20px; height: 1px; background: var(--line); margin: 0 4px; flex-shrink: 0; }

    /* Gene pill selector container */
    .gene-pill-container { background: var(--lab-white); border: 1px solid var(--line);
                           border-radius: var(--radius); padding: 12px 16px 8px 16px; margin-bottom: 12px; }
    .gene-pill-container .label { font-size: 0.78rem; color: var(--ink-faint); font-weight: 500; margin-bottom: 8px; }

    /* Accessibility: branded focus ring + honour reduced-motion */
    button:focus-visible, [role="tab"]:focus-visible,
    input:focus-visible, select:focus-visible, textarea:focus-visible {
        outline: 2px solid var(--cosmax-red);
        outline-offset: 2px;
    }
    @media (prefers-reduced-motion: reduce) {
        * { transition: none !important; animation: none !important; }
    }
</style>
""", unsafe_allow_html=True)

# ==================== CONSTANTS ====================
# COSMAX brand colors for PPT redesign
COSMAX_RED = "#EA1D22"  # Primary accent, emphasis
COSMAX_BLACK = "#000000"  # Main text
COSMAX_WHITE = "#FFFFFF"  # Background
COSMAX_LAB_WHITE = "#F3F0ED"  # Secondary background (off-white)
COSMAX_FROST_GREY = "#C1C6C7"  # Secondary elements, table headers
COSMAX_CREAM = "#D4CEC1"  # Secondary data series, neutral accents

# Font family with CJK (Korean) support for Plotly image export (kaleido)
# Cross-platform fallback: Linux → Windows → macOS → generic
_DEFAULT_CJK_FONTS = [
    "Noto Sans CJK KR", "NanumGothic", "Malgun Gothic",
    "Apple SD Gothic Neo", "AppleGothic",
]
_FALLBACK_FONTS = ["Arial", "sans-serif"]

# FIX-12: Validate CJK font availability at startup with graceful fallback
def _detect_available_fonts():
    """Check which CJK fonts are actually installed on this system."""
    try:
        from matplotlib import font_manager
        system_fonts = {f.name for f in font_manager.fontManager.ttflist}
        available_cjk = [f for f in _DEFAULT_CJK_FONTS if f in system_fonts]
        if not available_cjk:
            # No CJK fonts found — warn once via Streamlit toast
            # (deferred to first render since st isn't ready at import time)
            pass
        return available_cjk + _FALLBACK_FONTS
    except ImportError:
        # matplotlib not available — use full fallback list
        return _DEFAULT_CJK_FONTS + _FALLBACK_FONTS

PLOTLY_FONT_FAMILY = ", ".join(_detect_available_fonts())

CM_TO_PX = 37.7953
CM_TO_EMU = 360000


# ==================== SESSION STATE INIT ====================
# Helper function for well exclusion management
def get_well_exclusion_key(gene: str, sample: str) -> tuple:
    """Generate key for per-gene-sample well exclusions."""
    return (gene, sample)


def is_well_excluded(well: str, gene: str, sample: str) -> bool:
    """Check if a well is excluded for a specific gene-sample combination."""
    key = get_well_exclusion_key(gene, sample)
    if key in st.session_state.excluded_wells:
        return well in st.session_state.excluded_wells[key]
    return False


def exclude_well(well: str, gene: str, sample: str) -> None:
    """Exclude a well for a specific gene-sample combination."""
    key = get_well_exclusion_key(gene, sample)
    if key not in st.session_state.excluded_wells:
        st.session_state.excluded_wells[key] = set()
    st.session_state.excluded_wells[key].add(well)


def include_well(well: str, gene: str, sample: str) -> None:
    """Include a well (remove from exclusion) for a specific gene-sample combination."""
    key = get_well_exclusion_key(gene, sample)
    if key in st.session_state.excluded_wells:
        st.session_state.excluded_wells[key].discard(well)
        # FIX-18: Clean up orphaned empty sets to prevent memory leak
        if not st.session_state.excluded_wells[key]:
            del st.session_state.excluded_wells[key]


def get_excluded_wells_for_analysis(gene: str, sample: str) -> set:
    """Get excluded wells for a specific gene-sample combination."""
    key = get_well_exclusion_key(gene, sample)
    if key in st.session_state.excluded_wells:
        return st.session_state.excluded_wells[key]
    return set()


def get_all_excluded_wells() -> set:
    """Get all excluded wells across all gene-sample combinations as a set.

    Compatibility function for QualityControl methods that expect a set parameter.
    """
    excluded = set()
    for wells_set in st.session_state.excluded_wells.values():
        excluded.update(wells_set)
    return excluded


for key in [
    "data",
    "processed_data",
    "sample_mapping",
    "analysis_templates",
    "graphs",
    "excluded_wells",
    "excluded_samples",
    "selected_efficacy",
    "hk_gene",
]:
    if key not in st.session_state:
        st.session_state[key] = (
            {}
            if key
            in ["sample_mapping", "analysis_templates", "graphs"]
            else (
                dict()
                if key == "excluded_wells"
                else set()
                if "excluded" in key
                else None
            )
        )

if "mapping_finalized" not in st.session_state:
    st.session_state.mapping_finalized = False

if "analysis_stale" not in st.session_state:
    st.session_state.analysis_stale = False

if "gene_display_names" not in st.session_state:
    st.session_state.gene_display_names = {}

if "experiment_desc" not in st.session_state:
    st.session_state.experiment_desc = {"concentration": "1 ppm", "treatment_time": "24 h"}

# ==================== EFFICACY DATABASE ====================
EFFICACY_CONFIG = {
    "탄력": {
        "genes": ["COL1A1", "ELN", "FBN-1", "FBN1"],
        "cell": "HS68 fibroblast",
        "controls": {
            "negative": "Non-treated",
            "positive": "TGFb",
            "compare_to": "negative",
        },
        "description": "Elasticity - Non-treated vs TGFb (positive) vs Treatments",
    },
    "항노화": {
        "genes": ["COL1A1", "COL1", "MMP-1", "MMP1"],
        "cell": "HS68 fibroblast",
        "controls": {
            "baseline": "Non-treated (No UV)",
            "negative": "UVB only",
            "positive": "UVB+TGFb",
            "compare_to": "negative",
        },
        "description": "Anti-aging - COL1↑ (recovery), MMP1↓ (inhibition) after UVB damage",
        "expected_direction": {
            "COL1A1": "up",
            "COL1": "up",
            "MMP-1": "down",
            "MMP1": "down",
        },
    },
    "보습": {
        "genes": ["AQP3", "HAS3"],
        "cell": "HaCaT keratinocyte",
        "controls": {
            "negative": "Non-treated",
            "positive": "Retinoic acid",
            "compare_to": "negative",
        },
        "description": "Hydration - Non-treated vs Retinoic acid (positive) vs Treatments",
    },
    "장벽": {
        "genes": ["FLG", "CLDN", "IVL"],
        "cell": "HaCaT keratinocyte",
        "controls": {
            "negative": "Non-treated",
            "positive": "Retinoic acid",
            "compare_to": "negative",
        },
        "description": "Barrier function - Non-treated vs Retinoic acid (positive) vs Treatments",
    },
    "표피증식": {
        "genes": ["KI67", "PCNA"],
        "cell": "HaCaT keratinocyte",
        "controls": {
            "negative": "Non-treated",
            "positive": "TGFb or FBS",
            "compare_to": "negative",
        },
        "description": "Proliferation - Non-treated vs TGFb/FBS (positive) vs Treatments",
    },
    "멜라닌억제": {
        "genes": ["MITF", "TYR", "Melanin"],
        "cell": "B16F10 melanocyte",
        "controls": {
            "baseline": "Non-treated",
            "negative": "α-MSH only",
            "positive": "α-MSH+Arbutin",
            "compare_to": "negative",
        },
        "description": "Melanin inhibition - α-MSH induced vs α-MSH+Arbutin (positive) vs α-MSH+Treatments",
        "expected_direction": {"MITF": "down", "TYR": "down", "Melanin": "down"},
    },
    "진정": {
        "genes": ["IL1B", "IL-1β", "IL6", "TNFA", "TNFα"],
        "cell": "HaCaT keratinocyte",
        "controls": {
            "baseline": "Non-treated",
            "negative": "IL4+PolyIC (Inflammation)",
            "positive": "Inflammation+Dexamethasone",
            "compare_to": "negative",
        },
        "description": "Anti-inflammation - Reduce IL1β/IL6/TNFα (all should decrease)",
        "expected_direction": {
            "IL1B": "down",
            "IL-1β": "down",
            "IL6": "down",
            "TNFA": "down",
            "TNFα": "down",
        },
    },
    "지질억제": {
        "genes": ["SREBPA", "SREBPa", "SREBPC", "SREBPc", "PPARY", "PPARy"],
        "cell": "SZ95 sebocyte",
        "controls": {
            "baseline": "Non-treated",
            "negative": "IGF only",
            "positive": "IGF+Reference inhibitor",
            "compare_to": "negative",
        },
        "description": "Sebum inhibition - IGF induced vs IGF+Treatments",
        "expected_direction": {
            "SREBPA": "down",
            "SREBPa": "down",
            "SREBPC": "down",
            "SREBPc": "down",
            "PPARY": "down",
            "PPARy": "down",
        },
    },
    "냉감": {
        "genes": ["TRPM8", "CIRBP"],
        "cell": "HaCaT keratinocyte",
        "controls": {
            "negative": "Non-treated",
            "positive": "Menthol",
            "compare_to": "negative",
        },
        "description": "Cooling effect - Non-treated vs Menthol (positive) vs Treatments",
    },
    "모근 강화": {
        "genes": ["VEGF", "COL17A1", "HGF", "FGF7", "FLG"],
        "cell": "HFDPC / HaCaT",
        "controls": {
            "negative": "Non-treated",
            "positive": "Minoxidil",
            "compare_to": "negative",
        },
        "description": "Hair root strengthening - VEGF, COL17A1, HGF, FGF7, FLG expression",
    },
}


# ==================== ANALYSIS CONSTANTS ====================
class AnalysisConstants:
    MIN_REPLICATES_FOR_STATS = 2
    RECOMMENDED_REPLICATES = 3
    CT_UNDETERMINED_THRESHOLD = 40.0
    CT_HIGH_WARNING = 35.0
    CT_LOW_WARNING = 10.0
    CV_WARNING_THRESHOLD = 0.05
    CV_ERROR_THRESHOLD = 0.10
    HK_DEVIATION_WARNING = 1.0
    HK_DEVIATION_ERROR = 2.0
    P_VALUE_THRESHOLDS = {"*": 0.05, "**": 0.01, "***": 0.001}


# ==================== QC GRID STATE MANAGEMENT ====================
def get_grid_cell_key(gene: str, sample: str) -> str:
    """Generate unique key for grid cell identification.

    Creates a consistent, unique string key from gene and sample names.
    Used internally to identify grid cells and ensure consistent state management.

    Args:
        gene: Target gene name (e.g., "GENE1")
        sample: Sample name (e.g., "Sample_A")

    Returns:
        Unique string key (e.g., "GENE1::Sample_A")
    """
    return f"{gene}::{sample}"


def get_selected_cell(session_state) -> tuple[str, str] | None:
    """Retrieve currently selected (gene, sample) or None.

    Returns the currently selected cell from session state, or None if no cell
    is selected. The return value is a tuple of (gene, sample) strings.

    Args:
        session_state: Streamlit session state object

    Returns:
        Tuple of (gene, sample) if a cell is selected, None otherwise
    """
    if "qc_grid_selected_cell" not in session_state:
        return None

    cell = session_state["qc_grid_selected_cell"]
    if cell is None:
        return None

    return (cell["gene"], cell["sample"])


def set_selected_cell(session_state, gene: str, sample: str) -> None:
    """Store the selected cell in session state.

    Sets the currently selected cell. Only one cell can be selected at a time;
    calling this function overwrites any previous selection.

    Args:
        session_state: Streamlit session state object
        gene: Target gene name
        sample: Sample name
    """
    session_state["qc_grid_selected_cell"] = {"gene": gene, "sample": sample}


def clear_selected_cell(session_state) -> None:
    """Remove the current selection, returning the grid to an unselected state.

    Clears the currently selected cell from session state. Used when user clicks
    away or when filters change.

    Args:
        session_state: Streamlit session state object
    """
    session_state["qc_grid_selected_cell"] = None


def is_cell_selected(session_state, gene: str, sample: str) -> bool:
    """Check if the given gene and sample match the currently selected cell.

    Returns True if the specified cell is currently selected, False otherwise.
    Used to highlight the selected cell in the grid UI.

    Args:
        session_state: Streamlit session state object
        gene: Target gene name to check
        sample: Sample name to check

    Returns:
        True if this cell is selected, False otherwise
    """
    selected = get_selected_cell(session_state)
    if selected is None:
        return False

    return selected[0] == gene and selected[1] == sample


# ==================== ANALYSIS ENGINE ====================
class AnalysisEngine(_CoreAnalysisEngine):
    """App-facing engine: inherits pure compute methods from
    qpcr.analysis.AnalysisEngine; keeps the Streamlit-coupled orchestration."""

    @staticmethod
    def run_full_analysis(
        ref_sample_key: str, compare_sample_key: str,
        compare_sample_key_2: str = None, compare_sample_key_3: str = None,
    ):
        """
        Run ΔΔCt + statistical analysis and store results in st.session_state.
        Produces st.session_state.processed_data = {gene: DataFrame}.
        """
        try:
            data = st.session_state.get("data")
            mapping = st.session_state.get("sample_mapping", {})
            hk_gene = st.session_state.get("hk_gene")

            if data is None:
                st.error("No raw data loaded. Upload CSV files in the **Upload** tab first.")
                return False
            if not mapping:
                st.error("Sample mapping not found. Go to the **Mapping** tab to assign conditions.")
                return False
            if not hk_gene:
                st.error("Housekeeping gene not selected. Select one in the **Upload** tab.")
                return False

            ct_values = data["CT"]
            high_ct_count = (ct_values > 35).sum()
            low_ct_count = (ct_values < 10).sum()
            single_rep_samples = []

            for sample in data["Sample"].unique():
                for target in data["Target"].unique():
                    sample_target_data = data[
                        (data["Sample"] == sample) & (data["Target"] == target)
                    ]
                    if len(sample_target_data) == 1:
                        single_rep_samples.append(f"{sample}/{target}")

            if high_ct_count > 0:
                st.warning(
                    f"⚠️ {high_ct_count} CT values > 35 detected (possible low expression or failed reactions)"
                )
            if low_ct_count > 0:
                st.warning(
                    f"⚠️ {low_ct_count} CT values < 10 detected (unusually high expression - verify data)"
                )
            if single_rep_samples:
                st.info(
                    f"ℹ️ {len(single_rep_samples)} sample/target combinations have only 1 replicate (limited statistical power)"
                )

            ref_condition = mapping.get(ref_sample_key, {}).get(
                "condition", ref_sample_key
            )
            cmp_condition = mapping.get(compare_sample_key, {}).get(
                "condition", compare_sample_key
            )
            cmp_condition_2 = None
            if compare_sample_key_2:
                cmp_condition_2 = mapping.get(compare_sample_key_2, {}).get(
                    "condition", compare_sample_key_2
                )
            cmp_condition_3 = None
            if compare_sample_key_3:
                cmp_condition_3 = mapping.get(compare_sample_key_3, {}).get(
                    "condition", compare_sample_key_3
                )

            st.session_state.analysis_ref_condition = ref_condition
            st.session_state.analysis_cmp_condition = cmp_condition
            st.session_state.analysis_cmp_condition_2 = cmp_condition_2
            st.session_state.analysis_cmp_condition_3 = cmp_condition_3

            msg = f"Running full analysis using reference '{ref_condition}' and comparison '{cmp_condition}'"
            if cmp_condition_2:
                msg += f" + secondary comparison '{cmp_condition_2}'"
            if cmp_condition_3:
                msg += f" + tertiary comparison '{cmp_condition_3}'"

            with st.spinner(msg + "..."):
                # --- ΔΔCt calculation ---
                processed_df = AnalysisEngine.calculate_ddct(
                    data,
                    hk_gene,
                    ref_condition,
                    st.session_state.get("excluded_wells", {}),
                    st.session_state.get("excluded_samples", set()),
                    mapping,
                )

                if processed_df is None or processed_df.empty:
                    st.warning(
                        "⚠️ No ΔΔCt results produced. Check mapping and housekeeping gene."
                    )
                    return False

                # FIX-06: Display warnings for skipped genes/conditions
                _skipped = processed_df.attrs.get("_skipped_warnings", [])
                if _skipped:
                    st.warning(
                        f"⚠️ {len(_skipped)} condition(s) skipped due to all wells being excluded:\n"
                        + "\n".join(f"  • {w}" for w in _skipped)
                    )

                # --- Statistical test ---
                ttest_type = st.session_state.get("ttest_type", "welch")
                # Pre-filter raw_data by excluded samples so statistics
                # never sees data the user chose to exclude
                stats_raw = data[
                    ~data["Sample"].isin(st.session_state.get("excluded_samples", set()))
                ].copy()
                try:
                    processed_with_stats = AnalysisEngine.calculate_statistics(
                        processed_df,
                        cmp_condition,
                        cmp_condition_2,
                        cmp_condition_3,
                        raw_data=stats_raw,
                        hk_gene=hk_gene,
                        sample_mapping=mapping,
                        ttest_type=ttest_type,
                        excluded_wells=st.session_state.get("excluded_wells", {}),
                    )
                except TypeError as te:
                    st.warning(f"Statistics fallback triggered (secondary/tertiary comparisons dropped): {te}")
                    processed_with_stats = AnalysisEngine.calculate_statistics(
                        processed_df, cmp_condition, ttest_type=ttest_type,
                        raw_data=stats_raw, hk_gene=hk_gene,
                        sample_mapping=mapping,
                        excluded_wells=st.session_state.get("excluded_wells", {}),
                    )

                # FIX-16: Display warning when one-sample t-test was used
                _onesamp = processed_with_stats.attrs.get("_onesamp_warnings", [])
                if _onesamp:
                    st.info(
                        f"ℹ️ One-sample t-test used for {len(_onesamp)} comparison(s) "
                        f"(one group has n=1 replicate): "
                        + ", ".join(_onesamp[:5])
                        + (f" ... and {len(_onesamp) - 5} more" if len(_onesamp) > 5 else "")
                    )

                # --- Organize data for graphs ---
                gene_dict = {}
                if "Target" in processed_with_stats.columns:
                    for gene in processed_with_stats["Target"].unique():
                        gene_df = processed_with_stats[
                            processed_with_stats["Target"] == gene
                        ].copy()
                        gene_dict[gene] = gene_df.reset_index(drop=True)
                else:
                    gene_dict = {"results": processed_with_stats.reset_index(drop=True)}

                st.session_state.processed_data = gene_dict

                # Flush stale bar_colors_per_sample keys that no longer match any condition
                _valid_keys = set()
                for _g, _gdf in gene_dict.items():
                    for _c in _gdf["Condition"].unique():
                        _valid_keys.add(f"{_g}_{_c}")
                _bcs = st.session_state.graph_settings.get("bar_colors_per_sample", {})
                _stale = [k for k in _bcs if k not in _valid_keys]
                for k in _stale:
                    del _bcs[k]

                # Prune any per-gene state that no longer maps to a real gene in
                # the fresh analysis. Without this, an old gene's cached graph
                # or rename can carry over into PPT/Excel exports of the new run.
                _new_genes = set(gene_dict.keys())
                _graphs = st.session_state.get("graphs", {})
                for _stale_gene in [g for g in _graphs if g not in _new_genes]:
                    _graphs.pop(_stale_gene, None)
                _disp_map = st.session_state.get("gene_display_names", {})
                for _stale_gene in [g for g in _disp_map if g not in _new_genes]:
                    _disp_map.pop(_stale_gene, None)
                # Drop the cached PPT bytes so the next download reflects the
                # new analysis (and the new rename map) instead of the prior run.
                st.session_state.pop("_ppt_export", None)

                # Store exclusion snapshot for auto-rerun detection
                st.session_state['_exclusion_snapshot'] = {
                    'excluded_wells': {str(k): sorted(v) for k, v in st.session_state.get('excluded_wells', {}).items()},
                    'excluded_samples': sorted(st.session_state.get('excluded_samples', set())),
                    'ttest_type': st.session_state.get('ttest_type', 'welch')
                }

            st.success(
                "✅ Full analysis complete. Go to the Graphs tab to visualize results."
            )
            return True

        except Exception as e:
            import traceback
            st.error(f"Analysis failed: {e}")
            # Keep the full traceback available (collapsed) so production
            # failures can actually be diagnosed instead of a one-line message.
            with st.expander("Show technical details"):
                st.code(traceback.format_exc())
            return False


# ==================== REPORT GENERATOR (PPT) ====================
class ReportGenerator:
    SLIDE_WIDTH_INCHES = 13.333
    SLIDE_HEIGHT_INCHES = 7.5

    @staticmethod
    def _fig_to_image(fig: go.Figure, format: str = "png", scale: int = 2,
                      width: int = None, height: int = None) -> bytes:
        """Convert a Plotly figure to image bytes.

        Delegates to the shared, environment-hardened renderer in
        ``qpcr.export_utils`` (browser auto-detection + Chrome-download fallback).
        """
        return export_figure_to_bytes(fig, fmt=format, scale=scale, width=width, height=height)

    @staticmethod
    def create_presentation(
        graphs: Dict[str, go.Figure],
        processed_data: Dict[str, pd.DataFrame],
        analysis_params: dict,
        layout: str = "one_per_slide",
        include_title_slide: bool = True,
        include_summary: bool = True,
        method_text: str = "",
    ) -> bytes:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Emu
            from pptx.dml.color import RGBColor
        except ImportError:
            raise ImportError(
                "python-pptx is required. Install with: pip install python-pptx"
            )

        prs = Presentation()
        prs.slide_width = Inches(ReportGenerator.SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(ReportGenerator.SLIDE_HEIGHT_INCHES)

        blank_layout = prs.slide_layouts[6]

        if include_title_slide:
            ReportGenerator._add_title_slide(prs, analysis_params)

        for gene, fig in graphs.items():
            gene_data = processed_data.get(gene)
            ReportGenerator._add_gene_slide(
                prs, blank_layout, gene, fig, gene_data, method_text
            )

        if include_summary:
            ReportGenerator._add_summary_slide(
                prs, blank_layout, processed_data, analysis_params
            )

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()

    @staticmethod
    def _add_title_slide(prs, analysis_params: dict):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Set slide background to white
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White

        # Add logo placeholder in top-left corner
        logo_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(2), Inches(0.8)
        )
        logo_frame = logo_box.text_frame
        logo_frame.word_wrap = True
        logo_para = logo_frame.paragraphs[0]
        logo_para.text = "[Add Logo Here]"
        logo_para.font.size = Pt(12)
        logo_para.font.color.rgb = RGBColor(0xC1, 0xC6, 0xC7)  # Frost Grey

        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.33), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "qPCR Gene Expression Analysis"
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle with efficacy type
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.8), Inches(12.33), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        efficacy = analysis_params.get("Efficacy_Type", "Analysis")
        subtitle_para.text = f"{efficacy} Efficacy Study"
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Info box with date, HK gene, and reference sample
        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5), Inches(12.33), Inches(1)
        )
        info_frame = info_box.text_frame
        info_para = info_frame.paragraphs[0]
        date = analysis_params.get("Date", datetime.now().strftime("%Y-%m-%d"))
        hk = analysis_params.get("Housekeeping_Gene", "N/A")
        ref = analysis_params.get("Reference_Sample", "N/A")
        info_para.text = f"Date: {date}  |  HK Gene: {hk}  |  Reference: {ref}"
        info_para.font.size = Pt(14)
        info_para.font.color.rgb = RGBColor(0xC1, 0xC6, 0xC7)  # Frost Grey
        info_para.alignment = PP_ALIGN.CENTER

        # Add Cosmax Red accent bar at bottom (full width, 0.5" height)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(0xEA, 0x1D, 0x22)  # Cosmax Red
        accent_bar.line.fill.background()  # No border

    @staticmethod
    def _add_gene_slide(
        prs,
        layout,
        gene: str,
        fig: go.Figure,
        gene_data: pd.DataFrame = None,
        method_text: str = "",
    ):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(layout)

        # Set slide background to white
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White

        # Gene title (Black, bold)
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.2), Inches(12.73), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = f"{gene} Expression"
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black
        title_para.alignment = PP_ALIGN.LEFT

        # Method text (if provided) - Frost Grey, small font
        graph_top = Inches(0.9)
        if method_text:
            method_box = slide.shapes.add_textbox(
                Inches(0.3), Inches(0.7), Inches(12.73), Inches(0.4)
            )
            method_frame = method_box.text_frame
            method_frame.word_wrap = True
            method_para = method_frame.paragraphs[0]
            method_para.text = method_text
            method_para.font.size = Pt(11)
            method_para.font.color.rgb = RGBColor(0xC1, 0xC6, 0xC7)  # Frost Grey
            method_para.alignment = PP_ALIGN.LEFT
            # Adjust graph position if method text is present
            graph_top = Inches(1.2)

        fig_copy = go.Figure(fig)
        # Preserve the figure's bottom margin (may be large for angled labels)
        orig_margin = fig.layout.margin
        orig_b = orig_margin.b if orig_margin and orig_margin.b else 140
        fig_copy.update_layout(
            width=1000,
            height=550 + max(0, orig_b - 80),
            margin=dict(l=60, r=60, t=60, b=max(80, orig_b)),
            font=dict(size=14, family=PLOTLY_FONT_FAMILY, color="black"),
        )

        img_bytes = ReportGenerator._fig_to_image(fig_copy, format="png", scale=2)
        img_stream = io.BytesIO(img_bytes)

        left = Inches(0.5)
        top = graph_top
        width = Inches(9.5)
        slide.shapes.add_picture(img_stream, left, top, width=width)

        if gene_data is not None and not gene_data.empty:
            stats_box = slide.shapes.add_textbox(
                Inches(10.2), Inches(1.0), Inches(2.8), Inches(5.5)
            )
            stats_frame = stats_box.text_frame
            stats_frame.word_wrap = True

            header_para = stats_frame.paragraphs[0]
            header_para.text = "Statistics"
            header_para.font.size = Pt(14)
            header_para.font.bold = True

            for _, row in gene_data.iterrows():
                cond = row.get("Condition", "N/A")
                fc = row.get("Fold_Change", row.get("Relative_Expression", 0))
                pval = row.get("p_value", float("nan"))
                sig = row.get("significance", "")

                para = stats_frame.add_paragraph()
                para.text = f"{cond[:12]}"
                para.font.size = Pt(10)
                para.font.bold = True

                para2 = stats_frame.add_paragraph()
                fc_str = f"FC: {fc:.2f}" if pd.notna(fc) else "FC: N/A"
                p_str = f"p={pval:.3f}" if pd.notna(pval) else ""
                para2.text = f"  {fc_str} {sig}"
                para2.font.size = Pt(9)

        # Add Cosmax Red accent bar at bottom (full width, 0.5" height)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(0xEA, 0x1D, 0x22)  # Cosmax Red
        accent_bar.line.fill.background()  # No border

    @staticmethod
    def _add_dual_gene_slide(prs, layout, gene_pairs: list, processed_data: dict):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(layout)

        positions = [
            (Inches(0.3), Inches(0.8), Inches(6.2)),
            (Inches(6.8), Inches(0.8), Inches(6.2)),
        ]

        for idx, (gene, fig) in enumerate(gene_pairs):
            if idx >= 2:
                break

            left, top, width = positions[idx]

            title_box = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.5))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = f"{gene}"
            title_para.font.size = Pt(20)
            title_para.font.bold = True

            fig_copy = go.Figure(fig)
            fig_copy.update_layout(
                width=600,
                height=450,
                margin=dict(l=50, r=30, t=40, b=60),
                font=dict(size=11, family=PLOTLY_FONT_FAMILY, color="black"),
            )

            img_bytes = ReportGenerator._fig_to_image(fig_copy, format="png", scale=2)
            img_stream = io.BytesIO(img_bytes)

            slide.shapes.add_picture(img_stream, left, top, width=width)

    @staticmethod
    def _add_grid_slide(prs, layout, graphs: dict):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.1), Inches(12.73), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Gene Expression Overview"
        title_para.font.size = Pt(24)
        title_para.font.bold = True

        gene_list = list(graphs.items())
        n_genes = len(gene_list)

        if n_genes <= 2:
            cols, rows = 2, 1
        elif n_genes <= 4:
            cols, rows = 2, 2
        elif n_genes <= 6:
            cols, rows = 3, 2
        else:
            cols, rows = 4, 2

        cell_width = 12.5 / cols
        cell_height = 6.5 / rows

        for idx, (gene, fig) in enumerate(gene_list[: cols * rows]):
            col_idx = idx % cols
            row_idx = idx // cols

            left = Inches(0.4 + col_idx * cell_width)
            top = Inches(0.7 + row_idx * cell_height)

            fig_copy = go.Figure(fig)
            fig_copy.update_layout(
                width=350,
                height=280,
                margin=dict(l=40, r=20, t=35, b=40),
                font=dict(size=9, family=PLOTLY_FONT_FAMILY, color="black"),
                title=dict(text=gene, font=dict(size=12, family=PLOTLY_FONT_FAMILY, color="black")),
            )

            img_bytes = ReportGenerator._fig_to_image(fig_copy, format="png", scale=2)
            img_stream = io.BytesIO(img_bytes)

            slide.shapes.add_picture(
                img_stream, left, top, width=Inches(cell_width - 0.2)
            )

    @staticmethod
    def _add_summary_slide(prs, layout, processed_data: dict, analysis_params: dict):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.2), Inches(12.73), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Analysis Summary"
        title_para.font.size = Pt(28)
        title_para.font.bold = True

        all_results = (
            pd.concat(processed_data.values(), ignore_index=True)
            if processed_data
            else pd.DataFrame()
        )

        summary_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.0), Inches(6), Inches(5.5)
        )
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True

        header = summary_frame.paragraphs[0]
        header.text = "Experimental Parameters"
        header.font.size = Pt(16)
        header.font.bold = True

        params_text = [
            f"Efficacy Type: {analysis_params.get('Efficacy_Type', 'N/A')}",
            f"Housekeeping Gene: {analysis_params.get('Housekeeping_Gene', 'N/A')}",
            f"Reference Condition: {analysis_params.get('Reference_Sample', 'N/A')}",
            f"Comparison (*): {analysis_params.get('Compare_To', 'N/A')}",
        ]
        cmp2 = analysis_params.get("Compare_To_2")
        if cmp2:
            params_text.append(f"Comparison (#): {cmp2}")
        cmp3 = analysis_params.get("Compare_To_3")
        if cmp3:
            params_text.append(f"Comparison (\u2020): {cmp3}")
        params_text.extend([
            f"Genes Analyzed: {len(processed_data)}",
            f"Analysis Date: {analysis_params.get('Date', 'N/A')}",
        ])

        for text in params_text:
            para = summary_frame.add_paragraph()
            para.text = text
            para.font.size = Pt(12)

        if not all_results.empty:
            results_box = slide.shapes.add_textbox(
                Inches(7), Inches(1.0), Inches(5.8), Inches(5.5)
            )
            results_frame = results_box.text_frame
            results_frame.word_wrap = True

            header2 = results_frame.paragraphs[0]
            header2.text = "Key Findings"
            header2.font.size = Pt(16)
            header2.font.bold = True

            for gene, gene_df in processed_data.items():
                if gene_df.empty:
                    continue

                para = results_frame.add_paragraph()
                para.text = f"\n{gene}:"
                para.font.size = Pt(12)
                para.font.bold = True

                sig_results = gene_df[
                    gene_df.get("significance", pd.Series([""])).str.len() > 0
                ]
                if not sig_results.empty:
                    for _, row in sig_results.head(3).iterrows():
                        cond = row.get("Condition", "N/A")
                        fc = row.get("Fold_Change", row.get("Relative_Expression", 0))
                        sig = row.get("significance", "")
                        para2 = results_frame.add_paragraph()
                        para2.text = f"  {cond}: {fc:.2f}x {sig}"
                        para2.font.size = Pt(10)
                else:
                    para2 = results_frame.add_paragraph()
                    para2.text = "  No significant changes"
                    para2.font.size = Pt(10)


# ==================== PPT GENERATOR (NEW) ====================
class PPTGenerator:
    NAVY_BLUE = (0, 0, 0)  # Black (was navy #1B365D)
    WHITE = (255, 255, 255)

    @staticmethod
    def _get_color_rgb(rgb_tuple):
        from pptx.dml.color import RGBColor

        return RGBColor(*rgb_tuple)

    @staticmethod
    def create_title_slide(prs, analysis_params):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PPTGenerator._get_color_rgb(PPTGenerator.WHITE)

        # Title "유전자 발현 분석" (Gene Expression Analysis)
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.33), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "유전자 발현 분석 (Gene Expression Analysis)"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.name = "Malgun Gothic"
        p.alignment = PP_ALIGN.CENTER

        # Subtitle (Efficacy Type)
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.8), Inches(11.33), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        efficacy = analysis_params.get("Efficacy_Type", "Analysis")
        p.text = f"{efficacy} Efficacy Test"
        p.font.size = Pt(28)
        p.font.name = "Malgun Gothic"
        p.alignment = PP_ALIGN.CENTER

        # Date and Info
        info_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.0), Inches(11.33), Inches(1)
        )
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        date = analysis_params.get("Date", "")
        p.text = f"Date: {date}"
        p.font.size = Pt(18)
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

        # Logo Placeholder (Bottom Right)
        # X: 11.5, Y: 6.5 approx
        logo_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(6.5), Inches(1.5), Inches(0.5)
        )
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[Logo]"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.RIGHT

    @staticmethod
    def create_gene_slide(prs, gene, fig, gene_data, analysis_params, graph_settings=None, display_name=None):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        display_name = display_name or gene

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PPTGenerator._get_color_rgb(PPTGenerator.WHITE)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{display_name} Expression"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.name = "Malgun Gothic"

        # Navy Blue Line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = PPTGenerator._get_color_rgb(PPTGenerator.NAVY_BLUE)
        line.line.fill.background()

        # Graph — use per-gene dimensions from graph_settings if available
        gs = graph_settings or {}
        fb_w = max(int(gs.get(f"{gene}_figure_width", gs.get("figure_width", 28)) * CM_TO_PX), 800)
        fb_h = max(int(gs.get(f"{gene}_figure_height", gs.get("figure_height", 16)) * CM_TO_PX), 500)
        try:
            orig_m = fig.layout.margin
            extra_b = max(0, (orig_m.b if orig_m and orig_m.b else 0) - 120)
            img_bytes = ReportGenerator._fig_to_image(fig, format="png", scale=2, width=fb_w, height=fb_h + extra_b)
            image_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(
                image_stream,
                Inches(0.5),
                Inches(1.5),
                width=Inches(6.0),
                height=Inches(4.0),
            )
        except Exception as e:
            err_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(6.0), Inches(4.0)
            )
            err_box.text = f"Graph Error: {str(e)}"

        # Data Table (limit to 5 rows to prevent slide overflow — slide height is 7.5")
        if gene_data is not None and not gene_data.empty:
            max_table_rows = 5
            display_data = gene_data.head(max_table_rows) if len(gene_data) > max_table_rows else gene_data
            overflow_note = f" (+{len(gene_data) - max_table_rows} more in Excel)" if len(gene_data) > max_table_rows else ""
            rows = len(display_data) + 1
            cols = 4  # Sample, Condition, Fold Change, P-value
            table_height = Inches(0.3 * rows)
            table_shape = slide.shapes.add_table(
                rows, cols, Inches(0.5), Inches(5.8), Inches(6.0), table_height
            )
            table = table_shape.table

            # Headers
            headers = ["시료명 (Sample)", "Condition", "Fold Change", "P-value"]
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = h
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.bold = True

            # Data
            for i, row in enumerate(display_data.itertuples()):
                r = i + 1
                sample_val = getattr(row, "Original_Sample", getattr(row, "Sample", ""))
                cond_val = getattr(row, "Condition", "")
                fc_val = getattr(
                    row, "Fold_Change", getattr(row, "Relative_Expression", 0)
                )
                pval = getattr(row, "p_value", None)
                sig = getattr(row, "significance", "")

                table.cell(r, 0).text = str(sample_val)[:30]
                table.cell(r, 1).text = str(cond_val)[:30]
                table.cell(r, 2).text = f"{fc_val:.2f}" if pd.notna(fc_val) else "-"
                table.cell(r, 3).text = f"{pval:.4f} {sig}" if pd.notna(pval) else "-"

                for j in range(4):
                    table.cell(r, j).text_frame.paragraphs[0].font.size = Pt(9)

            if overflow_note:
                note_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(5.8 + 0.3 * rows), Inches(6.0), Inches(0.3)
                )
                note_box.text_frame.paragraphs[0].text = f"Showing {max_table_rows}/{len(gene_data)} conditions{overflow_note}"
                note_box.text_frame.paragraphs[0].font.size = Pt(8)
                note_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        # Metadata Box
        meta_box = slide.shapes.add_textbox(
            Inches(7.0), Inches(1.5), Inches(2.5), Inches(5.0)
        )
        tf = meta_box.text_frame
        tf.word_wrap = True

        def add_meta_line(text, bold=False):
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(11)
            p.font.bold = bold

        add_meta_line("Analysis Parameters", bold=True)
        add_meta_line(f"HK Gene: {analysis_params.get('Housekeeping_Gene', '-')}")
        add_meta_line(f"Ref Sample: {analysis_params.get('Reference_Sample', '-')}")
        add_meta_line(f"Compare (*): {analysis_params.get('Compare_To', '-')}")
        compare_2 = analysis_params.get("Compare_To_2")
        if compare_2:
            add_meta_line(f"Compare (#): {compare_2}")
        compare_3 = analysis_params.get("Compare_To_3")
        if compare_3:
            add_meta_line(f"Compare (\u2020): {compare_3}")
        add_meta_line("")
        add_meta_line("통계적 유의성 (Significance)", bold=True)
        add_meta_line("* p < 0.05   ** p < 0.01   *** p < 0.001")
        if compare_2:
            add_meta_line("# p < 0.05   ## p < 0.01   ### p < 0.001")
        if compare_3:
            add_meta_line("\u2020 p < 0.05   \u2020\u2020 p < 0.01   \u2020\u2020\u2020 p < 0.001")

        # Logo Placeholder
        logo_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(6.5), Inches(1.5), Inches(0.5)
        )
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[Logo]"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.RIGHT

    @staticmethod
    def _delete_slide(prs, index):
        slides_list = prs.slides._sldIdLst
        rId = slides_list[index].attrib[
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        ]
        prs.part.drop_rel(rId)
        slides_list.remove(slides_list[index])

    @staticmethod
    def _safe_copy_slide(prs, source_index):
        """Copy a slide by duplicating shapes via the spTree proxy.

        Uses python-pptx's spTree property to ensure shapes are properly
        recognized after copy. Copies relationships for images/media,
        then deep-copies each shape element from the source spTree.
        """
        import copy

        source = prs.slides[source_index]
        new_slide = prs.slides.add_slide(source.slide_layout)

        # Map source relationship IDs to new ones
        rid_map = {}
        for rid, rel in source.part.rels.items():
            rel_type = str(rel.reltype)
            if 'slideLayout' in rel_type:
                for new_rid, new_rel in new_slide.part.rels.items():
                    if 'slideLayout' in str(new_rel.reltype):
                        rid_map[rid] = new_rid
                        break
                continue
            if 'notesSlide' in rel_type:
                continue
            try:
                if rel.is_external:
                    new_rid = new_slide.part.rels.get_or_add_ext_rel(
                        rel.reltype, rel.target_ref
                    )
                else:
                    new_rid = new_slide.part.relate_to(
                        rel.target_part, rel.reltype
                    )
                rid_map[rid] = new_rid
            except Exception as e:
                import logging
                logging.warning(f"PPT slide copy: skipped relationship {rid}: {e}")
                # Surface to the user so a broken slide isn't silently shipped.
                try:
                    st.warning(
                        f"⚠️ PPT slide copy skipped a relationship "
                        f"({rel.reltype}). The exported slide may be missing an "
                        f"image or link. Details: {e}"
                    )
                except Exception:
                    pass

        source_spTree = source._element.spTree
        target_spTree = new_slide._element.spTree

        # Remove layout placeholder shapes from target
        shape_tags = ['}sp', '}pic', '}cxnSp', '}grpSp']
        for child in list(target_spTree):
            if any(child.tag.endswith(t) for t in shape_tags):
                target_spTree.remove(child)

        # Deep-copy shapes from source, remapping relationship IDs
        for child in source_spTree:
            if any(child.tag.endswith(t) for t in shape_tags):
                new_child = copy.deepcopy(child)
                for elem in new_child.iter():
                    for attr_key in list(elem.attrib.keys()):
                        val = elem.attrib[attr_key]
                        if val in rid_map:
                            elem.attrib[attr_key] = rid_map[val]
                target_spTree.append(new_child)

        return new_slide

    @staticmethod
    def _move_slide_to_end(prs, slide_index):
        slides_list = prs.slides._sldIdLst
        el = slides_list[slide_index]
        slides_list.remove(el)
        slides_list.append(el)

    @staticmethod
    def _populate_gene_slide(prs, slide, gene, fig, gene_data, analysis_params, graph_settings, display_name=None):
        """Populate a template-duplicated gene slide with gene name and graph image.

        `display_name` overrides `gene` for ALL user-visible text on the slide
        (title, "효능 평가" prefix). When the user renamed the gene in the Graphs
        tab, the rename is honored here so the slide matches the on-screen graph.
        """
        from pptx.util import Inches, Pt, Emu

        display_name = display_name or gene

        # Update text boxes on the template slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text

            # TextBox 17: "효능 평가" → "{display_name} 효능 평가"
            if "효능 평가" in full_text and "Results" not in full_text and "有" not in full_text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "효능" in run.text:
                            run.text = run.text.replace("효능 평가", f"{display_name} 효능 평가")
                            break
                    break

            # TextBox 14: "Results: 효능 有/無" → update with gene result
            elif "Results" in full_text or "효능 有" in full_text:
                has_efficacy = False
                if gene_data is not None and not gene_data.empty:
                    sig_vals = gene_data.get("significance", pd.Series())
                    has_efficacy = sig_vals.str.len().gt(0).any() if not sig_vals.empty else False
                result_text = f"Results: 효능 {'有' if has_efficacy else '無'}"
                for para in shape.text_frame.paragraphs:
                    runs = list(para.runs)
                    for i, run in enumerate(runs):
                        if i == 0:
                            run.text = result_text
                        else:
                            run.text = ""
                    break

            # TextBox 2: experiment details (top-right area) — fill template fields
            elif (shape.left is not None and shape.top is not None
                  and shape.left > 3500000 and shape.top < 500000):
                efficacy_t = analysis_params.get("Efficacy_Type", "")
                eff_cfg = EFFICACY_CONFIG.get(efficacy_t, {})
                field_values = {
                    "Date: ": analysis_params.get("Date", "")[:10],
                    "Cell line: ": eff_cfg.get("cell", ""),
                    "Sample concentration: ": analysis_params.get("concentration", "1 ppm"),
                    "Positive control: ": eff_cfg.get("controls", {}).get("positive", ""),
                    "Inducer: ": eff_cfg.get("controls", {}).get("negative", ""),
                    "Treatment time: ": analysis_params.get("treatment_time", "24 h"),
                    "Test method:": " qPCR (DDCt)",
                }
                for para in shape.text_frame.paragraphs:
                    if not para.runs:
                        continue
                    label = para.runs[0].text
                    for field_label, value in field_values.items():
                        if label.strip().startswith(field_label.strip()):
                            para.runs[0].text = f"{field_label}{value}"
                            for extra_run in para.runs[1:]:
                                extra_run.text = ""
                            break

        # Add graph image
        gs = graph_settings or {}
        w_cm = gs.get(f"{gene}_figure_width", gs.get("figure_width", 28))
        h_cm = gs.get(f"{gene}_figure_height", gs.get("figure_height", 16))

        fig_copy = go.Figure(fig)
        fig_copy.update_layout(
            width=max(int(w_cm * CM_TO_PX), 800),
            height=max(int(h_cm * CM_TO_PX), 500),
        )

        try:
            img_bytes = ReportGenerator._fig_to_image(fig_copy, format="png", scale=2)
            img_stream = io.BytesIO(img_bytes)

            w_emu = int(w_cm * CM_TO_EMU)
            h_emu = int(h_cm * CM_TO_EMU)
            # Available area: ~1.0" to ~6.6" vertically, ~12" wide
            max_w_emu = int(11.5 * 914400)
            max_h_emu = int(5.2 * 914400)

            if w_emu > max_w_emu:
                scale_f = max_w_emu / w_emu
                w_emu = max_w_emu
                h_emu = int(h_emu * scale_f)
            if h_emu > max_h_emu:
                scale_f = max_h_emu / h_emu
                h_emu = max_h_emu
                w_emu = int(w_emu * scale_f)

            slide_w = int(prs.slide_width)
            left = max(0, int((slide_w - w_emu) / 2))
            top = int(1.0 * 914400)  # 1.0" below header

            slide.shapes.add_picture(
                img_stream, left, top, width=Emu(w_emu), height=Emu(h_emu)
            )
        except Exception as e:
            err_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(11.0), Inches(4.0)
            )
            err_box.text_frame.paragraphs[0].text = f"Graph Error: {str(e)}"
            err_box.text_frame.paragraphs[0].font.size = Pt(14)

    @staticmethod
    def generate_presentation(graphs, processed_data, analysis_params, graph_settings=None, gene_display_names=None):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Emu
        except ImportError:
            st.error("python-pptx not installed")
            return None

        import os

        script_dir = os.path.dirname(os.path.abspath(__file__))
        template_path = os.path.join(
            script_dir, "251215 효능평가 결과 TEMPLATE.pptx"
        )
        use_template = os.path.isfile(template_path)

        gs = graph_settings or {}
        gene_display_names = gene_display_names or {}
        graphs = graphs or {}

        # Source of truth is processed_data (the analysis output). Stale entries
        # in `graphs` from earlier runs are ignored; missing figures are
        # rendered lazily below so PPT works even if the user never opened
        # the Graphs tab.
        gene_list = [
            g for g in (processed_data.keys() if processed_data else [])
            if processed_data.get(g) is not None and not processed_data[g].empty
        ]
        n_genes = len(gene_list)

        def _get_fig(gene):
            """Return a Plotly figure for `gene`, rendering on demand if missing."""
            fig = graphs.get(gene)
            if fig is not None:
                return fig
            try:
                gene_data_local = processed_data.get(gene)
                if gene_data_local is None or gene_data_local.empty:
                    return None
                disp = gene_display_names.get(gene, gene)
                return GraphGenerator.create_gene_graph(
                    gene_data_local,
                    gene,
                    gs or {},
                    EFFICACY_CONFIG.get(analysis_params.get("Efficacy_Type", ""), {}),
                    sample_order=st.session_state.get("sample_order"),
                    display_gene_name=disp,
                    ref_condition=st.session_state.get("analysis_ref_condition"),
                )
            except Exception:
                return None

        if use_template and n_genes > 0:
            prs = Presentation(template_path)

            # Template slides: [0]=Title, [1]=Description, [2]=Gene template, [3]=Closing
            # Step 1: Update title slide — replace "소재" with efficacy type
            efficacy = analysis_params.get("Efficacy_Type", "소재")
            title_slide = prs.slides[0]
            for shape in title_slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "소재" in run.text:
                                run.text = run.text.replace("소재", efficacy)

            # Step 2: Copy gene template (index 2) BEFORE deleting to avoid
            # ZIP part-name conflicts (delete frees a name that add_slide reuses)
            for _ in range(n_genes - 1):
                PPTGenerator._safe_copy_slide(prs, 2)
            # Now: [0]=Title, [1]=Desc, [2]=Gene, [3]=Closing, [4..N+2]=GeneCopies

            # Step 3: Delete slide 1 (description overview)
            PPTGenerator._delete_slide(prs, 1)
            # Now: [0]=Title, [1]=Gene, [2]=Closing, [3..N+1]=GeneCopies

            # Step 4: Move closing slide (index 2) to end
            PPTGenerator._move_slide_to_end(prs, 2)
            # Now: [0]=Title, [1]=Gene, [2..N]=GeneCopies, [N+1]=Closing

            # Step 5: Populate each gene slide with gene name and graph
            for i, gene in enumerate(gene_list):
                slide = prs.slides[1 + i]
                fig = _get_fig(gene)
                gene_data = processed_data.get(gene)
                PPTGenerator._populate_gene_slide(
                    prs, slide, gene, fig, gene_data, analysis_params, gs,
                    display_name=gene_display_names.get(gene, gene),
                )
        else:
            # Fallback: no template available
            prs = Presentation()
            prs.slide_width = Emu(12192000)
            prs.slide_height = Emu(6858000)
            PPTGenerator.create_title_slide(prs, analysis_params)

            for gene in gene_list:
                fig = _get_fig(gene)
                gene_data = processed_data.get(gene)
                PPTGenerator.create_gene_slide(
                    prs, gene, fig, gene_data, analysis_params, graph_settings=gs,
                    display_name=gene_display_names.get(gene, gene),
                )

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()


# ==================== EXPORT FUNCTIONS ====================
def _sanitize_sheet_name(name: str, used_names: set) -> str:
    """Sanitize Excel sheet name: remove invalid chars, truncate, deduplicate."""
    import re
    safe = re.sub(r'[\\/*\[\]:?]', '_', name)[:31]
    base = safe
    counter = 1
    while safe in used_names:
        suffix = f"_{counter}"
        safe = base[: 31 - len(suffix)] + suffix
        counter += 1
    used_names.add(safe)
    return safe


def export_to_excel(
    raw_data: pd.DataFrame,
    processed_data: Dict[str, pd.DataFrame],
    params: dict,
    mapping: dict,
    qc_stats: dict = None,
    replicate_stats: pd.DataFrame = None,
    excluded_wells=None,
    gene_display_names: dict = None,
) -> bytes:
    """Export comprehensive Excel with gene-by-gene sheets, QC report, and FC matrix.

    `gene_display_names` maps raw gene name -> user-edited display name (from Graphs
    tab). When provided, per-gene sheet names, the FC matrix index, and the chart
    Y-axis title all use the display name so Excel matches what the user sees on
    screen. Raw gene name is preserved as a `Target_Raw` column for traceability.
    """
    gene_display_names = gene_display_names or {}
    _disp = lambda g: str(gene_display_names.get(g, g))
    output = io.BytesIO()

    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        # Parameters sheet
        pd.DataFrame([params]).to_excel(
            writer, sheet_name="Analysis_Parameters", index=False
        )

        # Sample mapping sheet
        pd.DataFrame([{"Original": k, **v} for k, v in mapping.items()]).to_excel(
            writer, sheet_name="Sample_Mapping", index=False
        )

        # Raw data (include mapped Condition column reflecting sample mapping)
        raw_export = raw_data.copy()
        if mapping:
            raw_export["Condition"] = raw_export["Sample"].map(
                lambda x: mapping.get(x, {}).get("condition", x)
            )
        else:
            raw_export["Condition"] = raw_export["Sample"]
        raw_export = (
            raw_export[["Well", "Sample", "Condition", "Target", "CT", "Source_File"]]
            if "Source_File" in raw_export.columns
            else raw_export
        )
        raw_export.to_excel(writer, sheet_name="Raw_Data", index=False)

        # Gene-by-gene calculations with statistical test method column
        _used_sheet_names = {"Analysis_Parameters", "Raw_Data"}
        for gene, gene_data in processed_data.items():
            display = _disp(gene)
            sheet_name = _sanitize_sheet_name(f"{display}_Analysis", _used_sheet_names)
            gene_export = gene_data.copy()
            # Preserve raw gene name so an exported sheet can be cross-referenced
            # to raw_data["Target"] even when the user has renamed the gene.
            if "Target" in gene_export.columns and display != gene:
                gene_export.insert(
                    list(gene_export.columns).index("Target") + 1,
                    "Target_Raw",
                    gene,
                )
                gene_export["Target"] = display
            if "p_value" in gene_export.columns:
                ttest_type = params.get("ttest_type", "welch")
                gene_export["Stat_Test"] = ""
                for idx, row in gene_export.iterrows():
                    n_rep = row.get("n_replicates", 0)
                    if pd.notna(row.get("p_value")) and row.get("p_value") is not np.nan:
                        if n_rep >= 2:
                            gene_export.loc[idx, "Stat_Test"] = (
                                f"{'Welch' if ttest_type == 'welch' else 'Student'} t-test (n={n_rep})"
                            )
                        elif n_rep == 1:
                            gene_export.loc[idx, "Stat_Test"] = f"One-sample t-test (n={n_rep})"
                        else:
                            gene_export.loc[idx, "Stat_Test"] = "N/A"
            gene_export.to_excel(writer, sheet_name=sheet_name, index=False)

        # Summary sheet
        if processed_data:
            non_empty = []
            for g, df in processed_data.items():
                if df.empty:
                    continue
                d = df.copy()
                if "Target" in d.columns:
                    d["Target"] = _disp(g)
                non_empty.append(d)
            if non_empty:
                all_data = pd.concat(non_empty, ignore_index=True)
                agg_dict = {"Relative_Expression": ["mean", "std", "count"]}
                if "p_value" in all_data.columns:
                    agg_dict["p_value"] = "min"
                group_cols = ["Target"]
                if "Group" in all_data.columns:
                    group_cols.append("Group")
                summary = (
                    all_data.groupby(group_cols)
                    .agg(agg_dict)
                    .round(4)
                )
                summary.to_excel(writer, sheet_name="Summary")

        # Fold Change Matrix (pivot table)
        if processed_data:
            non_empty = []
            for g, df in processed_data.items():
                if df.empty:
                    continue
                d = df.copy()
                if "Target" in d.columns:
                    d["Target"] = _disp(g)
                non_empty.append(d)
            if non_empty:
                all_data = pd.concat(non_empty, ignore_index=True)
                if "Fold_Change" in all_data.columns and "Condition" in all_data.columns:
                    fc_matrix = all_data.pivot_table(
                        values="Fold_Change",
                        index="Target",
                        columns="Condition",
                        aggfunc="first",
                    )
                    fc_matrix = fc_matrix.round(4)
                    fc_matrix.to_excel(writer, sheet_name="FC_Matrix")

        # Replicate-level fold changes
        if raw_data is not None and excluded_wells is not None:
            hk_gene = params.get("Housekeeping_Gene")
            ref_sample = params.get("Reference_Sample")
            if hk_gene and ref_sample:
                try:
                    replicate_fc = AnalysisEngine.compute_replicate_fold_changes(
                        raw_data=raw_data,
                        hk_gene=hk_gene,
                        ref_sample=ref_sample,
                        sample_mapping=mapping,
                        excluded_wells=excluded_wells,
                    )
                    if not replicate_fc.empty:
                        replicate_fc.to_excel(writer, sheet_name="Replicate_FC", index=False)
                except Exception as e:
                    import logging
                    logging.warning(f"Replicate_FC sheet skipped: {e}")

        # QC Report sheet
        _write_qc_report_sheet(writer, qc_stats, replicate_stats)

    # Post-process: add gene chart sheets with openpyxl (supports rich text axis titles)
    output = _add_gene_chart_sheets(
        output, processed_data, params, gene_display_names=gene_display_names
    )

    return output.getvalue()


def _write_qc_report_sheet(writer, qc_stats=None, replicate_stats=None):
    """Write QC Report sheet with summary stats and replicate-level CV/outlier info."""
    rows = []
    if qc_stats:
        rows.append({"Metric": "Total Wells", "Value": qc_stats.get("total_wells", "")})
        rows.append({"Metric": "Excluded Wells", "Value": qc_stats.get("excluded_wells", "")})
        rows.append({"Metric": "Active Wells", "Value": qc_stats.get("active_wells", "")})
        rows.append({"Metric": "CT Mean", "Value": qc_stats.get("ct_mean", "")})
        rows.append({"Metric": "CT SD", "Value": qc_stats.get("ct_std", "")})
        rows.append({"Metric": "CT Range", "Value": f"{qc_stats.get('ct_min', '')}-{qc_stats.get('ct_max', '')}"})
        rows.append({"Metric": "High CT Wells (>35)", "Value": qc_stats.get("high_ct_count", "")})
        rows.append({"Metric": "Low CT Wells (<10)", "Value": qc_stats.get("low_ct_count", "")})
        rows.append({"Metric": "Total Triplicates", "Value": qc_stats.get("total_triplicates", "")})
        rows.append({"Metric": "Healthy Triplicates", "Value": qc_stats.get("healthy_triplicates", "")})
        rows.append({"Metric": "Warning Triplicates", "Value": qc_stats.get("warning_triplicates", "")})
        rows.append({"Metric": "Error Triplicates", "Value": qc_stats.get("error_triplicates", "")})
        rows.append({"Metric": "Avg CV%", "Value": qc_stats.get("avg_cv_pct", "")})
        rows.append({"Metric": "Max CV%", "Value": qc_stats.get("max_cv_pct", "")})
        rows.append({"Metric": "Health Score (%)", "Value": qc_stats.get("health_score", "")})
    if rows:
        pd.DataFrame(rows).to_excel(writer, sheet_name="QC_Report", index=False, startrow=0)
    if replicate_stats is not None and not replicate_stats.empty:
        start_row = len(rows) + 3 if rows else 0
        replicate_stats.to_excel(writer, sheet_name="QC_Report", index=False, startrow=start_row)


def _add_gene_chart_sheets(output_buf, processed_data, params, gene_display_names=None):
    """Post-process Excel bytes to add per-gene chart sheets using openpyxl.

    Two-phase approach:
    Phase 1: openpyxl API for chart creation (series, errBars, axis properties)
    Phase 2: zip-level XML post-processing for elements openpyxl can't set
             (axis label fonts, rich text Y-axis title, hidden gridlines)

    `gene_display_names` maps raw gene -> display name. The sheet name and the
    Y-axis title both use the display name so they always agree.
    """
    from openpyxl import load_workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.chart.data_source import NumRef, NumDataSource
    from openpyxl.chart.error_bar import ErrorBars
    from openpyxl.chart.shapes import GraphicalProperties
    from openpyxl.utils import get_column_letter
    from lxml import etree
    import zipfile

    if not processed_data:
        return output_buf

    gene_display_names = gene_display_names or {}
    _disp = lambda g: str(gene_display_names.get(g, g))

    output_buf.seek(0)
    wb = load_workbook(output_buf)
    hk_gene = params.get("Housekeeping_Gene", params.get("reference_gene", "\u03b2-Actin"))

    _used_sheet_names = set(wb.sheetnames)
    # Track display name per chart file so phase 2 can render the right
    # Y-axis title even when sheet names were truncated/deduplicated.
    display_for_chart_order = []

    for gene, gene_data in processed_data.items():
        if gene_data is None or gene_data.empty:
            continue

        display = _disp(gene)
        display_for_chart_order.append(display)
        sheet_name = _sanitize_sheet_name(f"{display}_Chart", _used_sheet_names)
        ws = wb.create_sheet(sheet_name)

        # Build data columns
        conditions = gene_data["Condition"].tolist() if "Condition" in gene_data.columns else gene_data.get("Group", pd.Series()).tolist()
        fold_changes = gene_data.get("Fold_Change", gene_data.get("Relative_Expression", pd.Series())).tolist()
        sems = gene_data.get("SEM", pd.Series([0] * len(gene_data))).tolist()
        p_values = gene_data.get("p_value", pd.Series()).tolist()
        significances = gene_data.get("significance", pd.Series([""]*len(gene_data))).tolist()
        has_p2 = "p_value_2" in gene_data.columns
        p_values_2 = gene_data.get("p_value_2", pd.Series()).tolist() if has_p2 else []
        sig_2 = gene_data.get("significance_2", pd.Series()).tolist() if has_p2 else []
        has_p3 = "p_value_3" in gene_data.columns
        p_values_3 = gene_data.get("p_value_3", pd.Series()).tolist() if has_p3 else []
        sig_3 = gene_data.get("significance_3", pd.Series()).tolist() if has_p3 else []

        # Fold-change-domain asymmetric error bounds (Livak). Match the on-screen
        # Plotly graph instead of plotting Ct-domain SEM symmetrically on the
        # linear fold-change axis. Fall back to symmetric SEM when unavailable.
        if "FC_Error_Upper" in gene_data.columns and "FC_Error_Lower" in gene_data.columns:
            fc_err_upper = gene_data["FC_Error_Upper"].fillna(0).tolist()
            fc_err_lower = gene_data["FC_Error_Lower"].fillna(0).tolist()
        else:
            fc_err_upper = [s if pd.notna(s) else 0 for s in sems]
            fc_err_lower = list(fc_err_upper)

        n_rows = len(conditions)
        if n_rows == 0:
            continue

        hdr_row, data_start = 5, 6
        headers = ["Condition", "Fold_Change", "SEM", "p_value", "significance"]
        if has_p2:
            headers.extend(["p_value_2", "significance_2"])
        if has_p3:
            headers.extend(["p_value_3", "significance_3"])
        for ci, h in enumerate(headers):
            ws.cell(row=hdr_row, column=3 + ci, value=h)
        # Asymmetric fold-change error columns, appended after the visible data.
        err_plus_col = 3 + len(headers)
        err_minus_col = err_plus_col + 1
        ws.cell(row=hdr_row, column=err_plus_col, value="FC_Err_Upper")
        ws.cell(row=hdr_row, column=err_minus_col, value="FC_Err_Lower")
        max_cond_len = max((len(str(c)) for c in conditions), default=10)
        ws.column_dimensions["C"].width = max(18, min(max_cond_len + 2, 45))
        for col_letter in ("D", "E", "F", "G"):
            ws.column_dimensions[col_letter].width = 14

        for i in range(n_rows):
            r = data_start + i
            ws.cell(row=r, column=3, value=conditions[i])
            fc = fold_changes[i]
            ws.cell(row=r, column=4, value=fc if pd.notna(fc) else 0)
            sem = sems[i]
            ws.cell(row=r, column=5, value=sem if pd.notna(sem) else 0)
            pv = p_values[i] if i < len(p_values) else None
            ws.cell(row=r, column=6, value=pv if pd.notna(pv) else None)
            sig = significances[i] if i < len(significances) else ""
            ws.cell(row=r, column=7, value=sig if pd.notna(sig) else "")
            if has_p2:
                pv2 = p_values_2[i] if i < len(p_values_2) else None
                ws.cell(row=r, column=8, value=pv2 if pd.notna(pv2) else None)
                s2 = sig_2[i] if i < len(sig_2) else ""
                ws.cell(row=r, column=9, value=s2 if pd.notna(s2) else "")
            p3_col_offset = 10 if has_p2 else 8
            if has_p3:
                pv3 = p_values_3[i] if i < len(p_values_3) else None
                ws.cell(row=r, column=p3_col_offset, value=pv3 if pd.notna(pv3) else None)
                s3_val = sig_3[i] if i < len(sig_3) else ""
                ws.cell(row=r, column=p3_col_offset + 1, value=s3_val if pd.notna(s3_val) else "")
            eu = fc_err_upper[i] if i < len(fc_err_upper) else 0
            el = fc_err_lower[i] if i < len(fc_err_lower) else 0
            ws.cell(row=r, column=err_plus_col, value=float(eu) if pd.notna(eu) else 0)
            ws.cell(row=r, column=err_minus_col, value=float(el) if pd.notna(el) else 0)

        last_data_row = data_start + n_rows - 1

        # ---- Phase 1: Create chart via openpyxl API ----
        chart = BarChart()
        chart.type = "col"
        chart.grouping = "clustered"
        chart.varyColors = False
        chart.legend = None
        chart.title = None
        chart.width = max(15, n_rows * 1.2)
        chart.height = 7.5
        chart.gapWidth = 219
        chart.overlap = -27

        cats = Reference(ws, min_col=3, min_row=data_start, max_row=last_data_row)
        vals = Reference(ws, min_col=4, min_row=data_start, max_row=last_data_row)
        chart.add_data(vals, titles_from_data=False)
        chart.set_categories(cats)

        series = chart.series[0]

        # Bar style: white fill (bg1), black outline (tx1) 0.5pt
        bar_sp_xml = (
            '<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:solidFill><a:schemeClr val="bg1"/></a:solidFill>'
            '<a:ln w="6350"><a:solidFill><a:schemeClr val="tx1"/></a:solidFill>'
            '<a:prstDash val="solid"/></a:ln></c:spPr>'
        )
        series.graphicalProperties = GraphicalProperties.from_tree(etree.fromstring(bar_sp_xml))

        # Error bars: custom asymmetric fold-change bounds (upper/lower), gray line
        plus_col_letter = get_column_letter(err_plus_col)
        minus_col_letter = get_column_letter(err_minus_col)
        plus_formula = f"'{sheet_name}'!${plus_col_letter}${data_start}:${plus_col_letter}${last_data_row}"
        minus_formula = f"'{sheet_name}'!${minus_col_letter}${data_start}:${minus_col_letter}${last_data_row}"
        eb_sp_xml = (
            '<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:noFill/>'
            '<a:ln w="9525" cap="flat" cmpd="sng" algn="ctr">'
            '<a:solidFill><a:schemeClr val="tx1">'
            '<a:lumMod val="65000"/><a:lumOff val="35000"/>'
            '</a:schemeClr></a:solidFill>'
            '<a:prstDash val="solid"/><a:round/></a:ln></c:spPr>'
        )
        eb_sp = GraphicalProperties.from_tree(etree.fromstring(eb_sp_xml))
        series.errBars = ErrorBars(
            errValType="cust", errBarType="both", noEndCap=False,
            minus=NumDataSource(numRef=NumRef(f=minus_formula)),
            plus=NumDataSource(numRef=NumRef(f=plus_formula)),
            spPr=eb_sp,
        )

        # X-axis: black line 0.5pt, no tick marks
        chart.x_axis.delete = False
        chart.x_axis.majorTickMark = "none"
        chart.x_axis.minorTickMark = "none"
        xax_sp_xml = (
            '<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:noFill/>'
            '<a:ln w="6350" cap="flat" cmpd="sng" algn="ctr">'
            '<a:solidFill><a:schemeClr val="tx1"/></a:solidFill>'
            '<a:prstDash val="solid"/><a:round/></a:ln></c:spPr>'
        )
        chart.x_axis.graphicalProperties = GraphicalProperties.from_tree(etree.fromstring(xax_sp_xml))

        # Y-axis: black line 0.5pt, tick out, min=0, no gridlines
        chart.y_axis.delete = False
        chart.y_axis.majorTickMark = "out"
        chart.y_axis.minorTickMark = "none"
        chart.y_axis.scaling.min = 0
        chart.y_axis.majorGridlines = None
        yax_sp_xml = (
            '<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:noFill/>'
            '<a:ln w="6350"><a:solidFill><a:schemeClr val="tx1"/></a:solidFill>'
            '<a:prstDash val="solid"/></a:ln></c:spPr>'
        )
        chart.y_axis.graphicalProperties = GraphicalProperties.from_tree(etree.fromstring(yax_sp_xml))

        ws.add_chart(chart, "I3")

    # ---- Phase 1 complete: save workbook to bytes ----
    phase1_buf = io.BytesIO()
    wb.save(phase1_buf)

    # ---- Phase 2: post-process chart XMLs in the XLSX zip ----
    C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    nsmap = {"c": C, "a": A}

    phase1_buf.seek(0)
    zf_in = zipfile.ZipFile(phase1_buf, "r")

    # Identify which chart files are gene charts (the ones we just added).
    # Sort NUMERICALLY (chart1, chart2, ..., chart10, chart11) — a plain sorted()
    # is lexicographic (chart1, chart10, chart11, ..., chart2) which mispairs
    # genes to charts as soon as there are 10+ genes, stamping the wrong gene
    # name onto a chart's Y-axis title.
    def _chart_file_num(name):
        m = re.search(r"chart(\d+)\.xml$", name)
        return int(m.group(1)) if m else 0

    all_chart_files = sorted(
        (n for n in zf_in.namelist()
         if n.startswith("xl/charts/chart") and n.endswith(".xml")),
        key=_chart_file_num,
    )
    # Gene chart files are the last N chart files (one per gene with data), in
    # creation order. Use the display-name order captured during phase 1 so the
    # Y-axis title matches the sheet name even when gene names were renamed.
    n_gene_charts = len(display_for_chart_order)
    gene_chart_files = all_chart_files[-n_gene_charts:] if n_gene_charts > 0 else []
    gene_for_chart = dict(zip(gene_chart_files, display_for_chart_order))

    result = io.BytesIO()
    zf_out = zipfile.ZipFile(result, "w", compression=zipfile.ZIP_DEFLATED)

    for item in zf_in.infolist():
        data = zf_in.read(item.filename)

        if item.filename in gene_for_chart:
            gene = gene_for_chart[item.filename]
            root = etree.fromstring(data)

            cat_ax = root.find(".//c:catAx", nsmap)
            val_ax = root.find(".//c:valAx", nsmap)

            # X-axis: add txPr for 9pt gray Arial labels
            if cat_ax is not None:
                cat_ax.append(_build_axis_txpr(C, A))

            if val_ax is not None:
                # Hidden gridlines (noFill)
                mgrid = etree.SubElement(val_ax, f"{{{C}}}majorGridlines")
                mg_sp = etree.SubElement(mgrid, f"{{{C}}}spPr")
                mg_ln = etree.SubElement(mg_sp, f"{{{A}}}ln")
                etree.SubElement(mg_ln, f"{{{A}}}noFill")

                # Y-axis tick label font
                val_ax.append(_build_axis_txpr(C, A))

                # Rich text Y-axis title
                _build_yaxis_title(val_ax, gene, hk_gene, C, A)

            # Chart area: white fill, light gray border
            chart_space = root
            cs_sp = etree.SubElement(chart_space, f"{{{C}}}spPr")
            sf_cs = etree.SubElement(cs_sp, f"{{{A}}}solidFill")
            etree.SubElement(sf_cs, f"{{{A}}}schemeClr").set("val", "bg1")
            ln_cs = etree.SubElement(cs_sp, f"{{{A}}}ln")
            ln_cs.set("w", "9525")
            sf_ln = etree.SubElement(ln_cs, f"{{{A}}}solidFill")
            sc_ln = etree.SubElement(sf_ln, f"{{{A}}}schemeClr")
            sc_ln.set("val", "tx1")
            etree.SubElement(sc_ln, f"{{{A}}}lumMod").set("val", "15000")
            etree.SubElement(sc_ln, f"{{{A}}}lumOff").set("val", "85000")

            data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

        zf_out.writestr(item, data)

    zf_in.close()
    zf_out.close()
    result.seek(0)
    return result


def _build_axis_txpr(C, A):
    """Build txPr element for axis tick labels: 9pt gray Arial."""
    from lxml import etree

    txpr = etree.Element(f"{{{C}}}txPr")
    bp = etree.SubElement(txpr, f"{{{A}}}bodyPr")
    bp.set("rot", "-60000000")
    bp.set("vert", "horz")
    bp.set("wrap", "square")
    etree.SubElement(txpr, f"{{{A}}}lstStyle")
    p = etree.SubElement(txpr, f"{{{A}}}p")
    pPr = etree.SubElement(p, f"{{{A}}}pPr")
    dRPr = etree.SubElement(pPr, f"{{{A}}}defRPr")
    dRPr.set("sz", "900")
    dRPr.set("b", "0")
    sf = etree.SubElement(dRPr, f"{{{A}}}solidFill")
    sc = etree.SubElement(sf, f"{{{A}}}schemeClr")
    sc.set("val", "tx1")
    etree.SubElement(sc, f"{{{A}}}lumMod").set("val", "65000")
    etree.SubElement(sc, f"{{{A}}}lumOff").set("val", "35000")
    lat = etree.SubElement(dRPr, f"{{{A}}}latin")
    lat.set("typeface", "Arial")
    cs = etree.SubElement(dRPr, f"{{{A}}}cs")
    cs.set("typeface", "Arial")
    return txpr


def _build_yaxis_title(val_ax, gene, hk_gene, C, A):
    """Build rich text Y-axis title: 'Relative mRNA expression level of' + gene(RED)/HK(black)."""
    from lxml import etree

    title_el = etree.SubElement(val_ax, f"{{{C}}}title")
    tx = etree.SubElement(title_el, f"{{{C}}}tx")
    rich = etree.SubElement(tx, f"{{{C}}}rich")
    body = etree.SubElement(rich, f"{{{A}}}bodyPr")
    body.set("rot", "-5400000")
    body.set("vert", "horz")
    body.set("wrap", "square")
    body.set("anchor", "ctr")
    body.set("anchorCtr", "1")
    etree.SubElement(rich, f"{{{A}}}lstStyle")

    # Paragraph 1: "Relative mRNA expression level of"
    p1 = etree.SubElement(rich, f"{{{A}}}p")
    p1Pr = etree.SubElement(p1, f"{{{A}}}pPr")
    etree.SubElement(p1Pr, f"{{{A}}}defRPr")
    r1 = etree.SubElement(p1, f"{{{A}}}r")
    rPr1 = etree.SubElement(r1, f"{{{A}}}rPr")
    rPr1.set("lang", "en-US")
    rPr1.set("sz", "800")
    rPr1.set("b", "1")
    sf1 = etree.SubElement(rPr1, f"{{{A}}}solidFill")
    sc1 = etree.SubElement(sf1, f"{{{A}}}sysClr")
    sc1.set("val", "windowText")
    sc1.set("lastClr", "000000")
    lat1 = etree.SubElement(rPr1, f"{{{A}}}latin")
    lat1.set("typeface", "Arial")
    cs1 = etree.SubElement(rPr1, f"{{{A}}}cs")
    cs1.set("typeface", "Arial")
    t1 = etree.SubElement(r1, f"{{{A}}}t")
    t1.text = "Relative mRNA expression level of"

    # Paragraph 2: gene name (RED) + /HK_gene (black)
    p2 = etree.SubElement(rich, f"{{{A}}}p")
    p2Pr = etree.SubElement(p2, f"{{{A}}}pPr")
    etree.SubElement(p2Pr, f"{{{A}}}defRPr")

    # Gene name run (RED #FF0000)
    r2 = etree.SubElement(p2, f"{{{A}}}r")
    rPr2 = etree.SubElement(r2, f"{{{A}}}rPr")
    rPr2.set("lang", "en-US")
    rPr2.set("sz", "800")
    rPr2.set("b", "1")
    sf2 = etree.SubElement(rPr2, f"{{{A}}}solidFill")
    etree.SubElement(sf2, f"{{{A}}}srgbClr").set("val", "FF0000")
    lat2 = etree.SubElement(rPr2, f"{{{A}}}latin")
    lat2.set("typeface", "Arial")
    cs2 = etree.SubElement(rPr2, f"{{{A}}}cs")
    cs2.set("typeface", "Arial")
    t2 = etree.SubElement(r2, f"{{{A}}}t")
    t2.text = gene

    # /HK_gene run (black)
    r3 = etree.SubElement(p2, f"{{{A}}}r")
    rPr3 = etree.SubElement(r3, f"{{{A}}}rPr")
    rPr3.set("lang", "en-US")
    rPr3.set("sz", "800")
    rPr3.set("b", "1")
    sf3 = etree.SubElement(rPr3, f"{{{A}}}solidFill")
    sc3 = etree.SubElement(sf3, f"{{{A}}}sysClr")
    sc3.set("val", "windowText")
    sc3.set("lastClr", "000000")
    lat3 = etree.SubElement(rPr3, f"{{{A}}}latin")
    lat3.set("typeface", "Arial")
    cs3 = etree.SubElement(rPr3, f"{{{A}}}cs")
    cs3.set("typeface", "Arial")
    t3 = etree.SubElement(r3, f"{{{A}}}t")
    t3.text = f"/{hk_gene}"

    etree.SubElement(title_el, f"{{{C}}}overlay").set("val", "0")


# ==================== UI ====================
render_sidebar_rail()

st.title("qPCR Analysis Suite")
st.caption("Gene-by-gene analysis with efficacy-specific workflows")

# Main tabs
tab1, tab_qc, tab2, tab3, tab4, tab5 = st.tabs(
    [
        "Upload",
        "QC Check",
        "Mapping",
        "Analysis",
        "Graphs",
        "Export",
    ]
)

# ==================== TAB 1: UPLOAD & FILTER ====================
with tab1:
    render_step_indicator("upload")
    st.header("Upload & Filter Data")

    uploaded_files = st.file_uploader(
        "Upload qPCR CSV files", type=["csv"], accept_multiple_files=True
    )

    if uploaded_files:
        # FIX-01: Use content hash instead of filename to detect re-uploads
        # This catches the case where the same filename is uploaded with different data
        current_file_hashes = sorted(
            [hashlib.md5(f.getvalue()).hexdigest() for f in uploaded_files]
        )
        # Reset file pointers after hashing
        for f in uploaded_files:
            f.seek(0)
        previous_file_hashes = st.session_state.get("_uploaded_file_hashes", [])
        is_new_upload = current_file_hashes != previous_file_hashes

        if is_new_upload:
            all_data = []
            for file in uploaded_files:
                parsed = QPCRParser.parse(file)
                # FIX-03: Reject empty DataFrames (parsed but no valid rows)
                if parsed is not None and not parsed.empty:
                    parsed["Source_File"] = file.name
                    all_data.append(parsed)
                    st.success(f"{file.name}: {len(parsed)} wells parsed")
                elif parsed is not None and parsed.empty:
                    st.warning(f"{file.name}: parsed successfully but contained no valid data rows.")

            if all_data:
                st.session_state.data = pd.concat(all_data, ignore_index=True)

                # FIX-02: Complete state reset on new upload to prevent stale data
                had_previous_analysis = bool(st.session_state.get("processed_data"))
                st.session_state.processed_data = {}
                st.session_state.graphs = {}
                st.session_state.sample_mapping = {}
                st.session_state.excluded_wells = {}
                st.session_state.excluded_wells_history = []
                st.session_state.excluded_samples = set()
                st.session_state.hk_gene = None
                st.session_state.selected_efficacy = None
                st.session_state.gene_display_names = {}
                st.session_state.experiment_desc = {}
                st.session_state.qc_reviewed = False
                for key in [
                    "selected_gene_idx",
                    "analysis_cmp_condition",
                    "analysis_cmp_condition_2",
                    "analysis_cmp_condition_3",
                    "analysis_ref_condition",
                    "mapping_finalized",
                    "_exclusion_snapshot",
                    "_last_ref_sample_key",
                    "_last_cmp_sample_key",
                    "_last_cmp_sample_key_2",
                    "_last_cmp_sample_key_3",
                    "hk_select",
                    "hk_select_manual",
                ]:
                    if key in st.session_state:
                        del st.session_state[key]
                # Re-initialize graph_settings to defaults
                st.session_state.graph_settings = {
                    "color_scheme": "plotly_white",
                    "title_size": 20,
                    "font_size": 14,
                    "sig_font_size": 18,
                    "figure_height": 16,
                    "figure_width": 28,
                    "show_error": True,
                    "show_significance": True,
                    "show_grid": True,
                    "show_legend": False,
                    "xlabel": "Condition",
                    "ylabel": "Relative mRNA Expression Level",
                    "bar_colors": {},
                    "bar_gap": 0.45,
                    "orientation": "v",
                    "bar_opacity": 0.85,
                    "marker_line_width": 1,
                    "y_log_scale": False,
                    "y_min": None,
                    "y_max": None,
                    "plot_bgcolor": "#FFFFFF",
                }

                st.session_state._uploaded_file_hashes = current_file_hashes
                if had_previous_analysis:
                    st.info("Previous analysis results cleared due to new data upload.")

                # FIX-05: Detect overlapping data across multiple files
                if len(all_data) > 1:
                    combined = st.session_state.data
                    dup_check_cols = ["Target", "Well", "Sample"]
                    duplicated_mask = combined.duplicated(subset=dup_check_cols, keep=False)
                    if duplicated_mask.any():
                        dup_rows = combined[duplicated_mask]
                        dup_sources = dup_rows.groupby(dup_check_cols)["Source_File"].apply(
                            lambda x: list(x.unique())
                        )
                        # Only warn if same combo appears in DIFFERENT files
                        cross_file_dups = dup_sources[dup_sources.apply(len) > 1]
                        if len(cross_file_dups) > 0:
                            n_dups = len(cross_file_dups)
                            st.warning(
                                f"⚠️ Found {n_dups} overlapping data point(s) across files "
                                f"(same Target+Well+Sample in different files). "
                                f"This may cause duplicated results. Consider deduplicating your input files."
                            )

                unique_samples = sorted(
                    st.session_state.data["Sample"].unique(), key=natural_sort_key
                )
                st.session_state.sample_order = unique_samples

    if st.session_state.data is not None:
        col1, col2, col3, col4 = st.columns(4)
        col1.metric("Wells", len(st.session_state.data))
        col2.metric("Samples", st.session_state.data["Sample"].nunique())
        col3.metric("Genes", st.session_state.data["Target"].nunique())

        KNOWN_HK_GENES = [
            "ACTIN",
            "B-ACTIN",
            "BACTIN",
            "GAPDH",
            "ACTB",
            "BETA-ACTIN",
            "BETAACTIN",
            "18S",
            "18S RRNA",
            "HPRT",
            "HPRT1",
            "B2M",
            "RPLP0",
            "TBP",
            "PPIA",
            "RPL13A",
            "YWHAZ",
            "SDHA",
            "HMBS",
            "UBC",
            "GUSB",
            "PGK1",
        ]
        all_genes = list(st.session_state.data["Target"].unique())
        hk_genes = [g for g in all_genes if g.upper() in KNOWN_HK_GENES]

        if hk_genes:
            default_idx = 0
            if st.session_state.get("hk_gene") in hk_genes:
                default_idx = hk_genes.index(st.session_state.hk_gene)
            st.session_state.hk_gene = st.selectbox(
                "🔬 Housekeeping Gene (auto-detected)",
                hk_genes,
                index=default_idx,
                key="hk_select",
            )
            col4.metric("HK Gene", st.session_state.hk_gene)
        else:
            st.warning(
                "No standard housekeeping gene detected. Please select one manually."
            )
            default_idx = 0
            if st.session_state.get("hk_gene") in all_genes:
                default_idx = all_genes.index(st.session_state.hk_gene)
            st.session_state.hk_gene = st.selectbox(
                "🔬 Select Housekeeping Gene",
                all_genes,
                index=default_idx,
                key="hk_select_manual",
                help="Select the reference/housekeeping gene for normalization",
            )
            col4.metric("HK Gene", st.session_state.hk_gene)

        st.subheader("Data Preview")
        st.dataframe(st.session_state.data.head(50), height=300)

        st.subheader("Data Validation")
        warnings_found = False

        data = st.session_state.data
        replicate_counts = data.groupby(["Sample", "Target"]).size()

        single_replicates = replicate_counts[replicate_counts < 2]
        if len(single_replicates) > 0:
            warnings_found = True
            st.warning(
                f"**Low Replicates**: {len(single_replicates)} sample-target combinations have only 1 replicate. Statistical analysis requires n>=2."
            )
            with st.expander("View affected samples"):
                st.dataframe(single_replicates.reset_index(name="n"))

        if st.session_state.get("hk_gene"):
            hk = st.session_state.hk_gene
            samples_with_hk = set(data[data["Target"] == hk]["Sample"].unique())
            all_samples = set(data["Sample"].unique())
            missing_hk = all_samples - samples_with_hk
            if missing_hk:
                warnings_found = True
                st.error(
                    f"❌ **Missing Housekeeping Gene**: {len(missing_hk)} samples have no {hk} data: {', '.join(str(s) for s in list(missing_hk)[:5])}{'...' if len(missing_hk) > 5 else ''}"
                )

        high_ct_count = len(data[data["CT"] > AnalysisConstants.CT_HIGH_WARNING])
        if high_ct_count > 0:
            warnings_found = True
            pct = high_ct_count / len(data) * 100
            st.warning(
                f"**High CT Values**: {high_ct_count} wells ({pct:.1f}%) have CT > {AnalysisConstants.CT_HIGH_WARNING} (low expression)"
            )

        if not warnings_found:
            st.success("Data validation passed. No issues detected.")

# ==================== TAB QC: QUALITY CONTROL ====================
with tab_qc:
    render_step_indicator("qc")
    st.header("Quality Control Check")

    # Honest QC-complete signal. Previously set True unconditionally on every
    # rerun (all tab bodies execute each run), so the workflow rail marked QC
    # "done" before the user reviewed anything. Now an explicit confirmation.
    st.session_state.qc_reviewed = st.checkbox(
        "Mark quality control as reviewed",
        value=st.session_state.get("qc_reviewed", False),
        help="Drives the workflow progress rail — check once you've reviewed replicate QC.",
    )

    if st.session_state.data is not None and not st.session_state.data.empty:
        data = st.session_state.data
        hk_gene = st.session_state.get("hk_gene")

        if "excluded_wells" not in st.session_state:
            st.session_state.excluded_wells = {}
        if "excluded_wells_history" not in st.session_state:
            st.session_state.excluded_wells_history = []

        # ==================== QC SUMMARY DASHBOARD ====================
        st.caption("Browse all CT values, review triplicates, and exclude outliers before analysis.")

        # Get comprehensive QC stats using helper to get all excluded wells
        qc_stats = QualityControl.get_qc_summary_stats(data, get_all_excluded_wells())

        # Summary metrics — 2 rows of 3 for better readability
        row1_cols = st.columns(3)
        row1_cols[0].metric("Total Wells", qc_stats.get("total_wells", 0))
        row1_cols[1].metric(
            "Active Wells",
            qc_stats.get("active_wells", 0),
            delta=f"-{qc_stats.get('excluded_wells', 0)}"
            if qc_stats.get("excluded_wells", 0) > 0
            else None,
        )
        row1_cols[2].metric("Triplicates", qc_stats.get("total_triplicates", 0))

        row2_cols = st.columns(3)
        row2_cols[0].metric(
            "Healthy",
            qc_stats.get("healthy_triplicates", 0),
            delta=f"{qc_stats.get('health_score', 0)}%"
            if qc_stats.get("health_score", 0) > 0
            else None,
            delta_color="normal",
        )
        row2_cols[1].metric(
            "Warnings",
            qc_stats.get("warning_triplicates", 0),
            delta=None if qc_stats.get("warning_triplicates", 0) == 0 else "review",
            delta_color="off",
        )
        row2_cols[2].metric(
            "Errors",
            qc_stats.get("error_triplicates", 0),
            delta=None if qc_stats.get("error_triplicates", 0) == 0 else "critical",
            delta_color="off",
        )

        # CT Distribution summary
        with st.expander("📈 CT Value Distribution", expanded=False):
            dist_cols = st.columns(4)
            dist_cols[0].metric("Mean CT", f"{qc_stats.get('ct_mean', 0):.1f}")
            dist_cols[1].metric(
                "CT Range",
                f"{qc_stats.get('ct_min', 0):.1f} - {qc_stats.get('ct_max', 0):.1f}",
            )
            dist_cols[2].metric(
                "High CT (>35)",
                qc_stats.get("high_ct_count", 0),
                delta="check" if qc_stats.get("high_ct_count", 0) > 0 else None,
                delta_color="off",
            )
            dist_cols[3].metric(
                "Avg CV%",
                f"{qc_stats.get('avg_cv_pct', 0):.1f}%",
                delta="high" if qc_stats.get("avg_cv_pct", 0) > 5 else None,
                delta_color="off" if qc_stats.get("avg_cv_pct", 0) > 5 else "normal",
            )

        # ==================== MAIN QC INTERFACE WITH TABS ====================
        qc_tab1, qc_tab2 = st.tabs(
            [
                "Triplicate Browser",
                "QC Overview",
            ]
        )

        # ==================== TAB 1: ENHANCED TRIPLICATE BROWSER ====================
        with qc_tab1:
            st.subheader("Per-Well Selection Browser")
            st.caption(
                "Browse triplicates, exclude outliers, and adjust QC settings — all in one place."
            )

            # ---- Inline QC Settings (collapsed) ----
            with st.expander("QC Settings", expanded=False):
                thresh_col1, thresh_col2 = st.columns(2)
                with thresh_col1:
                    new_ct_high = st.number_input(
                        "High CT Threshold",
                        min_value=25.0, max_value=45.0,
                        value=float(QualityControl.CT_HIGH_THRESHOLD),
                        step=0.5,
                        help="Wells with CT above this are flagged as low expression",
                        key="qc_settings_ct_high",
                    )
                    QualityControl.CT_HIGH_THRESHOLD = new_ct_high

                    new_ct_low = st.number_input(
                        "Low CT Threshold",
                        min_value=5.0, max_value=20.0,
                        value=float(QualityControl.CT_LOW_THRESHOLD),
                        step=0.5,
                        help="Wells with CT below this are flagged as unusually high expression",
                        key="qc_settings_ct_low",
                    )
                    QualityControl.CT_LOW_THRESHOLD = new_ct_low

                with thresh_col2:
                    new_cv = st.number_input(
                        "CV% Threshold",
                        min_value=1.0, max_value=20.0,
                        value=float(QualityControl.CV_THRESHOLD * 100),
                        step=0.5,
                        help="Replicates with CV above this are flagged as high variability",
                        key="qc_settings_cv",
                    )
                    QualityControl.CV_THRESHOLD = new_cv / 100

                    new_hk_var = st.number_input(
                        "HK Variation Threshold",
                        min_value=0.5, max_value=3.0,
                        value=float(QualityControl.HK_VARIATION_THRESHOLD),
                        step=0.1,
                        help="Housekeeping gene deviation threshold across samples",
                        key="qc_settings_hk_var",
                    )
                    QualityControl.HK_VARIATION_THRESHOLD = new_hk_var

                if st.button("🔄 Re-run QC with New Settings", type="primary", use_container_width=True, key="qc_rerun_settings"):
                    st.rerun()

            # ---- Auto-Exclude High SD Outliers ----
            st.markdown("---")
            auto_excl_cols = st.columns([2, 1, 1])
            with auto_excl_cols[0]:
                sd_thresh = st.slider(
                    "SD Threshold (Ct)",
                    min_value=0.3, max_value=1.0, value=0.5, step=0.1,
                    key="sd_threshold_slider",
                    help="Triplicates with CT SD above this will be flagged",
                )
            with auto_excl_cols[1]:
                find_all_btn = st.button(
                    "🔍 Find All Outliers", key="find_all_sd_outliers", use_container_width=True
                )
            with auto_excl_cols[2]:
                # Per-gene button — read gene filter from previous render cycle
                _prev_gene = st.session_state.get("qc_gene_filter", "All Genes")
                if _prev_gene and _prev_gene != "All Genes":
                    clean_gene_btn = st.button(
                        f"🧹 Clean {_prev_gene}", key="clean_gene_btn", use_container_width=True
                    )
                else:
                    clean_gene_btn = False

            # Handle find buttons
            if find_all_btn:
                suggestions = QualityControl.find_high_sd_outliers(
                    data, st.session_state.get("excluded_wells", {}), sd_threshold=sd_thresh,
                )
                st.session_state["_sd_outlier_suggestions"] = suggestions

            if clean_gene_btn:
                suggestions = QualityControl.find_high_sd_outliers(
                    data, st.session_state.get("excluded_wells", {}),
                    sd_threshold=sd_thresh, gene_filter=_prev_gene,
                )
                st.session_state["_sd_outlier_suggestions"] = suggestions

            # Display preview table
            suggestions = st.session_state.get("_sd_outlier_suggestions", [])
            if suggestions:
                st.markdown(f"**Found {len(suggestions)} high-SD outlier(s):**")
                preview_df = pd.DataFrame(suggestions)
                preview_df.insert(0, "Exclude", True)
                edited = st.data_editor(
                    preview_df[["Exclude", "Target", "Sample", "Well", "CT", "deviation", "group_sd", "n_replicates"]],
                    disabled=["Target", "Sample", "Well", "CT", "deviation", "group_sd", "n_replicates"],
                    key="sd_outlier_preview",
                    use_container_width=True,
                )

                apply_cols = st.columns([3, 1])
                with apply_cols[1]:
                    if st.button("Apply Selected", key="apply_sd_exclusions", type="primary", use_container_width=True):
                        selected = edited[edited["Exclude"]]
                        # Save undo history
                        if "excluded_wells_history" in st.session_state:
                            st.session_state.excluded_wells_history.append(
                                {k: v.copy() for k, v in st.session_state.excluded_wells.items()}
                            )
                        excl_dict = st.session_state.get("excluded_wells", {})
                        if not isinstance(excl_dict, dict):
                            excl_dict = {}
                        for _, row in selected.iterrows():
                            key = (row["Target"], row["Sample"])
                            if key not in excl_dict:
                                excl_dict[key] = set()
                            excl_dict[key].add(row["Well"])
                        st.session_state.excluded_wells = excl_dict
                        st.session_state.pop("_sd_outlier_suggestions", None)
                        st.success(f"Excluded {len(selected)} well(s).")
                        st.rerun()
                with apply_cols[0]:
                    if st.button("Dismiss", key="dismiss_sd_suggestions"):
                        st.session_state.pop("_sd_outlier_suggestions", None)
                        st.rerun()
            elif find_all_btn or clean_gene_btn:
                st.info("No triplicates exceed the SD threshold.")

            st.markdown("---")

            # ---- SD-based flagging for color-coded well labels ----
            # Flag wells whose triplicate SD exceeds the threshold
            # IMPORTANT: filter out already-excluded wells before computing SD
            _flagged_well_lookup = {}
            _flagged_genes = {}
            _excl_wells_flat = set()
            if isinstance(st.session_state.excluded_wells, dict):
                for ws in st.session_state.excluded_wells.values():
                    _excl_wells_flat.update(ws)
            _active_data = data[~data["Well"].isin(_excl_wells_flat)] if _excl_wells_flat else data

            for (gene_t, sample_t), grp in _active_data.groupby(["Target", "Sample"]):
                if len(grp) < 2:
                    continue
                grp_sd = grp["CT"].std(ddof=1)
                if pd.isna(grp_sd) or grp_sd <= sd_thresh:
                    continue
                # Find worst well (highest deviation from group mean)
                grp_mean = grp["CT"].mean()
                deviations = (grp["CT"] - grp_mean).abs()
                worst_idx = deviations.idxmax()
                worst_well = grp.loc[worst_idx, "Well"]
                _flagged_well_lookup[(gene_t, sample_t, worst_well)] = (
                    "error" if grp_sd > sd_thresh * 1.5 else "warning",
                    f"SD {grp_sd:.3f} > {sd_thresh}",
                )
                _flagged_genes[gene_t] = _flagged_genes.get(gene_t, 0) + 1

            if _flagged_well_lookup:
                _gene_parts = ", ".join(
                    f"{g} ({c})" for g, c in sorted(_flagged_genes.items())
                )
                st.warning(
                    f"⚠️ **{len(_flagged_well_lookup)} high-SD outliers** (SD > {sd_thresh}): {_gene_parts} — "
                    f"Click **Find All Outliers** above to review and exclude."
                )

            # ---- Filter Controls ----
            filter_col1, filter_col2 = st.columns([1, 1])

            all_genes = sorted(data["Target"].unique())
            all_samples = sorted(data["Sample"].unique(), key=natural_sort_key)

            with filter_col1:
                selected_gene_filter = st.selectbox(
                    "Filter by Gene", ["All Genes"] + all_genes, key="qc_gene_filter"
                )

            with filter_col2:
                selected_sample_filter = st.selectbox(
                    "Filter by Sample",
                    ["All Samples"] + all_samples,
                    key="qc_sample_filter",
                )

            # Build a combined dataframe of all wells with gene-sample context
            wells_all = data[["Well", "Sample", "Target", "CT"]].copy()
            wells_all = wells_all.sort_values(["Target", "Sample", "Well"])

            # Apply filters
            if selected_gene_filter != "All Genes":
                wells_all = wells_all[wells_all["Target"] == selected_gene_filter]
            if selected_sample_filter != "All Samples":
                wells_all = wells_all[wells_all["Sample"] == selected_sample_filter]

            if wells_all.empty:
                st.info("No wells match the current filters.")
            else:
                # Get genes to display
                display_genes = sorted(wells_all["Target"].unique())

                # Count excluded
                total_excluded = sum(
                    len(ws) for ws in st.session_state.excluded_wells.values()
                )
                st.markdown(
                    f"**{len(wells_all)} wells across {len(display_genes)} genes** | "
                    f"Excluded: {total_excluded}"
                )

                # Global quick actions
                action_cols = st.columns(3)
                with action_cols[0]:
                    if st.button("Include All Visible", key="qc_incl_all_visible", use_container_width=True):
                        st.session_state.excluded_wells_history.append(
                            {k: v.copy() for k, v in st.session_state.excluded_wells.items()}
                        )
                        for _, r in wells_all.iterrows():
                            include_well(r["Well"], r["Target"], r["Sample"])
                        st.rerun()
                with action_cols[1]:
                    if st.button("Exclude All Visible", key="qc_excl_all_visible", use_container_width=True):
                        st.session_state.excluded_wells_history.append(
                            {k: v.copy() for k, v in st.session_state.excluded_wells.items()}
                        )
                        for _, r in wells_all.iterrows():
                            exclude_well(r["Well"], r["Target"], r["Sample"])
                        st.rerun()
                with action_cols[2]:
                    can_undo = len(st.session_state.excluded_wells_history) > 0
                    if st.button("Undo Last Change", key="qc_undo_browser", use_container_width=True, disabled=not can_undo):
                        if st.session_state.excluded_wells_history:
                            st.session_state.excluded_wells = st.session_state.excluded_wells_history.pop()
                            st.rerun()

                # Render per-gene expandable sections
                for gene in display_genes:
                    gene_wells = wells_all[wells_all["Target"] == gene]
                    gene_samples = sorted(gene_wells["Sample"].unique(), key=natural_sort_key)

                    # Count excluded for this gene
                    gene_excluded = sum(
                        1 for _, r in gene_wells.iterrows()
                        if is_well_excluded(r["Well"], gene, r["Sample"])
                    )

                    _gene_flag_count = _flagged_genes.get(gene, 0)
                    if gene_excluded > 0 and _gene_flag_count > 0:
                        gene_label = f"{gene}  ({len(gene_wells)} wells, {gene_excluded} excluded, ⚠️ {_gene_flag_count} high-SD)"
                    elif gene_excluded > 0:
                        gene_label = f"{gene}  ({len(gene_wells)} wells, {gene_excluded} excluded)"
                    elif _gene_flag_count > 0:
                        gene_label = f"{gene}  ({len(gene_wells)} wells, ⚠️ {_gene_flag_count} high-SD)"
                    else:
                        gene_label = f"{gene}  ({len(gene_wells)} wells)"

                    _auto_expand = (selected_gene_filter != "All Genes") or (_gene_flag_count > 0)
                    with st.expander(gene_label, expanded=_auto_expand):
                        # Build editable dataframe for this gene
                        gene_editor_df = gene_wells.copy()
                        gene_editor_df["Include"] = gene_editor_df.apply(
                            lambda r: not is_well_excluded(r["Well"], gene, r["Sample"]),
                            axis=1,
                        )
                        # Compute triplicate mean and deviation per sample
                        sample_means = gene_wells.groupby("Sample")["CT"].transform("mean")
                        gene_editor_df["Deviation"] = (gene_editor_df["CT"] - sample_means).round(3)
                        gene_editor_df["CT"] = gene_editor_df["CT"].round(2)

                        # Reorder
                        gene_editor_df = gene_editor_df[["Include", "Well", "Sample", "CT", "Deviation"]]

                        with st.form(key=f"well_form_{gene}"):
                            st.caption("Check/uncheck wells to include or exclude, then click Apply.")
                            checkbox_states = {}
                            for s_idx, sample_name in enumerate(gene_samples):
                                sample_rows = gene_editor_df[gene_editor_df["Sample"] == sample_name]
                                n_total = len(sample_rows)
                                n_included = int(sample_rows["Include"].sum())
                                included_cts = pd.to_numeric(
                                    sample_rows[sample_rows["Include"] == True]["CT"], errors="coerce"
                                )
                                if len(included_cts) >= 2:
                                    s_mean = included_cts.mean()
                                    s_sd = included_cts.std(ddof=1)
                                    s_cv = (s_sd / s_mean * 100) if s_mean != 0 else 0.0
                                    stats_str = f"Mean {s_mean:.2f} · SD {s_sd:.3f} · CV {s_cv:.1f}%"
                                elif len(included_cts) == 1:
                                    stats_str = f"Mean {included_cts.mean():.2f} · SD N/A"
                                else:
                                    stats_str = "No wells included"
                                # Check if any well in this sample has high SD
                                _sample_has_flag = any(
                                    (gene, sample_name, r["Well"]) in _flagged_well_lookup
                                    for _, r in sample_rows.iterrows()
                                )
                                _sample_flag_indicator = " ⚠️ high SD" if _sample_has_flag else ""
                                st.markdown(
                                    f'<div style="padding:8px 0 4px 0;font-weight:600;font-size:0.9rem;color:#1d1d1f;">'
                                    f'{sample_name}{_sample_flag_indicator}'
                                    f'<span style="font-weight:400;color:#86868b;font-size:0.8rem;margin-left:8px;">'
                                    f'{n_included}/{n_total} included · {stats_str}</span></div>',
                                    unsafe_allow_html=True,
                                )
                                for idx, row in sample_rows.iterrows():
                                    well = row["Well"]
                                    ct_val = row["CT"]
                                    dev_val = row["Deviation"]
                                    # Color-coded prefix based on QC flag severity
                                    _well_flag = _flagged_well_lookup.get((gene, sample_name, well))
                                    if _well_flag is not None:
                                        _sev, _issues = _well_flag
                                        if _sev == "error":
                                            _prefix = "🔴 "
                                        elif _sev == "warning":
                                            _prefix = "🟡 "
                                        else:
                                            _prefix = ""
                                        label = f"{_prefix}{well}  ·  CT {ct_val:.2f}  ·  Dev {dev_val:+.3f}  ·  {_issues}"
                                    else:
                                        label = f"{well}  ·  CT {ct_val:.2f}  ·  Dev {dev_val:+.3f}"
                                    checkbox_states[idx] = st.checkbox(
                                        label,
                                        value=row["Include"],
                                        key=f"cb_{gene}_{well}_{sample_name}_{idx}",
                                    )
                                if s_idx < len(gene_samples) - 1:
                                    st.markdown(
                                        '<hr style="margin:8px 0;border:none;border-top:1px solid #e5e5e7;">',
                                        unsafe_allow_html=True,
                                    )
                            submitted = st.form_submit_button("Apply Changes")

                        if submitted:
                            changed = False
                            for idx, include in checkbox_states.items():
                                row = gene_editor_df.loc[idx]
                                well = row["Well"]
                                sample = row["Sample"]
                                if not include and not is_well_excluded(well, gene, sample):
                                    if not changed:
                                        st.session_state.excluded_wells_history.append(
                                            {k: v.copy() for k, v in st.session_state.excluded_wells.items()}
                                        )
                                        changed = True
                                    exclude_well(well, gene, sample)
                                elif include and is_well_excluded(well, gene, sample):
                                    if not changed:
                                        st.session_state.excluded_wells_history.append(
                                            {k: v.copy() for k, v in st.session_state.excluded_wells.items()}
                                        )
                                        changed = True
                                    include_well(well, gene, sample)
                            if changed:
                                st.rerun()



        # ==================== TAB 2: QC OVERVIEW (CONSOLIDATED) ====================
        with qc_tab2:
            st.subheader("QC Overview")
            st.caption("Plate heatmap, auto-flagged wells, and pre-analysis summary at a glance.")

            # ---- Section 1: Plate Heatmap ----
            st.markdown("### Plate Heatmap")
            heatmap_col1, heatmap_col2 = st.columns([2, 1])

            with heatmap_col1:
                plate_fig = QualityControl.create_plate_heatmap(
                    data, value_col="CT", excluded_wells=get_all_excluded_wells()
                )
                st.plotly_chart(plate_fig, use_container_width=True)
                st.caption(
                    "Red = High CT (low expression) | Green = Low CT (high expression) | X = Excluded"
                )

                # Per-gene heatmap
                all_genes = sorted(data["Target"].dropna().unique().tolist())
                if all_genes:
                    selected_gene = st.selectbox(
                        "Filter heatmap by gene",
                        options=["(All)"] + all_genes,
                        key="heatmap_gene_select",
                    )
                    if selected_gene != "(All)":
                        gene_data = data[data["Target"] == selected_gene]
                        gene_fig = QualityControl.create_plate_heatmap(
                            gene_data, value_col="CT", excluded_wells=get_all_excluded_wells()
                        )
                        gene_fig.update_layout(title=f"Plate Heatmap — {selected_gene}")
                        st.plotly_chart(gene_fig, use_container_width=True)

                # Per-sample heatmap
                all_samples = sorted(data["Sample"].dropna().unique().tolist(), key=natural_sort_key)
                if all_samples:
                    selected_sample = st.selectbox(
                        "Filter heatmap by sample",
                        options=["(All)"] + all_samples,
                        key="heatmap_sample_select",
                    )
                    if selected_sample != "(All)":
                        sample_data = data[data["Sample"] == selected_sample]
                        sample_fig = QualityControl.create_plate_heatmap(
                            sample_data, value_col="CT", excluded_wells=get_all_excluded_wells()
                        )
                        sample_fig.update_layout(title=f"Plate Heatmap — {selected_sample}")
                        st.plotly_chart(sample_fig, use_container_width=True)

            with heatmap_col2:
                st.markdown("**Replicate Statistics**")
                rep_stats = QualityControl.get_replicate_stats(data)
                if not rep_stats.empty:
                    def highlight_status(row):
                        if row["Status"] == "High CV":
                            return ["background-color: #fff3cd"] * len(row)
                        elif row["Status"] == "Low Expression":
                            return ["background-color: #f8d7da"] * len(row)
                        elif row["Status"] == "Check Signal":
                            return ["background-color: #cce5ff"] * len(row)
                        return [""] * len(row)

                    styled_stats = rep_stats.style.apply(highlight_status, axis=1)
                    st.dataframe(styled_stats, height=400, use_container_width=True)

            # ---- Section 2: Flagged Wells (SD-based, matches Triplicate Browser) ----
            st.markdown("### Flagged Wells")

            # Use same SD-based flagging logic as the Triplicate Browser
            _ov_excl_flat = set()
            if isinstance(st.session_state.excluded_wells, dict):
                for _ws in st.session_state.excluded_wells.values():
                    _ov_excl_flat.update(_ws)
            _ov_active = data[~data["Well"].isin(_ov_excl_flat)] if _ov_excl_flat else data

            _ov_flagged_rows = []
            for (gene_t, sample_t), grp in _ov_active.groupby(["Target", "Sample"]):
                if len(grp) < 2:
                    continue
                grp_sd = grp["CT"].std(ddof=1)
                if pd.isna(grp_sd) or grp_sd <= sd_thresh:
                    continue
                grp_mean = grp["CT"].mean()
                deviations = (grp["CT"] - grp_mean).abs()
                worst_idx = deviations.idxmax()
                worst_row = grp.loc[worst_idx]
                severity = "error" if grp_sd > sd_thresh * 1.5 else "warning"
                _ov_flagged_rows.append({
                    "Well": worst_row["Well"],
                    "Sample": sample_t,
                    "Target": gene_t,
                    "CT": worst_row["CT"],
                    "Issues": f"SD {grp_sd:.3f} > {sd_thresh}",
                    "Severity": severity,
                })

            if _ov_flagged_rows:
                flagged_overview = pd.DataFrame(_ov_flagged_rows)
                flag_ov_col1, flag_ov_col2 = st.columns([3, 1])
                with flag_ov_col1:
                    st.warning(f"Found {len(flagged_overview)} wells with high triplicate SD")
                with flag_ov_col2:
                    if st.button(
                        "Exclude All Flagged",
                        use_container_width=True,
                        key="qc_overview_exclude_flagged",
                    ):
                        st.session_state.excluded_wells_history.append(
                            {k: v.copy() for k, v in st.session_state.excluded_wells.items()}
                        )
                        for _, row in flagged_overview.iterrows():
                            exclude_well(row["Well"], row["Target"], row["Sample"])
                        st.rerun()

                st.dataframe(
                    flagged_overview,
                    hide_index=True,
                    use_container_width=True,
                    column_config={
                        "CT": st.column_config.NumberColumn("CT", format="%.2f"),
                    },
                )
            else:
                st.success("No quality issues detected. All wells pass SD threshold.")

            # ---- Section 3: Pre-Analysis Summary ----
            st.markdown("### Pre-Analysis Summary")
            st.caption(
                "Verify which wells will be used in analysis for each gene-sample combination."
            )

            all_genes_summary = sorted(data["Target"].unique())
            all_samples_summary = sorted(data["Sample"].unique(), key=natural_sort_key)

            hk_gene_name = st.session_state.get("hk_gene", None)

            summary_rows = []
            for gene in all_genes_summary:
                is_hk = False
                if hk_gene_name and gene.upper() == hk_gene_name.upper():
                    is_hk = True

                for sample in all_samples_summary:
                    gene_sample_wells = data[
                        (data["Target"] == gene) & (data["Sample"] == sample)
                    ]
                    if gene_sample_wells.empty:
                        continue

                    total_wells = len(gene_sample_wells)
                    well_ids = gene_sample_wells["Well"].tolist()

                    excluded_set = st.session_state.excluded_wells.get(
                        (gene, sample), set()
                    )
                    excluded_ids = [w for w in well_ids if w in excluded_set]
                    included_ids = [w for w in well_ids if w not in excluded_set]

                    n_included = len(included_ids)
                    n_excluded = len(excluded_ids)

                    included_cts = gene_sample_wells[
                        gene_sample_wells["Well"].isin(included_ids)
                    ]["CT"]
                    mean_ct = included_cts.mean() if n_included > 0 else None
                    sd_ct = (
                        included_cts.std() if n_included > 1 else 0.0
                    )

                    summary_rows.append(
                        {
                            "Gene": gene,
                            "Role": "HK" if is_hk else "Target",
                            "Sample": sample,
                            "Total Wells": total_wells,
                            "Included": n_included,
                            "Excluded": n_excluded,
                            "Included Wells": ", ".join(str(w) for w in included_ids),
                            "Excluded Wells": ", ".join(str(w) for w in excluded_ids) if excluded_ids else "—",
                            "Mean CT": round(mean_ct, 2) if mean_ct is not None else "—",
                            "CT SD": round(sd_ct, 3) if sd_ct is not None else "—",
                        }
                    )

            if summary_rows:
                summary_df = pd.DataFrame(summary_rows)

                total_exclusions = summary_df["Excluded"].sum()
                low_rep = summary_df[summary_df["Included"] < 2]

                sum_col1, sum_col2, sum_col3 = st.columns(3)
                sum_col1.metric("Gene-Sample Groups", len(summary_df))
                sum_col2.metric(
                    "Total Well-Exclusions",
                    int(total_exclusions),
                    delta=f"-{int(total_exclusions)}" if total_exclusions > 0 else None,
                    delta_color="inverse" if total_exclusions > 0 else "off",
                )
                sum_col3.metric(
                    "Low Replicate (n<2)",
                    len(low_rep),
                    delta="check" if len(low_rep) > 0 else None,
                    delta_color="off",
                )

                if len(low_rep) > 0:
                    st.warning(
                        f"{len(low_rep)} gene-sample group(s) have fewer than 2 included wells. "
                        "Statistics (SD, SEM, p-values) will be unreliable or zero."
                    )

                for gene in all_genes_summary:
                    gene_rows = summary_df[summary_df["Gene"] == gene]
                    if gene_rows.empty:
                        continue

                    gene_excl = gene_rows["Excluded"].sum()
                    gene_label = f"{gene}  (HK)" if gene_rows.iloc[0]["Role"] == "HK" else gene
                    if gene_excl > 0:
                        gene_label += f"  — {int(gene_excl)} exclusion(s)"

                    with st.expander(gene_label, expanded=(gene_excl > 0)):
                        display_cols = [
                            "Sample", "Included", "Excluded",
                            "Included Wells", "Excluded Wells",
                            "Mean CT", "CT SD",
                        ]
                        st.dataframe(
                            gene_rows[display_cols].reset_index(drop=True),
                            use_container_width=True,
                            hide_index=True,
                        )
            else:
                st.info("No data available. Upload and parse data first.")

        # ==================== BOTTOM STATUS BAR ====================
        # Count total excluded wells across all gene-sample combinations
        excluded_count = sum(
            len(wells_set) for wells_set in st.session_state.excluded_wells.values()
        )

        status_cols = st.columns([2, 1, 1])
        with status_cols[0]:
            if excluded_count > 0:
                st.info(
                    f"ℹ️ **{excluded_count} wells** will be excluded from analysis. Review in Triplicate Browser or proceed to Mapping tab."
                )
            else:
                st.success(
                    "✅ QC check complete. All wells will be included. Proceed to Mapping tab."
                )

        with status_cols[1]:
            can_undo = len(st.session_state.excluded_wells_history) > 0
            if st.button(
                "↩️ Undo Last",
                use_container_width=True,
                disabled=not can_undo,
                key="global_undo",
            ):
                if st.session_state.excluded_wells_history:
                    st.session_state.excluded_wells = (
                        st.session_state.excluded_wells_history.pop()
                    )
                    st.rerun()

        with status_cols[2]:
            if st.button(
                "🔄 Refresh QC", use_container_width=True, key="global_refresh"
            ):
                st.rerun()

    else:
        st.info("Upload data first in the Upload tab.")

# ==================== TAB 2: SAMPLE MAPPING ====================
with tab2:
    render_step_indicator("mapping")
    st.header("Map Samples to Conditions")

    if st.session_state.data is not None:
        # Efficacy type selection
        detected_genes = set(st.session_state.data["Target"].unique())
        # FIX-07: Case-insensitive efficacy gene matching
        detected_genes_upper = {g.upper() for g in detected_genes}
        suggested = None
        for eff, cfg in EFFICACY_CONFIG.items():
            if any(g.upper() in detected_genes_upper for g in cfg["genes"]):
                suggested = eff
                break

        efficacy = st.selectbox(
            "🎯 Efficacy Test Type",
            list(EFFICACY_CONFIG.keys()),
            index=list(EFFICACY_CONFIG.keys()).index(suggested) if suggested else 0,
        )
        st.session_state.selected_efficacy = efficacy

        config = EFFICACY_CONFIG[efficacy]
        st.info(f"**{config['description']}**")
        st.caption(f"Cell line: {config['cell']} | Genes: {', '.join(config['genes'])}")

        # Show control structure
        with st.expander("📋 Control Structure for this Test"):
            for ctrl_type, ctrl_name in config["controls"].items():
                st.markdown(f"- **{ctrl_type.title()}**: {ctrl_name}")

        # Sample mapping interface with professional layout
        st.markdown("### Sample Condition Mapping")

        # FIX-04: Always sync sample_order with actual data to prevent desync
        # New samples from data that aren't in sample_order yet get appended
        current_data_samples = set(st.session_state.data["Sample"].unique())
        if "sample_order" not in st.session_state or not st.session_state.sample_order:
            st.session_state.sample_order = sorted(
                current_data_samples, key=natural_sort_key
            )
        else:
            # Add any new samples not yet in sample_order (preserving user's custom order)
            existing_order = st.session_state.sample_order
            new_samples = sorted(
                [s for s in current_data_samples if s not in existing_order],
                key=natural_sort_key,
            )
            if new_samples:
                st.session_state.sample_order = existing_order + new_samples
            # Remove samples no longer in data
            st.session_state.sample_order = [
                s for s in st.session_state.sample_order if s in current_data_samples
            ]

        # Group type options
        group_types = ["Negative Control", "Positive Control", "Treatment"]
        if "baseline" in config["controls"]:
            group_types.insert(0, "Baseline")

        # Smart default group: match a new sample's name against the efficacy
        # category's known control labels so e.g. "Non-treated" pre-selects
        # "Negative Control" instead of "Treatment". Always overridable below.
        _role_to_group = {
            "baseline": "Baseline",
            "negative": "Negative Control",
            "positive": "Positive Control",
        }

        def _suggest_group(sample_name):
            sl = str(sample_name).strip().lower()
            if not sl:
                return "Treatment"
            controls = config.get("controls", {})
            # Exact label match first (highest confidence), then substring.
            for role in ("baseline", "negative", "positive"):
                label = controls.get(role)
                if label and sl == str(label).strip().lower():
                    grp = _role_to_group[role]
                    return grp if grp in group_types else "Treatment"
            for role in ("baseline", "negative", "positive"):
                label = controls.get(role)
                ll = str(label).strip().lower() if label else ""
                if ll and (ll in sl or sl in ll):
                    grp = _role_to_group[role]
                    return grp if grp in group_types else "Treatment"
            return "Treatment"

        # Ensure all samples in sample_order have mapping
        for sample in st.session_state.sample_order:
            if sample not in st.session_state.sample_mapping:
                st.session_state.sample_mapping[sample] = {
                    "condition": sample,
                    "group": _suggest_group(sample),
                    "concentration": "",
                    "include": True,
                }
            if "include" not in st.session_state.sample_mapping[sample]:
                st.session_state.sample_mapping[sample]["include"] = True

        # Header row matching column proportions below
        hdr0, hdr_ord, hdr1, hdr2, hdr3 = st.columns([0.6, 0.5, 2, 3.5, 2.5])
        hdr0.markdown("**Use**")
        hdr_ord.markdown("**#**")
        hdr1.markdown("**Original**")
        hdr2.markdown("**Condition Name**")
        hdr3.markdown("**Group**")

        # FIXED: Display ALL samples in sample_order (including excluded ones)
        display_samples = st.session_state.sample_order.copy()

        # Sample rows with improved spacing
        for i, sample in enumerate(display_samples):
            # FIX-04: Guard against samples not in mapping (desync protection)
            if sample not in st.session_state.sample_mapping:
                continue
            # Container for each row
            with st.container():
                col0, col_order, col1, col2, col3 = st.columns(
                    [0.6, 0.5, 2, 3.5, 2.5]
                )

                # Include checkbox
                with col0:
                    include = st.checkbox(
                        "Include sample",
                        value=st.session_state.sample_mapping[sample].get(
                            "include", True
                        ),
                        key=f"include_{sample}_{i}",
                        label_visibility="collapsed",
                    )
                    st.session_state.sample_mapping[sample]["include"] = include

                # Order number
                with col_order:
                    st.markdown(
                        f"<div style='text-align: center; padding-top: 10px;'><b>{i + 1}</b></div>",
                        unsafe_allow_html=True,
                    )

                # Original sample name (non-editable)
                with col1:
                    st.text_input(
                        "Original",
                        sample,
                        key=f"orig_{sample}_{i}",
                        disabled=True,
                        label_visibility="collapsed",
                    )

                # Condition name (editable)
                with col2:
                    cond = st.text_input(
                        "Condition",
                        st.session_state.sample_mapping[sample]["condition"],
                        key=f"cond_{sample}_{i}",
                        label_visibility="collapsed",
                        placeholder="Enter condition name...",
                        max_chars=50,
                    )
                    st.session_state.sample_mapping[sample]["condition"] = cond

                # Group selector
                with col3:
                    grp_idx = 0
                    previous_group = st.session_state.sample_mapping[sample].get("group", "Treatment")
                    try:
                        grp_idx = group_types.index(previous_group)
                    except (ValueError, KeyError):
                        # FIX-08: Warn when group type resets due to efficacy change
                        st.caption(f"⚠️ '{previous_group}' not available; reset to '{group_types[0]}'")
                        st.session_state.sample_mapping[sample]["group"] = group_types[0]

                    grp = st.selectbox(
                        "Group",
                        group_types,
                        index=grp_idx,
                        key=f"grp_{sample}_{i}",
                        label_visibility="collapsed",
                    )
                    st.session_state.sample_mapping[sample]["group"] = grp

                st.markdown("")

        if st.button("Finalize Mapping", type="primary", use_container_width=True, key="finalize_mapping"):
            st.session_state.mapping_finalized = True
            st.session_state.analysis_stale = False
            st.rerun()
        st.caption("Lock in condition names and groups, then run analysis below.")

        if st.session_state.get("mapping_finalized", False):
            edit_col1, edit_col2 = st.columns([3, 1])
            with edit_col2:
                if st.button("✏️ Edit Mapping", key="edit_mapping", use_container_width=True):
                    st.session_state.mapping_finalized = False
                    st.session_state.analysis_stale = True
                    st.rerun()
            st.markdown("### 🔀 Drag to Reorder Samples")
            st.caption("Drag included samples to set the display order for graphs and exports.")

            current_order = st.session_state.sample_order
            included_order = [
                s for s in current_order
                if st.session_state.sample_mapping.get(s, {}).get("include", True)
            ]
            excluded_samples_list = [
                s for s in current_order
                if not st.session_state.sample_mapping.get(s, {}).get("include", True)
            ]

            reorder_mode = st.radio(
                "Reorder method",
                ["Arrow buttons (fast)", "Drag-and-drop"],
                horizontal=True,
                key="reorder_mode",
                label_visibility="collapsed",
            )

            if reorder_mode == "Arrow buttons (fast)":
                for i, sample in enumerate(included_order):
                    cond = st.session_state.sample_mapping.get(sample, {}).get("condition", sample)
                    cond_short = cond if len(cond) <= 25 else cond[:22] + "..."
                    sample_short = sample if len(sample) <= 20 else sample[:17] + "..."
                    lbl = f"{cond_short}  ({sample_short})" if cond != sample else sample_short
                    btn_cols = st.columns([0.5, 0.5, 5])
                    with btn_cols[0]:
                        if i > 0 and st.button("\u2191", key=f"up_{sample}", help="Move up"):
                            included_order[i], included_order[i - 1] = included_order[i - 1], included_order[i]
                            st.session_state.sample_order = included_order + excluded_samples_list
                            st.rerun()
                    with btn_cols[1]:
                        if i < len(included_order) - 1 and st.button("\u2193", key=f"dn_{sample}", help="Move down"):
                            included_order[i], included_order[i + 1] = included_order[i + 1], included_order[i]
                            st.session_state.sample_order = included_order + excluded_samples_list
                            st.rerun()
                    with btn_cols[2]:
                        st.markdown(f"<span style='font-size:0.9rem;'>{lbl}</span>", unsafe_allow_html=True)
            elif sort_items is not None:
                # Build unique labels: "Condition Name  (Original)" format
                order_labels = []
                label_to_sample = {}
                for s in included_order:
                    cond = st.session_state.sample_mapping.get(s, {}).get("condition", s)
                    label = f"{cond}  ({s})" if cond != s else s
                    if label in label_to_sample:
                        label = f"{cond}  ({s}) #{included_order.index(s) + 1}"
                    order_labels.append(label)
                    label_to_sample[label] = s

                st.markdown("""
                <style>
                .sortable-item {
                    padding: 6px 14px !important;
                    margin: 2px 0 !important;
                    border-radius: 8px !important;
                    font-size: 0.85rem !important;
                    max-width: 400px !important;
                }
                </style>
                """, unsafe_allow_html=True)

                if not order_labels:
                    st.info("No included samples to reorder. Include samples above first.")
                    sorted_labels = []
                else:
                    col_drag, col_spacer = st.columns([2, 3])
                    with col_drag:
                        sorted_labels = sort_items(items=order_labels, direction="vertical")

                new_included_order = [label_to_sample[lbl] for lbl in sorted_labels]
                new_order = new_included_order + excluded_samples_list

                if new_order != st.session_state.sample_order:
                    st.session_state.sample_order = new_order
            else:
                st.warning(
                    "Drag-and-drop requires `streamlit-sortables`. "
                    "Install with: `pip install streamlit-sortables`"
                )

        # Update excluded_samples from include flags
        st.session_state.excluded_samples = set(
            [
                s
                for s, v in st.session_state.sample_mapping.items()
                if not v.get("include", True)
            ]
        )

        # FIX-11: Warn on duplicate condition names across included samples
        included_conditions = [
            v["condition"]
            for s, v in st.session_state.sample_mapping.items()
            if v.get("include", True)
        ]
        seen_conditions = {}
        for cond in included_conditions:
            seen_conditions[cond] = seen_conditions.get(cond, 0) + 1
        duplicate_conditions = {c: n for c, n in seen_conditions.items() if n > 1}
        if duplicate_conditions:
            dup_list = ", ".join(f"'{c}' (×{n})" for c, n in duplicate_conditions.items())
            st.info(
                f"ℹ️ Duplicate condition names detected: {dup_list}. "
                f"Samples with the same condition name will be treated as biological replicates."
            )

        # Summary with styled cards
        st.subheader("Mapping Summary")

        col_card1, col_card2, col_card3, col_card4 = st.columns(4)

        total_samples = len(st.session_state.sample_order)
        included = sum(
            1
            for s in st.session_state.sample_order
            if st.session_state.sample_mapping[s].get("include", True)
        )
        excluded = total_samples - included

        with col_card1:
            st.metric("Total Samples", total_samples)
        with col_card2:
            st.metric(
                "Included",
                included,
                delta=None if included == total_samples else f"-{excluded}",
            )
        with col_card3:
            st.metric(
                "Excluded", excluded, delta=None if excluded == 0 else f"+{excluded}"
            )
        with col_card4:
            groups = set(
                v["group"]
                for v in st.session_state.sample_mapping.values()
                if v.get("include", True)
            )
            st.metric("Groups", len(groups))

        # Detailed table view
        with st.expander("📋 View Detailed Mapping Table"):
            mapping_df = pd.DataFrame(
                [
                    {
                        "Order": idx + 1,
                        "Include": "✅"
                        if st.session_state.sample_mapping[s].get("include", True)
                        else "❌",
                        "Original": s,
                        "Condition": st.session_state.sample_mapping[s]["condition"],
                        "Group": st.session_state.sample_mapping[s]["group"],
                    }
                    for idx, s in enumerate(st.session_state.sample_order)
                ]
            )
            st.dataframe(mapping_df, use_container_width=True, hide_index=True)

        # Run analysis
        st.subheader("Run Full Analysis (DDCt + Statistics)")

        # Build condition list from mapping
        condition_list = []
        sample_to_condition = {}

        for sample in st.session_state.get("sample_order", []):
            if sample in st.session_state.sample_mapping:
                mapping_info = st.session_state.sample_mapping[sample]
                if mapping_info.get("include", True):
                    condition = mapping_info.get("condition", sample)
                    condition_list.append(condition)
                    # Keep first sample for each condition (not last)
                    if condition not in sample_to_condition:
                        sample_to_condition[condition] = sample
        # Deduplicate condition_list while preserving order
        seen = set()
        condition_list = [c for c in condition_list if not (c in seen or seen.add(c))]

        if condition_list:
            # Enhanced layout with clear separation
            st.markdown("#### 📊 Analysis Configuration")

            col_info1, col_info2 = st.columns(2)
            with col_info1:
                st.info(
                    "**ΔΔCt Reference:** Used to calculate fold changes. All samples will be relative to this (Fold Change = 1.0)"
                )
            with col_info2:
                st.info(
                    "**P-value References:** Used for statistical comparison (t-test). Choose one or two conditions for comparison."
                )

            col_r1, col_r2, col_r3 = st.columns(3)
            with col_r1:
                ref_condition = st.selectbox(
                    "🎯 ΔΔCt Reference Condition",
                    condition_list,
                    index=0,
                    key="ref_choice_ddct",
                    help="Baseline for relative expression calculation",
                )
                ref_sample_key = sample_to_condition[ref_condition]
                st.caption(f"→ Sample: **{ref_sample_key}**")

            with col_r2:
                cmp_condition = st.selectbox(
                    "📈 P-value Reference 1 (*)",
                    condition_list,
                    index=0,
                    key="cmp_choice_pval",
                    help="Primary control group for statistical testing (asterisk symbols)",
                )
                cmp_sample_key = sample_to_condition[cmp_condition]
                st.caption(f"→ Sample: **{cmp_sample_key}**")

            with col_r3:
                # Add option for second p-value comparison
                use_second_comparison = st.checkbox(
                    "Enable 2nd comparison (#)",
                    value=False,
                    key="use_second_pval",
                    help="Add a second statistical comparison with hashtag symbols",
                )

                if use_second_comparison:
                    condition_list_2 = [c for c in condition_list if c != cmp_condition]
                    if condition_list_2:
                        cmp_condition_2 = st.selectbox(
                            "P-value Reference 2 (#)",
                            condition_list_2,
                            index=0,
                            key="cmp_choice_pval_2",
                            help="Secondary control group for statistical testing (hashtag symbols)",
                        )
                        cmp_sample_key_2 = sample_to_condition[cmp_condition_2]
                        st.caption(f"→ Sample: **{cmp_sample_key_2}**")
                    else:
                        st.warning("Need at least 3 conditions for dual comparison")
                        use_second_comparison = False
                        cmp_sample_key_2 = None
                else:
                    cmp_sample_key_2 = None

                # Third comparison (†)
                use_third_comparison = st.checkbox(
                    "Enable 3rd comparison (†)",
                    value=False,
                    key="use_third_pval",
                    help="Add a third statistical comparison with dagger symbols",
                ) if use_second_comparison else False

                if use_third_comparison:
                    used = {cmp_condition, cmp_condition_2} if use_second_comparison and cmp_sample_key_2 else {cmp_condition}
                    condition_list_3 = [c for c in condition_list if c not in used]
                    if condition_list_3:
                        cmp_condition_3 = st.selectbox(
                            "P-value Reference 3 (†)",
                            condition_list_3,
                            index=0,
                            key="cmp_choice_pval_3",
                            help="Third control group for statistical testing (dagger symbols)",
                        )
                        cmp_sample_key_3 = sample_to_condition[cmp_condition_3]
                        st.caption(f"→ Sample: **{cmp_sample_key_3}**")
                    else:
                        st.warning("Need at least 4 conditions for triple comparison")
                        use_third_comparison = False
                        cmp_sample_key_3 = None
                else:
                    cmp_sample_key_3 = None

            # Statistical options
            st.markdown("#### ⚙️ Statistical Options")

            # Statistics Information Box
            with st.expander("ℹ️ Understanding Statistics Options", expanded=False):
                st.markdown("""
                **📊 Standard Deviation (SD)**  
                - Calculated using Excel-like `=STDEV()` function on CT values
                - Reflects **sample selection** from QC Check tab (excluded wells are NOT included)
                - Shows **data variability** within replicates (2-3 values typically)
                - Formula: `SD = sqrt(Σ(xi - mean)² / (n-1))`
                - Displayed in results as `Target_Ct_SD`, `HK_Ct_SD`, and error bar option
                
                **📈 Coefficient of Variation (CV%)**  
                - Calculated as: `CV% = (SD / Mean) × 100`
                - Normalized measure of variability (allows comparison across different CT ranges)
                - Lower CV% indicates better reproducibility
                
                **🎯 P-value Testing**  
                - Statistical significance between conditions
                - Uses t-test to compare treatment vs control groups
                - Symbols: `*` p<0.05, `**` p<0.01, `***` p<0.001
                - Respects QC exclusions (uses only included wells)
                
                **📏 Error Bars**  
                - **SEM** (Standard Error of Mean): Shows precision of mean estimate = `SD / √n`
                - **SD** (Standard Deviation): Shows data spread/variability
                - Choose SEM for publication (smaller bars, shows reliability)
                - Choose SD to show full data variation
                """)

            stat_col1, stat_col2 = st.columns(2)

            with stat_col1:
                ttest_type = st.radio(
                    "T-test Type",
                    ["welch", "student"],
                    format_func=lambda x: "Welch's t-test (unequal variance)"
                    if x == "welch"
                    else "Student's t-test (equal variance)",
                    index=0,
                    key="ttest_type",
                    help="Welch's is recommended - more robust when group variances differ",
                )
                # Note: Widget with key="ttest_type" auto-syncs to session state

            with stat_col2:
                error_bar_type = st.radio(
                    "Error Bar Type",
                    ["sem", "sd"],
                    format_func=lambda x: "SEM (Standard Error of Mean)"
                    if x == "sem"
                    else "SD (Standard Deviation)",
                    index=0,
                    key="error_bar_type",
                    help="SEM shows precision of mean estimate; SD shows data variability",
                )
                # Note: Widget with key="error_bar_type" auto-syncs to session state

            # Visual summary
            col_sum1, col_sum2, col_sum3 = st.columns([1, 2, 1])
            with col_sum2:
                summary_html = f"""
                <div style='background: #fafafa; padding: 20px; border-radius: 12px; text-align: center; border: 1px solid #f0f0f0;'>
                    <h4>Analysis Summary</h4>
                    <p><b>Fold Changes:</b> Relative to <code>{ref_condition}</code></p>
                    <p><b>P-values (*):</b> Compared to <code>{cmp_condition}</code></p>
                """
                if use_second_comparison and cmp_sample_key_2:
                    summary_html += f"<p><b>P-values (#):</b> Compared to <code>{cmp_condition_2}</code></p>"
                if use_third_comparison and cmp_sample_key_3:
                    summary_html += f"<p><b>P-values (†):</b> Compared to <code>{cmp_condition_3}</code></p>"
                summary_html += "</div>"
                st.markdown(summary_html, unsafe_allow_html=True)

            # Store analysis parameters for auto-rerun
            st.session_state['_last_ref_sample_key'] = ref_sample_key
            st.session_state['_last_cmp_sample_key'] = cmp_sample_key
            st.session_state['_last_cmp_sample_key_2'] = cmp_sample_key_2 if use_second_comparison else None
            st.session_state['_last_cmp_sample_key_3'] = cmp_sample_key_3 if use_third_comparison else None

            # Run button
            if st.button("Run Full Analysis", type="primary", use_container_width=True):
                ok = AnalysisEngine.run_full_analysis(
                    ref_sample_key,
                    cmp_sample_key,
                    cmp_sample_key_2 if use_second_comparison else None,
                    cmp_sample_key_3 if use_third_comparison else None,
                )
                if ok:
                    success_msg = f"Analysis complete!\n\n- Fold changes relative to: **{ref_condition}**\n- P-values (*) vs: **{cmp_condition}**"
                    if use_second_comparison and cmp_sample_key_2:
                        success_msg += f"\n- P-values (#) vs: **{cmp_condition_2}**"
                    if use_third_comparison and cmp_sample_key_3:
                        success_msg += f"\n- P-values (†) vs: **{cmp_condition_3}**"
                    st.session_state.analysis_stale = False
                    st.success(success_msg)
                    st.rerun()
                else:
                    st.error("Analysis failed. Check messages above.")
        else:
            st.warning("No samples available for analysis.")

# ==================== TAB 3: ANALYSIS ====================
with tab3:
    # Auto-rerun if exclusion state changed since last analysis
    if (st.session_state.get('processed_data') and
        '_exclusion_snapshot' in st.session_state and
        '_last_ref_sample_key' in st.session_state):

        current_snapshot = {
            'excluded_wells': {str(k): sorted(v) for k, v in st.session_state.get('excluded_wells', {}).items()},
            'excluded_samples': sorted(st.session_state.get('excluded_samples', set())),
            'ttest_type': st.session_state.get('ttest_type', 'welch')
        }

        if current_snapshot != st.session_state['_exclusion_snapshot']:
            AnalysisEngine.run_full_analysis(
                st.session_state['_last_ref_sample_key'],
                st.session_state['_last_cmp_sample_key'],
                st.session_state.get('_last_cmp_sample_key_2'),
                st.session_state.get('_last_cmp_sample_key_3'),
            )
            st.info("Analysis auto-updated to reflect changes in QC exclusions or settings.")

    if st.session_state.get("analysis_stale", False):
        st.warning("⚠️ Mapping changed since last analysis. Results may be outdated.")
        if st.button("Re-run Analysis", key="rerun_stale_analysis", type="primary"):
            ref_key = st.session_state.get("_last_ref_sample_key")
            cmp_key = st.session_state.get("_last_cmp_sample_key")
            if ref_key and cmp_key:
                ok = AnalysisEngine.run_full_analysis(ref_key, cmp_key,
                    st.session_state.get("_last_cmp_sample_key_2"),
                    st.session_state.get("_last_cmp_sample_key_3"))
                if ok:
                    st.session_state.analysis_stale = False
                    st.rerun()

    render_step_indicator("analysis")
    st.header("Analysis Results")

    if st.session_state.processed_data:
        st.subheader("Analysis Summary")

        # Summary metrics
        all_results = pd.concat(
            st.session_state.processed_data.values(), ignore_index=True
        )

        col1, col2, col3, col4 = st.columns(4)
        col1.metric("Genes", len(st.session_state.processed_data))
        col2.metric("Conditions", all_results["Condition"].nunique())
        sig_count = (all_results["p_value"] < 0.05).sum()
        total_testable = all_results["p_value"].notna().sum()
        col3.metric("Sig. (p<0.05)", f"{sig_count}/{total_testable}")
        excluded_well_count = sum(len(ws) for ws in st.session_state.excluded_wells.values()) if isinstance(st.session_state.excluded_wells, dict) else len(st.session_state.excluded_wells)
        col4.metric("Excluded", excluded_well_count)

        # Reproducibility / provenance record (MIQE-style) — travels with exports.
        with st.expander("🔬 Analysis Provenance (reproducibility record)", expanded=False):
            st.code(format_provenance_text(_current_provenance()), language=None)
            st.caption(
                "This record — method, reference gene/condition, comparisons, test, "
                "and every excluded well — is downloadable on the Export tab."
            )

        # Show results per gene
        st.subheader("Gene-by-Gene Results")

        for gene, gene_df in st.session_state.processed_data.items():
            with st.expander(f"📍 {gene}", expanded=False):
                # Show expected direction if available
                efficacy_config = EFFICACY_CONFIG.get(
                    st.session_state.selected_efficacy, {}
                )
                if "expected_direction" in efficacy_config:
                    direction = efficacy_config["expected_direction"].get(gene)
                    if direction:
                        st.caption(
                            f"Expected: {'↑ Increase' if direction == 'up' else '↓ Decrease'}"
                        )

                # Display columns
                display_cols = [
                    "Condition",
                    "Group",
                    "Fold_Change",
                    "p_value",
                    "significance",
                    "n_replicates",
                    "Target_Ct_Mean",
                    "Target_Ct_SD",
                    "Target_Ct_CV%",
                    "HK_Ct_Mean",
                    "HK_Ct_SD",
                    "HK_Ct_CV%",
                    "Delta_Ct",
                    "SEM",
                    "SD",
                ]

                # Filter to existing columns
                display_df = gene_df[[c for c in display_cols if c in gene_df.columns]]

                # Style the dataframe
                styled = display_df.style.background_gradient(
                    subset=["Fold_Change"], cmap="RdYlGn", vmin=0, vmax=3
                ).format(
                    {
                        "Fold_Change": "{:.3f}",
                        "p_value": "{:.4f}",
                        "Target_Ct_Mean": "{:.2f}",
                        "Target_Ct_SD": "{:.3f}",
                        "Target_Ct_CV%": "{:.1f}%",
                        "HK_Ct_Mean": "{:.2f}",
                        "HK_Ct_SD": "{:.3f}",
                        "HK_Ct_CV%": "{:.1f}%",
                        "Delta_Ct": "{:.2f}",
                        "SEM": "{:.3f}",
                        "SD": "{:.3f}",
                    },
                    na_rep="—",
                )

                st.dataframe(styled, use_container_width=True)
                st.download_button(
                    "⬇️ Download CSV",
                    data=display_df.to_csv(index=False).encode("utf-8-sig"),
                    file_name=f"{gene}_results_{datetime.now().strftime('%Y%m%d')}.csv",
                    mime="text/csv", key=f"csv_{gene}", use_container_width=True,
                )

        st.success("Results ready. Go to **Graphs** tab to visualize.")

    else:
        st.info(
            "⏳ No analysis results yet. Go to the 'Mapping' tab and click 'Run Full Analysis'."
        )

# ==================== TAB 4: GRAPHS ====================
with tab4:
    render_step_indicator("graphs")
    st.header("Individual Gene Graphs")

    if st.session_state.get("analysis_stale", False):
        st.warning("⚠️ Mapping changed since last analysis. Go to Analysis tab to re-run.")

    # Graph tab styles handled by global theme

    if st.session_state.processed_data:
        if "graph_settings" not in st.session_state:
            st.session_state.graph_settings = {
                "title_size": 20,
                "font_size": 14,
                "sig_font_size": 18,
                "figure_width": 28,
                "figure_height": 16,
                "color_scheme": "plotly_white",
                "show_error": True,
                "show_significance": True,
                "show_grid": True,
                "xlabel": "Condition",
                "ylabel": "Relative mRNA Expression Level",
                "bar_colors": {},
                "orientation": "v",
                "bar_opacity": 0.85,
                "bar_gap": 0.45,
                "marker_line_width": 1,
                "show_legend": False,
                "y_log_scale": False,
                "y_min": None,
                "y_max": None,
            }

        if "graphs" not in st.session_state:
            st.session_state.graphs = {}

        efficacy_config = EFFICACY_CONFIG.get(st.session_state.selected_efficacy, {})
        gene_list = list(st.session_state.processed_data.keys())

        if "selected_gene_idx" not in st.session_state:
            st.session_state.selected_gene_idx = 0
        if st.session_state.selected_gene_idx >= len(gene_list):
            st.session_state.selected_gene_idx = 0

        # Gene pill selector with container
        st.markdown('<div class="gene-pill-container"><div class="label">Select Gene</div></div>', unsafe_allow_html=True)
        gene_cols = st.columns(min(len(gene_list), 6)) if gene_list else [st.container()]
        for idx, gene in enumerate(gene_list):
            with gene_cols[idx % len(gene_cols)]:
                if st.button(
                    f"{'✓ ' if idx == st.session_state.selected_gene_idx else ''}{gene}",
                    key=f"gene_btn_{gene}",
                    use_container_width=True,
                    type="primary"
                    if idx == st.session_state.selected_gene_idx
                    else "secondary",
                ):
                    st.session_state.selected_gene_idx = idx
                    st.rerun()

        current_gene = gene_list[st.session_state.selected_gene_idx]

        gene_data = st.session_state.processed_data[current_gene]

        show_sig_key = f"{current_gene}_show_sig"
        show_err_key = f"{current_gene}_show_err"
        bar_gap_key = f"{current_gene}_bar_gap"

        if show_sig_key not in st.session_state.graph_settings:
            st.session_state.graph_settings[show_sig_key] = True
        if show_err_key not in st.session_state.graph_settings:
            st.session_state.graph_settings[show_err_key] = True
        if bar_gap_key not in st.session_state.graph_settings:
            st.session_state.graph_settings[bar_gap_key] = 0.45

        conditions_list = gene_data["Condition"].tolist()
        ref_options = ["None"] + conditions_list
        ref_line_key = f"{current_gene}_ref_line"
        if ref_line_key not in st.session_state.graph_settings:
            st.session_state.graph_settings[ref_line_key] = "None"
        if st.session_state.graph_settings.get(ref_line_key, "None") not in ref_options:
            st.session_state.graph_settings[ref_line_key] = "None"

        # Initialize color preset key
        color_preset_key = f"{current_gene}_color_preset"
        if color_preset_key not in st.session_state.graph_settings:
            st.session_state.graph_settings[color_preset_key] = "Classic"

        # Initialize data points toggle key
        show_dp_key = f"{current_gene}_show_data_points"
        if show_dp_key not in st.session_state.graph_settings:
            st.session_state.graph_settings[show_dp_key] = False

        # Toolbar row 1: preset, toggles, bar gap, reset
        tb_row1 = st.columns([1.5, 1.0, 1.0, 1.0, 1.5, 0.7])
        with tb_row1[0]:
            preset_names = list(GRAPH_PRESETS.keys()) + ["Custom"]
            current_preset = st.session_state.graph_settings.get(color_preset_key, "Classic")
            if current_preset not in preset_names:
                current_preset = "Custom"
            selected_preset = st.selectbox(
                "Color Preset",
                preset_names,
                index=preset_names.index(current_preset),
                key=f"preset_{current_gene}",
            )
            st.session_state.graph_settings[color_preset_key] = selected_preset
        with tb_row1[1]:
            sig_on = st.toggle(
                "Sig. */#/\u2020",
                st.session_state.graph_settings[show_sig_key],
                key=f"tgl_sig_{current_gene}",
            )
            st.session_state.graph_settings[show_sig_key] = sig_on
        with tb_row1[2]:
            err_on = st.toggle(
                "Error Bars",
                st.session_state.graph_settings[show_err_key],
                key=f"tgl_err_{current_gene}",
            )
            st.session_state.graph_settings[show_err_key] = err_on
        with tb_row1[3]:
            dp_on = st.toggle(
                "Data Points",
                st.session_state.graph_settings[show_dp_key],
                key=f"tgl_dp_{current_gene}",
            )
            st.session_state.graph_settings[show_dp_key] = dp_on
        with tb_row1[4]:
            gap_val = st.select_slider(
                "Bar Gap",
                options=[round(x * 0.05, 2) for x in range(2, 21)],  # 0.1 to 1.0 step 0.05
                value=st.session_state.graph_settings[bar_gap_key],
                key=f"gap_sl_{current_gene}",
            )
            st.session_state.graph_settings[bar_gap_key] = gap_val
        with tb_row1[5]:
            if st.button("Reset", key=f"reset_all_{current_gene}", use_container_width=True):
                if f"{current_gene}_bar_settings" in st.session_state:
                    del st.session_state[f"{current_gene}_bar_settings"]
                st.session_state.graph_settings[show_sig_key] = True
                st.session_state.graph_settings[show_err_key] = True
                st.session_state.graph_settings[bar_gap_key] = 0.45
                st.session_state.graph_settings[show_dp_key] = False
                st.session_state.graph_settings.pop(color_preset_key, None)
                for per_gene_key in [
                    f"{current_gene}_figure_width", f"{current_gene}_figure_height",
                    f"{current_gene}_font_size", f"{current_gene}_bar_opacity",
                    f"{current_gene}_marker_line_width", f"{current_gene}_tick_size",
                    f"{current_gene}_ylabel_size", f"{current_gene}_bg_color",
                    f"{current_gene}_y_min", f"{current_gene}_y_max",
                    f"{current_gene}_label_mode",
                ]:
                    st.session_state.graph_settings.pop(per_gene_key, None)
                st.rerun()

        # Toolbar row 2: ref line + label mode
        label_mode_key = f"{current_gene}_label_mode"
        label_modes = ["Auto-wrap", "Angled 45\u00b0", "Angled 90\u00b0", "Horizontal"]
        if label_mode_key not in st.session_state.graph_settings:
            st.session_state.graph_settings[label_mode_key] = "Auto-wrap"

        tb_row2 = st.columns([2, 2])
        with tb_row2[0]:
            selected_ref = st.selectbox(
                "Reference Line",
                ref_options,
                index=ref_options.index(st.session_state.graph_settings.get(ref_line_key, "None"))
                if st.session_state.graph_settings.get(ref_line_key, "None") in ref_options
                else 0,
                key=f"ref_sel_{current_gene}",
                help="Horizontal dashed line at a condition's expression level",
            )
            st.session_state.graph_settings[ref_line_key] = selected_ref
        with tb_row2[1]:
            label_mode = st.selectbox(
                "Label Mode",
                label_modes,
                index=label_modes.index(st.session_state.graph_settings.get(label_mode_key, "Auto-wrap")),
                key=f"lbl_mode_{current_gene}",
                help="How x-axis labels handle long text",
            )
            st.session_state.graph_settings[label_mode_key] = label_mode

        if f"{current_gene}_bar_settings" not in st.session_state:
            st.session_state[f"{current_gene}_bar_settings"] = {}
        if "bar_colors_per_sample" not in st.session_state.graph_settings:
            st.session_state.graph_settings["bar_colors_per_sample"] = {}

        for idx, (_, row) in enumerate(gene_data.iterrows()):
            condition = row["Condition"]
            bar_key = f"{current_gene}_{condition}"
            if bar_key not in st.session_state[f"{current_gene}_bar_settings"]:
                st.session_state[f"{current_gene}_bar_settings"][bar_key] = {
                    "color": "#FFFFFF",
                    "show_sig": True,
                    "show_sig_1": True,
                    "show_sig_2": True,
                    "show_sig_3": True,
                    "show_err": True,
                }
            # Migrate old entries missing per-symbol keys
            bs = st.session_state[f"{current_gene}_bar_settings"][bar_key]
            for sk in ("show_sig_1", "show_sig_2", "show_sig_3"):
                if sk not in bs:
                    bs[sk] = bs.get("show_sig", True)

        with st.expander("Settings & Colors", expanded=False):
                s_col1, s_col2, s_col3, s_col4 = st.columns(4)

                with s_col1:
                    st.caption("**Display**")
                    gene_display = st.text_input(
                        "Gene display name",
                        value=st.session_state.gene_display_names.get(current_gene, current_gene),
                        key=f"gene_display_{current_gene}",
                        placeholder=current_gene,
                    )
                    gene_display_stripped = (gene_display.strip() if gene_display else "")[:50]
                    if gene_display_stripped and gene_display_stripped != current_gene:
                        st.session_state.gene_display_names[current_gene] = gene_display_stripped
                    elif current_gene in st.session_state.gene_display_names:
                        del st.session_state.gene_display_names[current_gene]

                    bg_color = st.color_picker(
                        "Background",
                        st.session_state.graph_settings.get(f"{current_gene}_bg_color",
                            st.session_state.graph_settings.get("plot_bgcolor", "#FFFFFF")),
                        key=f"bg_{current_gene}",
                    )
                    st.session_state.graph_settings[f"{current_gene}_bg_color"] = bg_color

                with s_col2:
                    st.caption("**Figure Size**")
                    size_preset_key = f"{current_gene}_size_preset"
                    size_preset_names = list(FIGURE_SIZE_PRESETS.keys()) + ["Custom"]
                    if size_preset_key not in st.session_state.graph_settings:
                        st.session_state.graph_settings[size_preset_key] = "PPT Full"
                    current_size_preset = st.session_state.graph_settings.get(size_preset_key, "PPT Full")
                    if current_size_preset not in size_preset_names:
                        current_size_preset = "Custom"
                    selected_size = st.selectbox(
                        "Size Preset",
                        size_preset_names,
                        index=size_preset_names.index(current_size_preset),
                        key=f"size_preset_{current_gene}",
                    )
                    prev_size_preset = st.session_state.graph_settings.get(size_preset_key, "PPT Full")
                    if selected_size != "Custom" and selected_size in FIGURE_SIZE_PRESETS:
                        preset = FIGURE_SIZE_PRESETS[selected_size]
                        st.session_state.graph_settings[f"{current_gene}_figure_width"] = preset["width"]
                        st.session_state.graph_settings[f"{current_gene}_figure_height"] = preset["height"]
                    st.session_state.graph_settings[size_preset_key] = selected_size
                    if selected_size != prev_size_preset:
                        st.rerun()
                    fig_width_cm = st.slider(
                        "Width (cm)", 10.0, 40.0,
                        value=float(st.session_state.graph_settings.get(f"{current_gene}_figure_width",
                            st.session_state.graph_settings.get("figure_width", 28))),
                        step=0.5, key=f"fig_w_{current_gene}",
                    )
                    fig_height_cm = st.slider(
                        "Height (cm)", 6.0, 25.0,
                        value=float(st.session_state.graph_settings.get(f"{current_gene}_figure_height",
                            st.session_state.graph_settings.get("figure_height", 16))),
                        step=0.5, key=f"fig_h_{current_gene}",
                    )
                    st.session_state.graph_settings[f"{current_gene}_figure_width"] = fig_width_cm
                    st.session_state.graph_settings[f"{current_gene}_figure_height"] = fig_height_cm
                    if selected_size != "Custom" and selected_size in FIGURE_SIZE_PRESETS:
                        preset = FIGURE_SIZE_PRESETS[selected_size]
                        if fig_width_cm != preset["width"] or fig_height_cm != preset["height"]:
                            st.session_state.graph_settings[size_preset_key] = "Custom"

                with s_col3:
                    st.caption("**Fonts**")
                    global_font = st.slider(
                        "Global font", 8, 28,
                        value=st.session_state.graph_settings.get(f"{current_gene}_font_size",
                            st.session_state.graph_settings.get("font_size", 14)),
                        key=f"gf_{current_gene}",
                    )
                    st.session_state.graph_settings[f"{current_gene}_font_size"] = global_font
                    tick_size = st.slider(
                        "X-tick labels", 8, 24,
                        value=st.session_state.graph_settings.get(f"{current_gene}_tick_size", 12),
                        key=f"ts_{current_gene}",
                    )
                    st.session_state.graph_settings[f"{current_gene}_tick_size"] = tick_size
                    ylabel_size = st.slider(
                        "Y-axis label", 8, 24,
                        value=st.session_state.graph_settings.get(f"{current_gene}_ylabel_size", 14),
                        key=f"ys_{current_gene}",
                    )
                    st.session_state.graph_settings[f"{current_gene}_ylabel_size"] = ylabel_size

                with s_col4:
                    st.caption("**Bar Style**")
                    bar_opacity = st.slider(
                        "Bar opacity", 0.3, 1.0,
                        value=float(st.session_state.graph_settings.get(f"{current_gene}_bar_opacity",
                            st.session_state.graph_settings.get("bar_opacity", 0.85))),
                        step=0.05, key=f"bo_{current_gene}",
                    )
                    st.session_state.graph_settings[f"{current_gene}_bar_opacity"] = bar_opacity
                    outline_width = st.slider(
                        "Outline width", 0, 3,
                        value=int(st.session_state.graph_settings.get(f"{current_gene}_marker_line_width",
                            st.session_state.graph_settings.get("marker_line_width", 1))),
                        key=f"ow_{current_gene}",
                    )
                    st.session_state.graph_settings[f"{current_gene}_marker_line_width"] = outline_width
                    y_min_key = f"{current_gene}_y_min"
                    y_max_key = f"{current_gene}_y_max"
                    y_min_input = st.number_input(
                        "Y-axis min",
                        value=st.session_state.graph_settings.get(y_min_key),
                        placeholder="Auto",
                        key=f"ymin_{current_gene}",
                        help="Leave empty for auto range",
                    )
                    y_max_input = st.number_input(
                        "Y-axis max",
                        value=st.session_state.graph_settings.get(y_max_key),
                        placeholder="Auto",
                        key=f"ymax_{current_gene}",
                        help="Leave empty for auto range",
                    )
                    if (y_min_input is not None and y_max_input is not None
                            and y_max_input > 0 and y_min_input >= y_max_input):
                        st.warning("Y-axis min must be less than max. Using auto range.")
                        y_min_input = None
                        y_max_input = None
                    if y_min_input is not None:
                        st.session_state.graph_settings[y_min_key] = y_min_input
                    elif y_min_key in st.session_state.graph_settings:
                        del st.session_state.graph_settings[y_min_key]
                    if y_max_input is not None:
                        st.session_state.graph_settings[y_max_key] = y_max_input
                    elif y_max_key in st.session_state.graph_settings:
                        del st.session_state.graph_settings[y_max_key]

                st.markdown("---")

                st.markdown("**Per-Bar Settings**")
                # Sort gene_data to match graph bar order (sample_order → condition)
                _bar_display_data = gene_data.copy()
                _so = st.session_state.get("sample_order", [])
                _sm = st.session_state.get("sample_mapping", {})
                if _so and _sm:
                    _cond_order = []
                    _seen = set()
                    for _s in _so:
                        if _sm.get(_s, {}).get("include", True):
                            _c = _sm.get(_s, {}).get("condition", _s)
                            if _c in _bar_display_data["Condition"].values and _c not in _seen:
                                _cond_order.append(_c)
                                _seen.add(_c)
                    for _c in _bar_display_data["Condition"].unique():
                        if _c not in _seen:
                            _cond_order.append(_c)
                    _bar_display_data["Condition"] = pd.Categorical(
                        _bar_display_data["Condition"], categories=_cond_order, ordered=True
                    )
                    _bar_display_data = _bar_display_data.sort_values("Condition").reset_index(drop=True)

                # Header row
                hdr = st.columns([3, 0.8, 2.5])
                hdr[0].markdown("<small>**Condition**</small>", unsafe_allow_html=True)
                hdr[1].markdown("<small>**Color**</small>", unsafe_allow_html=True)
                with hdr[2]:
                    _opt_hdr = st.columns(4)
                    for _hi, _hlbl in enumerate([
                        "<small>**✱**<br><span style='color:#888'>Sig.1</span></small>",
                        "<small>**#**<br><span style='color:#888'>Sig.2</span></small>",
                        "<small>**†**<br><span style='color:#888'>Sig.3</span></small>",
                        "<small>**±**<br><span style='color:#888'>Err</span></small>",
                    ]):
                        _opt_hdr[_hi].markdown(_hlbl, unsafe_allow_html=True)

                option_labels = ["✱", "#", "†", "±"]
                option_keys = ["show_sig_1", "show_sig_2", "show_sig_3", "show_err"]

                for idx, (_, row) in enumerate(_bar_display_data.iterrows()):
                    condition = row["Condition"]
                    group = row.get("Group", "Treatment")
                    bar_key = f"{current_gene}_{condition}"
                    bs = st.session_state[f"{current_gene}_bar_settings"][bar_key]

                    rc = st.columns([3, 0.8, 2.5])
                    lbl = condition if len(condition) <= 22 else condition[:19] + "..."
                    rc[0].markdown(
                        f"<small>{lbl} <span style='color:#888;'>({group})</span></small>",
                        unsafe_allow_html=True,
                    )
                    # Color picker — show preset tone or custom
                    _active_pn = st.session_state.graph_settings.get(f"{current_gene}_color_preset", "Classic")
                    _cp_key = f"cp_{current_gene}_{condition}"
                    _desired_key = f"_desired_{_cp_key}"
                    if _active_pn != "Custom" and _active_pn in GRAPH_PRESETS:
                        _preset = GRAPH_PRESETS[_active_pn]
                        _is_ref = (condition == st.session_state.get("analysis_ref_condition"))
                        _display_color = _preset["ref"] if _is_ref else _preset["color"]
                        # Sync widget to preset color only when preset (or ref detection) changes.
                        # This avoids overwriting user interactions while keeping the picker in
                        # sync when the user switches presets.
                        if st.session_state.get(_desired_key) != _display_color:
                            st.session_state[_cp_key] = _display_color
                            st.session_state[_desired_key] = _display_color
                    else:
                        _display_color = st.session_state.graph_settings.get(
                            "bar_colors_per_sample", {}
                        ).get(bar_key, bs.get("color", "#FFFFFF"))

                    new_color = rc[1].color_picker(
                        "c", _display_color,
                        key=_cp_key,
                        label_visibility="collapsed",
                    )

                    # Update state based on current mode
                    if _active_pn != "Custom":
                        if new_color != _display_color:
                            # User picked a different color → switch to Custom mode
                            st.session_state.graph_settings[f"{current_gene}_color_preset"] = "Custom"
                            st.session_state[_desired_key] = new_color
                            bs["color"] = new_color
                            st.session_state.graph_settings.setdefault("bar_colors_per_sample", {})[bar_key] = new_color
                        else:
                            # No change: keep bs and bar_colors in sync with preset
                            bs["color"] = _display_color
                            st.session_state.graph_settings.setdefault("bar_colors_per_sample", {})[bar_key] = _display_color
                    else:
                        # Custom mode: widget value is authoritative
                        st.session_state[_desired_key] = new_color
                        bs["color"] = new_color
                        st.session_state.graph_settings.setdefault("bar_colors_per_sample", {})[bar_key] = new_color

                    with rc[2]:
                        opt_cols = st.columns(4)
                        for oi, (ol, ok) in enumerate(zip(option_labels, option_keys)):
                            bs[ok] = opt_cols[oi].checkbox(
                                ol, bs.get(ok, True),
                                key=f"{ok}_{current_gene}_{condition}",
                                label_visibility="collapsed",
                            )
                    # Keep legacy show_sig in sync (all-or-nothing fallback)
                    bs["show_sig"] = bs["show_sig_1"] or bs["show_sig_2"] or bs["show_sig_3"]

                _bar_settings = st.session_state[f"{current_gene}_bar_settings"]

        current_settings = st.session_state.graph_settings.copy()
        current_settings["show_significance"] = st.session_state.graph_settings.get(
            show_sig_key, True
        )
        current_settings["show_error"] = st.session_state.graph_settings.get(
            show_err_key, True
        )
        current_settings["bar_gap"] = st.session_state.graph_settings.get(
            bar_gap_key, 0.45
        )
        # Override global defaults with per-gene values
        gs = st.session_state.graph_settings
        current_settings["figure_width"] = gs.get(f"{current_gene}_figure_width", gs.get("figure_width", 28))
        current_settings["figure_height"] = gs.get(f"{current_gene}_figure_height", gs.get("figure_height", 16))
        current_settings["font_size"] = gs.get(f"{current_gene}_font_size", gs.get("font_size", 14))
        current_settings["bar_opacity"] = gs.get(f"{current_gene}_bar_opacity", gs.get("bar_opacity", 0.85))
        current_settings["marker_line_width"] = gs.get(f"{current_gene}_marker_line_width", gs.get("marker_line_width", 1))
        current_settings["y_min"] = gs.get(f"{current_gene}_y_min")
        current_settings["y_max"] = gs.get(f"{current_gene}_y_max")
        current_settings["label_mode"] = gs.get(f"{current_gene}_label_mode", "Auto-wrap")

        display_gene_name = st.session_state.gene_display_names.get(current_gene, current_gene)

        ref_condition = st.session_state.graph_settings.get(f"{current_gene}_ref_line", "None")
        ref_line_val = None
        ref_line_lbl = None
        if ref_condition != "None":
            ref_rows = gene_data[gene_data["Condition"] == ref_condition]
            if not ref_rows.empty:
                _rv = ref_rows["Relative_Expression"].values[0]
                if pd.notna(_rv):
                    ref_line_val = _rv
                    ref_line_lbl = f"{ref_condition}: {_rv:.2f}"

        # Prepare data points overlay if enabled
        replicate_df = None
        show_dp = st.session_state.graph_settings.get(f"{current_gene}_show_data_points", False)
        if show_dp:
            raw = st.session_state.get("data")
            hk = st.session_state.get("hk_gene")
            ref = st.session_state.get("analysis_ref_condition")
            mapping = st.session_state.get("sample_mapping", {})
            excl = st.session_state.get("excluded_wells", set())
            if raw is not None and hk and ref:
                all_replicates = AnalysisEngine.compute_replicate_fold_changes(raw, hk, ref, mapping, excl)
                replicate_df = all_replicates[all_replicates["Target"] == current_gene]

        fig = GraphGenerator.create_gene_graph(
            gene_data,
            current_gene,
            current_settings,
            efficacy_config,
            sample_order=st.session_state.get("sample_order"),
            per_sample_overrides=None,
            display_gene_name=display_gene_name,
            ref_line_value=ref_line_val,
            ref_line_label=ref_line_lbl,
            show_data_points=show_dp,
            replicate_data=replicate_df,
            color_preset=st.session_state.graph_settings.get(f"{current_gene}_color_preset", "Classic"),
            ref_condition=st.session_state.get("analysis_ref_condition"),
        )

        st.plotly_chart(fig, use_container_width=True, key=f"main_fig_{current_gene}")
        st.session_state.graphs[current_gene] = fig

        with st.expander("📊 All Gene Graphs (Quick View)", expanded=False):
            _qv_genes = [g for g in gene_list if g != current_gene]
            all_gene_cols = st.columns(min(len(_qv_genes), 2)) if _qv_genes else [st.container()]
            for idx, gene in enumerate(_qv_genes):
                gd = st.session_state.processed_data[gene]

                if f"{gene}_bar_settings" not in st.session_state:
                    st.session_state[f"{gene}_bar_settings"] = {}
                    for _, row in gd.iterrows():
                        condition = row["Condition"]
                        bar_key = f"{gene}_{condition}"
                        st.session_state[f"{gene}_bar_settings"][bar_key] = {
                            "color": "#FFFFFF",
                            "show_sig": True,
                            "show_err": True,
                        }

                gs = st.session_state.graph_settings.copy()
                gs["show_significance"] = gs.get(f"{gene}_show_sig", True)
                gs["show_error"] = gs.get(f"{gene}_show_err", True)
                gs["bar_gap"] = gs.get(f"{gene}_bar_gap", 0.45)
                gs["figure_height"] = gs.get(f"{gene}_figure_height", gs.get("figure_height", 16))
                gs["figure_width"] = gs.get(f"{gene}_figure_width", gs.get("figure_width", 28))
                gs["font_size"] = gs.get(f"{gene}_font_size", gs.get("font_size", 14))
                gs["bar_opacity"] = gs.get(f"{gene}_bar_opacity", gs.get("bar_opacity", 0.85))
                gs["marker_line_width"] = gs.get(f"{gene}_marker_line_width", gs.get("marker_line_width", 1))
                gs["y_min"] = gs.get(f"{gene}_y_min")
                gs["y_max"] = gs.get(f"{gene}_y_max")
                gs["label_mode"] = gs.get(f"{gene}_label_mode", "Auto-wrap")
                # sig_style removed — direct mode only

                qv_ref_condition = st.session_state.graph_settings.get(f"{gene}_ref_line", "None")
                qv_ref_val = None
                qv_ref_lbl = None
                if qv_ref_condition != "None":
                    qv_ref_rows = gd[gd["Condition"] == qv_ref_condition]
                    if not qv_ref_rows.empty:
                        _qv_rv = qv_ref_rows["Relative_Expression"].values[0]
                        if pd.notna(_qv_rv):
                            qv_ref_val = _qv_rv
                            qv_ref_lbl = f"{qv_ref_condition}: {_qv_rv:.2f}"

                # Data points for quick view
                qv_show_dp = gs.get(f"{gene}_show_data_points", False)
                qv_replicate_df = None
                if qv_show_dp:
                    raw = st.session_state.get("data")
                    hk = st.session_state.get("hk_gene")
                    ref = st.session_state.get("analysis_ref_condition")
                    mapping = st.session_state.get("sample_mapping", {})
                    excl = st.session_state.get("excluded_wells", set())
                    if raw is not None and hk and ref:
                        all_reps = AnalysisEngine.compute_replicate_fold_changes(raw, hk, ref, mapping, excl)
                        qv_replicate_df = all_reps[all_reps["Target"] == gene]

                f = GraphGenerator.create_gene_graph(
                    gd,
                    gene,
                    gs,
                    efficacy_config,
                    sample_order=st.session_state.get("sample_order"),
                    display_gene_name=st.session_state.gene_display_names.get(gene, gene),
                    ref_line_value=qv_ref_val,
                    ref_line_label=qv_ref_lbl,
                    show_data_points=qv_show_dp,
                    replicate_data=qv_replicate_df,
                    color_preset=gs.get(f"{gene}_color_preset", "Classic"),
                    ref_condition=st.session_state.get("analysis_ref_condition"),
                )

                with all_gene_cols[idx % len(all_gene_cols)]:
                    st.markdown(f"**{gene}**")
                    st.plotly_chart(f, use_container_width=True, key=f"mini_{gene}")
                    st.session_state.graphs[gene] = f
    else:
        st.info(
            "⏳ No analysis results yet. Go to the 'Mapping' tab and click 'Run Full Analysis'."
        )

# ==================== TAB 5: EXPORT ====================
with tab5:
    render_step_indicator("export")
    st.header("Export Results")

    # Detect when settings have drifted from the snapshot taken during the last
    # analysis run. If they have, the significance markers / fold changes in
    # processed_data no longer reflect what the user currently has selected,
    # and any Excel/PPT exported now will carry the stale values.
    if st.session_state.get("processed_data") and "_exclusion_snapshot" in st.session_state:
        _snap = st.session_state["_exclusion_snapshot"]
        _cur = {
            "excluded_wells": {
                str(k): sorted(v)
                for k, v in st.session_state.get("excluded_wells", {}).items()
            },
            "excluded_samples": sorted(st.session_state.get("excluded_samples", set())),
            "ttest_type": st.session_state.get("ttest_type", "welch"),
        }
        if _cur != _snap:
            st.warning(
                "⚠️ Settings changed since the last analysis run "
                "(QC exclusions, excluded samples, or t-test type). The exports "
                "below still reflect the previous run. Open the **Analysis** tab "
                "to auto-rerun, or click **Re-run Analysis** there before exporting."
            )
        if st.session_state.get("analysis_stale", False):
            st.warning(
                "⚠️ Sample mapping or comparison group changed since the last "
                "analysis. Significance values may be outdated — re-run analysis "
                "before exporting."
            )

        # Cross-check that the stored ref/cmp condition names still exist in the
        # current mapping. If the user renamed a condition after running
        # analysis, processed_data labels diverge from the mapping and exports
        # will carry the old name.
        _mapped_conditions = {
            (st.session_state.get("sample_mapping", {}) or {}).get(s, {}).get("condition", s)
            for s in (st.session_state.get("sample_mapping", {}) or {})
        }
        _ref = st.session_state.get("analysis_ref_condition")
        _missing = [
            (_label, _val)
            for _label, _val in (
                ("Reference", _ref),
                ("Compare (*)", st.session_state.get("analysis_cmp_condition")),
                ("Compare (#)", st.session_state.get("analysis_cmp_condition_2")),
                ("Compare (†)", st.session_state.get("analysis_cmp_condition_3")),
            )
            if _val and _mapped_conditions and _val not in _mapped_conditions
        ]
        if _missing:
            _items = ", ".join(f"{lbl}={v!r}" for lbl, v in _missing)
            st.warning(
                f"⚠️ Some analysis comparison conditions are no longer present "
                f"in the current sample mapping: {_items}. The exports will "
                f"label charts with the old condition names. Re-run analysis "
                f"so labels match the current mapping."
            )

    if st.session_state.processed_data:
        analysis_params = {
            "Date": datetime.now().strftime("%Y-%m-%d %H:%M"),
            "Efficacy_Type": st.session_state.selected_efficacy,
            "Housekeeping_Gene": st.session_state.hk_gene,
            "Reference_Sample": st.session_state.get("analysis_ref_condition", "N/A"),
            "Compare_To": st.session_state.get("analysis_cmp_condition", "N/A"),
            "Compare_To_2": st.session_state.get("analysis_cmp_condition_2"),
            "Compare_To_3": st.session_state.get("analysis_cmp_condition_3"),
            "Excluded_Wells": sum(len(ws) for ws in st.session_state.excluded_wells.values()) if isinstance(st.session_state.excluded_wells, dict) else len(st.session_state.excluded_wells),
            "Excluded_Samples": len(st.session_state.excluded_samples),
            "Genes_Analyzed": len(st.session_state.processed_data),
            "concentration": st.session_state.experiment_desc.get("concentration", "1 ppm"),
            "treatment_time": st.session_state.experiment_desc.get("treatment_time", "24 h"),
            "ttest_type": st.session_state.get("ttest_type", "welch"),
        }

        def _build_export_extras():
            """Build qc_stats, replicate_stats, and excluded_wells for export."""
            qc_stats = None
            replicate_stats_df = None
            excl = st.session_state.get("excluded_wells", {})

            if st.session_state.get("data") is not None:
                excl_flat = set()
                if isinstance(excl, dict):
                    for well_set in excl.values():
                        excl_flat.update(well_set)
                elif excl:
                    excl_flat = set(excl)

                qc_data = st.session_state.data.copy()
                if excl_flat:
                    qc_data = qc_data[~qc_data["Well"].isin(excl_flat)]

                replicate_stats_df = QualityControl.get_replicate_stats(qc_data)
                qc_summary = QualityControl.get_qc_summary_stats(st.session_state.data, excl_flat or None)
                qc_stats = qc_summary if qc_summary else None

            return qc_stats, replicate_stats_df, excl

        efficacy = st.session_state.selected_efficacy
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")

        # PPT experiment description (collapsed)
        with st.expander("Experiment Description (for PPT)", expanded=False):
            conc_val = st.text_input(
                "Sample concentration",
                value=st.session_state.experiment_desc.get("concentration", "1 ppm"),
                key="exp_concentration",
            )
            st.session_state.experiment_desc["concentration"] = conc_val
            time_val = st.text_input(
                "Treatment time",
                value=st.session_state.experiment_desc.get("treatment_time", "24 h"),
                key="exp_treatment_time",
            )
            st.session_state.experiment_desc["treatment_time"] = time_val

        # ---- Reports ----
        st.subheader("Reports")

        rpt_cols = st.columns(2)
        with rpt_cols[0]:
            # Gated behind a button + spinner: export_to_excel post-processes chart
            # XML per gene, so running it on every tab rerun (slider moves, expander
            # toggles) is wasteful. Mirrors the PowerPoint pattern below.
            if st.button("Generate Excel Report", key="gen_excel", use_container_width=True):
                with st.spinner("Building Excel report..."):
                    try:
                        _qc, _rep, _excl = _build_export_extras()
                        st.session_state["_excel_export"] = export_to_excel(
                            st.session_state.data, st.session_state.processed_data,
                            analysis_params, st.session_state.sample_mapping,
                            qc_stats=_qc, replicate_stats=_rep, excluded_wells=_excl,
                            gene_display_names=st.session_state.get("gene_display_names", {}),
                        )
                    except Exception as e:
                        st.error(f"Excel generation failed: {e}")
            if "_excel_export" in st.session_state:
                st.download_button(
                    "Download Excel Report", data=st.session_state["_excel_export"],
                    file_name=f"qPCR_{efficacy}_{timestamp}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                )

        with rpt_cols[1]:
            if st.button("Generate PowerPoint", key="gen_ppt", use_container_width=True):
                with st.spinner("Generating PPT..."):
                    try:
                        ppt_bytes = PPTGenerator.generate_presentation(
                            st.session_state.graphs, st.session_state.processed_data,
                            analysis_params, graph_settings=st.session_state.get("graph_settings"),
                            gene_display_names=st.session_state.get("gene_display_names", {}),
                        )
                        if ppt_bytes:
                            st.session_state["_ppt_export"] = ppt_bytes
                    except Exception as e:
                        st.error(f"PPT generation failed: {e}")
            if "_ppt_export" in st.session_state:
                ppt_data = st.session_state["_ppt_export"]
                st.download_button(
                    "Download PowerPoint",
                    data=ppt_data.getvalue() if hasattr(ppt_data, "getvalue") else ppt_data,
                    file_name=f"qPCR_Report_{efficacy}_{timestamp}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )

        rpt_cols2 = st.columns(2)
        with rpt_cols2[0]:
            if st.session_state.graphs:
                try:
                    html_parts = [
                        "<html><head><title>qPCR Graphs</title></head><body>",
                        f"<h1>{efficacy} Analysis</h1>",
                        f"<p>Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}</p>",
                    ]
                    _disp_map = st.session_state.get("gene_display_names", {})
                    for gene, fig in st.session_state.graphs.items():
                        html_parts.append(f"<h2>{_disp_map.get(gene, gene)}</h2>")
                        html_parts.append(fig.to_html(full_html=False, include_plotlyjs="cdn"))
                    html_parts.append("</body></html>")
                    combined_html = "\n".join(html_parts)
                    st.download_button(
                        "Download Interactive HTML", data=combined_html,
                        file_name=f"qPCR_graphs_{efficacy}_{timestamp}.html",
                        mime="text/html", use_container_width=True,
                    )
                except Exception as e:
                    st.error(f"HTML generation failed: {e}")

        with rpt_cols2[1]:
            try:
                config_data = {
                    "analysis_params": analysis_params,
                    "sample_mapping": st.session_state.sample_mapping,
                    "graph_settings": st.session_state.graph_settings,
                    "excluded_wells": {f"{k[0]}|{k[1]}": list(v) for k, v in st.session_state.excluded_wells.items()} if isinstance(st.session_state.excluded_wells, dict) else list(st.session_state.excluded_wells),
                }
                st.download_button(
                    "Download Analysis Config",
                    data=json.dumps(config_data, indent=2),
                    file_name=f"config_{timestamp}.json",
                    mime="application/json", use_container_width=True,
                )
            except Exception as e:
                st.error(f"Config generation failed: {e}")

        # ---- Provenance (reproducibility record) ----
        if st.session_state.get("processed_data"):
            try:
                st.download_button(
                    "⬇️ Download Provenance (JSON)",
                    data=json.dumps(_current_provenance(), indent=2),
                    file_name=f"qPCR_provenance_{efficacy}_{timestamp}.json",
                    mime="application/json", use_container_width=True, key="dl_provenance",
                )
            except Exception as e:
                st.error(f"Provenance generation failed: {e}")

        # ---- One-click full bundle ----
        st.markdown("---")
        if st.session_state.graphs:
            if st.button("📦 Generate Complete Report Bundle (ZIP)", key="gen_bundle",
                         use_container_width=True,
                         help="Excel + PowerPoint + interactive HTML + PNG images in one download"):
                bundle = {}
                _disp_map = st.session_state.get("gene_display_names", {})
                with st.spinner("Building complete report bundle..."):
                    # Provenance (reproducibility record)
                    try:
                        bundle[f"provenance_{efficacy}_{timestamp}.json"] = json.dumps(
                            _current_provenance(), indent=2)
                    except Exception as e:
                        st.warning(f"Provenance skipped: {e}")
                    # Excel
                    try:
                        _qc, _rep, _excl = _build_export_extras()
                        bundle[f"qPCR_{efficacy}_{timestamp}.xlsx"] = export_to_excel(
                            st.session_state.data, st.session_state.processed_data,
                            analysis_params, st.session_state.sample_mapping,
                            qc_stats=_qc, replicate_stats=_rep, excluded_wells=_excl,
                            gene_display_names=_disp_map,
                        )
                    except Exception as e:
                        st.warning(f"Excel skipped: {e}")
                    # PowerPoint
                    try:
                        ppt_bytes = PPTGenerator.generate_presentation(
                            st.session_state.graphs, st.session_state.processed_data,
                            analysis_params, graph_settings=st.session_state.get("graph_settings"),
                            gene_display_names=_disp_map,
                        )
                        bundle[f"qPCR_Report_{efficacy}_{timestamp}.pptx"] = ppt_bytes
                    except Exception as e:
                        st.warning(f"PowerPoint skipped: {e}")
                    # Interactive HTML
                    try:
                        _html = ["<html><head><meta charset='utf-8'><title>qPCR Graphs</title></head><body>",
                                 f"<h1>{efficacy} Analysis</h1>",
                                 f"<p>Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}</p>"]
                        for gene, fig in st.session_state.graphs.items():
                            _html.append(f"<h2>{_disp_map.get(gene, gene)}</h2>")
                            _html.append(fig.to_html(full_html=False, include_plotlyjs="cdn"))
                        _html.append("</body></html>")
                        bundle[f"qPCR_graphs_{efficacy}_{timestamp}.html"] = "\n".join(_html)
                    except Exception as e:
                        st.warning(f"HTML skipped: {e}")
                    # PNG images (high-res)
                    for gene, fig in st.session_state.graphs.items():
                        try:
                            _fc = go.Figure(fig)
                            _om = fig.layout.margin
                            _pb = max(180, _om.b if _om and _om.b else 180)
                            _fc.update_layout(width=1200, height=800 + max(0, _pb - 180),
                                margin=dict(b=_pb), font=dict(size=14, family=PLOTLY_FONT_FAMILY, color="black"))
                            bundle[f"images/{gene}.png"] = ReportGenerator._fig_to_image(
                                _fc, format="png", scale=3, width=1200, height=800 + max(0, _pb - 180))
                        except Exception as e:
                            st.warning(f"Image skipped ({gene}): {e}")
                if bundle:
                    st.session_state["_bundle_export"] = build_zip(bundle)
                    st.success(f"Bundle ready — {len(bundle)} files.")
            if "_bundle_export" in st.session_state:
                st.download_button(
                    "⬇️ Download Complete Report (ZIP)", data=st.session_state["_bundle_export"],
                    file_name=f"qPCR_full_report_{efficacy}_{timestamp}.zip",
                    mime="application/zip", use_container_width=True,
                )

        # ---- Gene Images ----
        st.markdown("---")
        st.subheader("Gene Images")

        img_col1, img_col2, img_col3 = st.columns(3)
        with img_col1:
            img_format = st.selectbox("Format", ["PNG (300 DPI)", "SVG (Vector)", "PDF (Vector)"], key="pub_img_format")
        with img_col2:
            img_width = st.number_input("Width (px)", min_value=400, max_value=3000, value=1200, step=100)
        with img_col3:
            img_height = st.number_input("Height (px)", min_value=300, max_value=2000, value=800, step=100)

        if st.session_state.graphs:
            fmt = "png" if "PNG" in img_format else "svg" if "SVG" in img_format else "pdf"
            mime = {"png": "image/png", "svg": "image/svg+xml", "pdf": "application/pdf"}[fmt]
            scale = 3 if fmt == "png" else 1
            # Rendering each figure launches headless Chrome (~seconds each), so gate
            # it behind an explicit button instead of re-rendering on every rerun.
            if st.button(f"🖼️ Generate Images ({fmt.upper()})", key="gen_images", use_container_width=True):
                rendered, failed = {}, []
                prog = st.progress(0.0, text="Rendering images...")
                genes = list(st.session_state.graphs.items())
                for idx, (gene, fig) in enumerate(genes):
                    try:
                        fig_copy = go.Figure(fig)
                        _orig_m = fig.layout.margin
                        _pub_b = max(180, _orig_m.b if _orig_m and _orig_m.b else 180)
                        _adj_h = img_height + max(0, _pub_b - 180)
                        fig_copy.update_layout(width=img_width, height=_adj_h, margin=dict(b=_pub_b),
                            font=dict(size=14, family=PLOTLY_FONT_FAMILY, color="black"))
                        rendered[gene] = ReportGenerator._fig_to_image(
                            fig_copy, format=fmt, scale=scale, width=img_width, height=_adj_h)
                    except Exception as e:
                        failed.append(f"{gene}: {e}")
                    prog.progress((idx + 1) / len(genes), text=f"Rendered {idx + 1}/{len(genes)}")
                prog.empty()
                st.session_state["_gene_images"] = {"fmt": fmt, "images": rendered}
                if failed:
                    st.warning("Some images failed:\n" + "\n".join(failed))

            cached = st.session_state.get("_gene_images")
            if cached and cached.get("images"):
                _cfmt = cached["fmt"]
                _cmime = {"png": "image/png", "svg": "image/svg+xml", "pdf": "application/pdf"}[_cfmt]
                today = datetime.now().strftime("%Y%m%d")
                st.download_button(
                    f"⬇️ Download All Images (.zip, {len(cached['images'])})",
                    data=build_zip({f"{g}.{_cfmt}": b for g, b in cached["images"].items()}),
                    file_name=f"qPCR_images_{efficacy}_{today}.zip",
                    mime="application/zip", use_container_width=True, key="dl_images_zip",
                )
                img_cols = st.columns(min(len(cached["images"]), 4))
                for idx, (gene, img_bytes) in enumerate(cached["images"].items()):
                    with img_cols[idx % len(img_cols)]:
                        st.download_button(label=f"{gene}.{_cfmt}", data=img_bytes,
                            file_name=f"{gene}_{today}.{_cfmt}",
                            mime=_cmime, key=f"img_{gene}", use_container_width=True)
    else:
        st.warning("Complete analysis first.")

# ==================== FOOTER ====================
st.markdown("---")

footer_html = """
<div style='text-align: center; color: #86868b; padding: 24px 0; font-size: 0.85rem;'>
    <p style='margin: 0;'>qPCR Analysis Suite v3.1</p>
</div>
"""

st.markdown(footer_html, unsafe_allow_html=True)
